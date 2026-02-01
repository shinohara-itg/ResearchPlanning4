import os
import tempfile
import zipfile
from pathlib import Path
import streamlit as st
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image
from openai import AzureOpenAI
from dotenv import load_dotenv
import re
from pptx.dml.color import RGBColor
import json  # ← 先頭へ



# =========================
# ページ設定
# =========================
st.set_page_config(page_title="みんなの知恵", layout="wide")
st.title("みんなの知恵")

# =========================
# セッション初期化
# =========================
if "selected_mode" not in st.session_state:
    st.session_state["selected_mode"] = None
if "message_center" not in st.session_state:
    st.session_state["message_center"] = ""
if "message_right" not in st.session_state:
    st.session_state["message_right"] = ""
if "uploaded_docs" not in st.session_state:
    st.session_state["uploaded_docs"] = []
if "pptx_path" not in st.session_state:
    st.session_state["pptx_path"] = None
if "edited_texts" not in st.session_state:
    st.session_state["edited_texts"] = {}
if "orien_outline_text" not in st.session_state:
    st.session_state["orien_outline_text"] = ""
if "orien_company_text" not in st.session_state:
    st.session_state["orien_company_text"] = ""
if "final_pptx_path" not in st.session_state:
    st.session_state["final_pptx_path"] = None
# =========================
# Revision（課題ピボット）用セッション初期化
# =========================
if "proposal_revisions" not in st.session_state:
    st.session_state["proposal_revisions"] = []  # list[dict]
if "active_rev_id" not in st.session_state:
    st.session_state["active_rev_id"] = None
if "kickoff_selected_purpose" not in st.session_state:
    st.session_state["kickoff_selected_purpose"] = None
if "orien_auto_generated" not in st.session_state:
    st.session_state["orien_auto_generated"] = False
if "kickoff_purpose_free_editor" not in st.session_state:
    st.session_state["kickoff_purpose_free_editor"] = ""
if "orien_outline_ai_draft" not in st.session_state:
    st.session_state["orien_outline_ai_draft"] = ""
if "orien_outline_manual" not in st.session_state:
    st.session_state["orien_outline_manual"] = ""

if "revisions_version" not in st.session_state:
    st.session_state["revisions_version"] = 0
# =========================
# 課題変換（前処理）用セッション初期化
# =========================
if "problem_reframe_generated" not in st.session_state:
    st.session_state["problem_reframe_generated"] = False

if "problem_reframe_output" not in st.session_state:
    st.session_state["problem_reframe_output"] = {}  # JSON固定

if "true_problem_text" not in st.session_state:
    st.session_state["true_problem_text"] = ""  # 採用課題のみ抜粋（キックオフ生成に注入）
if "reframe_logic_map" not in st.session_state:
    st.session_state["reframe_logic_map"] = ""
if "reframe_c6_user_notes" not in st.session_state:
    st.session_state["reframe_c6_user_notes"] = ""  # 6観点目：ユーザー任意追記
if "orien_outline_ai_draft_store" not in st.session_state:
    st.session_state["orien_outline_ai_draft_store"] = ""
if "orien_outline_manual_store" not in st.session_state:
    st.session_state["orien_outline_manual_store"] = ""
# =========================
# オリエン（永続データ）用キー：widget key と分離
# =========================
if "data_orien_outline_ai_draft" not in st.session_state:
    st.session_state["data_orien_outline_ai_draft"] = ""
if "data_orien_outline_manual" not in st.session_state:
    st.session_state["data_orien_outline_manual"] = ""

# UIウィジェット用（表示時に data -> ui へ注入する）
if "ui_orien_outline_ai_draft" not in st.session_state:
    st.session_state["ui_orien_outline_ai_draft"] = st.session_state["data_orien_outline_ai_draft"]
if "ui_orien_outline_manual" not in st.session_state:
    st.session_state["ui_orien_outline_manual"] = st.session_state["data_orien_outline_manual"]

if "dbg_logs" not in st.session_state: st.session_state["dbg_logs"] = []

if "dbg_enabled" not in st.session_state:
    st.session_state["dbg_enabled"] = True


# if "__dbg_before_upsert" in st.session_state:
#     st.warning(f"DEBUG before_upsert: {st.session_state['__dbg_before_upsert']}")
# if "__dbg_after_upsert" in st.session_state:
#     st.warning(f"DEBUG after_upsert: {st.session_state['__dbg_after_upsert']}")



from pathlib import Path
import streamlit as st


#案内画像用
CHAR_IMG_PATH = Path(__file__).parent / "assets" / "character.png"
CHAR_IMG_PATH2 = Path(__file__).parent / "assets" / "character2.png"
CHAR_IMG_PATH3 = Path(__file__).parent / "assets" / "character3.png"
CHAR_IMG_PATH4 = Path(__file__).parent / "assets" / "character4.png"

def render_character_guide(title: str, body_md: str, *, img_width: int = 300, kind: str = "info"):
    """
    kind: "info" | "warning" | "success" | "error"
    """
    col1, col2 = st.columns([1,3], gap="medium")
    with col1:
        if CHAR_IMG_PATH.exists():
            st.image(str(CHAR_IMG_PATH), width=img_width)
        else:
            st.caption("（character.png が見つかりません）")

    with col2:
        st.markdown(f"### {title}")
        if kind == "warning":
            st.warning(body_md)
        elif kind == "success":
            st.success(body_md)
        elif kind == "error":
            st.error(body_md)
        else:
            st.info(body_md)


def render_character_guide2(title: str, body_md: str, *, img_width: int = 300, kind: str = "info"):
    """
    kind: "info" | "warning" | "success" | "error"
    """
    col1, col2 = st.columns([1,5], gap="medium")
    with col1:
        if CHAR_IMG_PATH.exists():
            st.image(str(CHAR_IMG_PATH2), width=img_width)
        else:
            st.caption("（character.png が見つかりません）")

    with col2:
        st.markdown(f"### {title}")
        if kind == "warning":
            st.warning(body_md)
        elif kind == "success":
            st.success(body_md)
        elif kind == "error":
            st.error(body_md)
        else:
            st.info(body_md)


def render_character_guide3(title: str, body_md: str, *, img_width: int = 300, kind: str = "info"):
    """
    kind: "info" | "warning" | "success" | "error"
    """
    col1, col2 = st.columns([1,5], gap="medium")
    with col1:
        if CHAR_IMG_PATH.exists():
            st.image(str(CHAR_IMG_PATH3), width=img_width)
        else:
            st.caption("（character.png が見つかりません）")

    with col2:
        st.markdown(f"### {title}")
        if kind == "warning":
            st.warning(body_md)
        elif kind == "success":
            st.success(body_md)
        elif kind == "error":
            st.error(body_md)
        else:
            st.info(body_md)



def render_character_guide4(title: str, body_md: str, *, img_width: int = 300, kind: str = "info"):
    """
    kind: "info" | "warning" | "success" | "error"
    """
    col1, col2 = st.columns([1,5], gap="medium")
    with col1:
        if CHAR_IMG_PATH.exists():
            st.image(str(CHAR_IMG_PATH4), width=img_width)
        else:
            st.caption("（character.png が見つかりません）")

    with col2:
        st.markdown(f"### {title}")
        if kind == "warning":
            st.warning(body_md)
        elif kind == "success":
            st.success(body_md)
        elif kind == "error":
            st.error(body_md)
        else:
            st.info(body_md)


# =========================
# Azure OpenAI 設定
# =========================
load_dotenv()
client = AzureOpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)
DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o")



# import datetime  ← これは削除

from datetime import datetime, timedelta


def dbg_issue_conv(msg: str):
    if not st.session_state.get("dbg_enabled", True):
        return
    ts = datetime.now().strftime("%H:%M:%S.%f")[:-3]
    st.session_state["dbg_logs"].append(f"{ts} | {msg}")
    st.session_state["dbg_logs"] = st.session_state["dbg_logs"][-200:]


def render_dbg_sidebar():
    dbg_issue_conv("app rerun")

    with st.sidebar.expander("Debug log", expanded=True):
        col1, col2 = st.columns(2)
        with col1:
            st.session_state["dbg_enabled"] = st.toggle(
                "Enable", value=st.session_state.get("dbg_enabled", True)
            )
        with col2:
            if st.button("Clear"):
                st.session_state["dbg_logs"] = []
        logs = st.session_state.get("dbg_logs", [])
        st.text("\n".join(logs[-80:]) if logs else "(no logs)")
def init_session_state():
    # 既存の初期化（mode や dbg_log 等）
    st.session_state.setdefault("mode", "オリエン内容の整理")
    st.session_state.setdefault("dbg_log", [])
    # --- ここから追記：課題変換 data/ui の初期化 ---
    REFRAME_FIELDS = [
        "c1_next_action",
        "c2_exec_summary",
        "c4_business_brand",
        "c6_user_notes",
    ]

    # 永続データ（正本）
    for f in REFRAME_FIELDS:
        st.session_state.setdefault(f"data_reframe_{f}", "")

    # UI用（表示用コピー）
    for f in REFRAME_FIELDS:
        st.session_state.setdefault(
            f"ui_reframe_{f}",
            st.session_state.get(f"data_reframe_{f}", "")
        )

REFRAME_FIELDS = [
    "c1_next_action",
    "c2_exec_summary",
    "c4_business_brand",
    "c6_user_notes",
]

def sync_reframe_from_ui():
    for f in REFRAME_FIELDS:
        st.session_state[f"data_reframe_{f}"] = st.session_state.get(f"ui_reframe_{f}", "")

def hydrate_reframe_ui_from_data_if_empty():
    """
    課題変換画面に入ったとき、
    ui が空のものだけ data を流し込む（編集の上書きを避けるため）
    """
    for f in REFRAME_FIELDS:
        ui_key = f"ui_reframe_{f}"
        if not (st.session_state.get(ui_key) or "").strip():
            st.session_state[ui_key] = st.session_state.get(f"data_reframe_{f}", "")

def debug_log(msg: str):
    st.session_state.setdefault("dbg_log", [])
    ts = datetime.now().strftime("%H:%M:%S.%f")[:-3]
    st.session_state["dbg_log"].append(f"{ts} | {msg}")
    print(f"{ts} | {msg}")

# --- imports の直後あたりに追加 ---
def init_state():
    if "mode" not in st.session_state:
        st.session_state.mode = "Tab1"   # 初期モード（あなたの実態に合わせて）
    # もし他にも「消えると困る」状態があるならここで同様に初期化

# init_state()





# =========================
# 古いセッションの自動クリーンアップ
# =========================
def cleanup_old_sessions(days: int = 1):
    """
    最終アクセスから days 日以上経過したセッションディレクトリを削除
    """
    if not BASE_ROOT.exists():
        return

    now = datetime.now()
    for child in BASE_ROOT.iterdir():
        if not child.is_dir():
            continue

        marker = child / ".last_access"
        try:
            if marker.exists():
                ts = datetime.fromisoformat(marker.read_text(encoding="utf-8"))
            else:
                # マーカーがない場合はディレクトリの更新時刻を使う
                ts = datetime.fromtimestamp(child.stat().st_mtime)

            if now - ts > timedelta(days=days):
                shutil.rmtree(child, ignore_errors=True)
        except Exception:
            # 読み取り・削除で何かあってもアプリを落とさない
            continue



# =========================
# セッション専用ディレクトリのヘルパーを作る
# =========================
import uuid
import shutil
from datetime import datetime, timedelta

BASE_ROOT = Path("/home/streamlit_workspace")


def get_session_dir() -> Path:
    """
    セッションごとに一意の作業ディレクトリを返す。
    例）/home/streamlit_workspace/20250201_120000_ab12cd34/
    """
    cleanup_old_sessions(days=1) 

    if "session_id" not in st.session_state:
        sid = datetime.now().strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:8]
        st.session_state["session_id"] = sid

    session_dir = BASE_ROOT / st.session_state["session_id"]
    session_dir.mkdir(parents=True, exist_ok=True)

    # 最終アクセス時刻を記録しておく（自動クリーンアップ用）
    (session_dir / ".last_access").write_text(datetime.now().isoformat(), encoding="utf-8")

    return session_dir


def get_orien_context_text() -> str:
    docs_text = "\n".join(st.session_state.get("uploaded_docs", []) or []).strip()
    manual_text = (st.session_state.get("data_orien_outline_manual") or "").strip()
    ai_text = (st.session_state.get("data_orien_outline_ai_draft") or "").strip()


    parts = []
    if docs_text:
        parts.append("【オリエン資料（アップロード抽出テキスト）】\n" + docs_text)

    # 手入力が最優先
    if manual_text:
        parts.append("【オリエン内容レビュー（手入力：最優先）】\n" + manual_text)
    elif ai_text:
        # 手入力が無い場合のみAI整理結果を補助として使う
        parts.append("【オリエン内容の整理（AI）】\n" + ai_text)

    return "\n\n".join(parts).strip()


def safe_extract_zip(zip_path: str, extract_to: str) -> None:
    """
    Zip Slip対策：ZIP内のパスが extract_to 配下に収まる場合のみ展開する
    """
    extract_base = Path(extract_to).resolve()
    with zipfile.ZipFile(zip_path, "r") as z:
        for member in z.infolist():
            dest = (extract_base / member.filename).resolve()
            if not str(dest).startswith(str(extract_base)):
                raise ValueError(f"Unsafe zip entry detected: {member.filename}")
        z.extractall(extract_to)

# =========================
# ファイル読込関数
# =========================
def process_orien_upload(uploaded_files):
    """
    uploaded_files: st.file_uploader(..., accept_multiple_files=True) の戻り
    抽出テキストを st.session_state["uploaded_docs"] に格納し、
    変更検知や自動生成フラグも更新する。
    """
    if "orien_upload_sig" not in st.session_state:
        st.session_state["orien_upload_sig"] = None

    if not uploaded_files:
        return

    upload_sig = tuple(sorted((f.name, getattr(f, "size", None)) for f in uploaded_files))

    if st.session_state["orien_upload_sig"] != upload_sig:
        st.session_state["orien_upload_sig"] = upload_sig
        st.session_state["orien_auto_generated"] = False

        texts = []

        with tempfile.TemporaryDirectory() as tempdir:
            for file in uploaded_files:
                path = os.path.join(tempdir, file.name)
                with open(path, "wb") as f:
                    f.write(file.getbuffer())

                low = path.lower()
                if low.endswith(".pdf"):
                    texts.append(read_pdf(path))
                elif low.endswith(".pptx"):
                    texts.append(read_pptx_text(path))
                elif low.endswith(".txt"):
                    texts.append(read_txt(path))
                elif low.endswith(".docx"):
                    texts.append(read_docx(path))
                elif low.endswith(".xlsx"):
                    texts.append(read_xlsx(path))
                elif low.endswith(".zip"):
                    try:
                        safe_extract_zip(path, tempdir)  # ★Cの安全化をここでも使う
                    except Exception as e:
                        st.warning(f"ZIP展開をスキップしました（安全性チェックNG）: {e}")
                        continue

                    for root, _, files in os.walk(tempdir):
                        for fn in files:
                            fp = os.path.join(root, fn)
                            low2 = fp.lower()
                            if low2.endswith(".pdf"):
                                texts.append(read_pdf(fp))
                            elif low2.endswith(".pptx"):
                                texts.append(read_pptx_text(fp))
                            elif low2.endswith(".txt"):
                                texts.append(read_txt(fp))
                            elif low2.endswith(".docx"):
                                texts.append(read_docx(fp))
                            elif low2.endswith(".xlsx"):
                                texts.append(read_xlsx(fp))

        texts = [t for t in texts if isinstance(t, str) and t.strip()]
        st.session_state["uploaded_docs"] = texts

        st.success(f"資料を共有しました。（{len(uploaded_files)}件 / 抽出テキスト {len(texts)}件）")

def read_txt(path):
    for enc in ("utf-8", "utf-8-sig", "cp932"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""

def read_pdf(path):
    try:
        doc = fitz.open(path)
        return "\n".join(page.get_text("text") for page in doc)
    except Exception:
        return ""

def read_pptx_text(path):
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text:
                    texts.append(shp.text)
        return "\n".join(texts)
    except Exception:
        return ""
from docx import Document
import openpyxl

def read_docx(path: str) -> str:
    """Word(.docx) から本文テキストを抽出"""
    try:
        doc = Document(path)
        parts = []

        # 段落
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)

        # 表（テーブル）
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [(c.text or "").replace("\n", " ").strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))

        return "\n".join(parts).strip()
    except Exception:
        return ""


def read_xlsx(path: str) -> str:
    """Excel(.xlsx) から全シートのセル内容をテキスト化"""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)  # 数式は計算結果を読む
        parts = []
        for ws in wb.worksheets:
            parts.append(f"【Sheet】{ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = []
                for v in row:
                    if v is None:
                        continue
                    s = str(v).strip()
                    if s:
                        vals.append(s)
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts).strip()
    except Exception:
        return ""





# =========================
# PPT → 画像変換関数
# =========================
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import io, os

def pptx_to_images(pptx_path: Path) -> list[Image.Image]:
    """
    PowerPointファイルをスライドレイアウト通りに簡易描画して画像リストで返す。
    - 日本語フォント対応
    - テキスト・画像を元の位置(left, top, width, height)に再配置
    """
    images: list[Image.Image] = []

    # ---- 日本語フォント設定 ----
    FONT_CANDIDATES = [
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/fonts-japanese-gothic.ttf",
        "C:/Windows/Fonts/meiryo.ttc",
        "/System/Library/Fonts/ヒラギノ角ゴシック W4.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    font_path = next((f for f in FONT_CANDIDATES if os.path.exists(f)), None)
    if font_path:
        font_small = ImageFont.truetype(font_path, 20)
    else:
        font_small = ImageFont.load_default()

    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            # スライドサイズ（EMU → px換算）
            width_px = int(prs.slide_width / 9525)
            height_px = int(prs.slide_height / 9525)

            # 白背景キャンバス
            img = Image.new("RGB", (width_px, height_px), "white")
            draw = ImageDraw.Draw(img)

            # === スライド上の図形を順に描画 ===
            for shp in slide.shapes:
                left = int(shp.left / 9525)
                top = int(shp.top / 9525)
                width = int(shp.width / 9525)
                height = int(shp.height / 9525)

                # 図形タイプで分岐
                stype = shp.shape_type

                # 画像
                if stype == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = io.BytesIO(shp.image.blob)
                        pic = Image.open(image_bytes).convert("RGB")
                        pic = pic.resize((width, height))
                        img.paste(pic, (left, top))
                    except Exception:
                        draw.rectangle([left, top, left + width, top + height], outline="gray")
                        draw.text((left + 4, top + 4), "画像読み込み失敗", font=font_small, fill="red")

                # テキスト付き図形
                elif getattr(shp, "has_text_frame", False):
                    text = shp.text.strip()
                    if text:
                        # テキスト枠（背景塗り）
                        draw.rectangle([left, top, left + width, top + height], outline="lightgray", fill=None)
                        # テキスト（簡易左寄せ）
                        lines = text.replace("\r", "").split("\n")
                        y = top + 5
                        for line in lines:
                            draw.text((left + 8, y), line[:40], font=font_small, fill="black")
                            y += 24

                # 図形（塗りつぶしのみ）
                else:
                    draw.rectangle([left, top, left + width, top + height], outline="lightgray", fill=None)

            # スライド番号
            draw.text((20, height_px - 40), f"Slide {i+1}", font=font_small, fill="gray")

            images.append(img)

        return images

    except Exception as e:
        st.error(f"PPT変換エラー: {e}")
        return []
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import base64

def emu_to_percent(val_emu, total_emu):
    """EMU単位をスライド全体に対する%へ変換"""
    try:
        return float(val_emu) / float(total_emu) * 100.0
    except Exception:
        return 0.0


def color_to_css(rgb):
    """RGBColor → CSSカラーコード"""
    if not rgb:
        return None
    if isinstance(rgb, RGBColor):
        return f"#{rgb.rgb:06X}"
    try:
        return f"#{int(rgb):06X}"
    except Exception:
        return None


def extract_slide_model(prs, slide_index=0):
    """
    PowerPointスライド内の図形を走査し、
    Streamlit用のHTML描画モデルに変換する。
    - PICTURE：画像
    - TEXT：テキストボックス
    - TABLE：セルの文字を連結して1つのテキストブロックとして描画（★追加）
    """
    slide = prs.slides[slide_index]
    sw, sh = prs.slide_width, prs.slide_height
    blocks = []

    def add_block(shape, offset_left=0, offset_top=0):
        stype = shape.shape_type
        name = getattr(shape, "name", "")
        editable = name.startswith("Edit_") or name.startswith("EDIT_")

        left = shape.left + offset_left
        top = shape.top + offset_top
        width = shape.width
        height = shape.height

        base = {
            "name": name,
            "editable": editable,
            "left": emu_to_percent(left, sw),
            "top": emu_to_percent(top, sh),
            "width": emu_to_percent(width, sw),
            "height": emu_to_percent(height, sh),
        }

        # グループ処理
        if stype == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                add_block(child, offset_left=left, offset_top=top)
            return

        # 画像処理
        if stype == MSO_SHAPE_TYPE.PICTURE:
            try:
                content_type = getattr(shape.image, "content_type", "image/png")
                b64 = base64.b64encode(shape.image.blob).decode("ascii")
                base["type"] = "picture"
                base["src"] = f"data:{content_type};base64,{b64}"
            except Exception:
                pass
            blocks.append(base)
            return

        # 塗り（背景色）
        fill_css = None
        try:
            if shape.fill and shape.fill.type == 1:  # solid fill
                fill_css = color_to_css(shape.fill.fore_color.rgb)
        except Exception:
            pass

        # ★ TABLE（表）の処理を追加：セルのテキストを連結して1ブロックとして描画
        if stype == MSO_SHAPE_TYPE.TABLE:
            try:
                table = shape.table
                rows_text = []
                for row in table.rows:
                    cells = [
                        cell.text.replace("\n", " ").strip()
                        for cell in row.cells
                    ]
                    # 全部空なら無視
                    if any(cells):
                        rows_text.append(" | ".join(cells))
                text = "\n".join(rows_text).strip()
                if text:
                    blocks.append(
                        {
                            **base,
                            "type": "text",
                            "text": text,
                            "fill": fill_css,
                        }
                    )
            except Exception:
                # 失敗したら単なるボックスとして描画
                blocks.append({**base, "type": "box", "fill": fill_css})
            return

        # テキスト付き図形
        if getattr(shape, "has_text_frame", False):
            blocks.append({**base, "type": "text", "text": shape.text, "fill": fill_css})
        else:
            # 図形（塗りだけ）
            blocks.append({**base, "type": "box", "fill": fill_css})

    for shape in slide.shapes:
        add_block(shape, 0, 0)

    return {"blocks": blocks}


def render_slide_html(model, edited_texts):
    """
    extract_slide_model()で抽出した構造をもとに、
    Streamlit内でスライドの見た目を再現するHTMLを生成。
    """
    blocks = model["blocks"]

    html = """
    <div style="position:relative; width:100%; padding-top:56.25%; background:#f8f9fb;
                border-radius:14px; box-shadow:0 4px 16px rgba(0,0,0,0.08); overflow:hidden;">
      <div style="position:absolute; inset:0; background:white;">
    """

    for b in blocks:
        style = (
            f"position:absolute; left:{b['left']}%; top:{b['top']}%; "
            f"width:{b['width']}%; height:{b['height']}%;"
        )
        content = ""

        # 画像
        if b.get("type") == "picture" and b.get("src"):
            content = (
                f'<img src="{b["src"]}" style="width:100%;height:100%;object-fit:contain;">'
            )

        # テキスト
        elif b.get("type") == "text":
            text_val = edited_texts.get(b["name"], b.get("text", ""))
            bg = f'background:{b["fill"]};' if b.get("fill") else ""
            content = (
                f'<div style="{bg}padding:6px;font-family:\'Noto Sans JP\',sans-serif;'
                f'font-size:13px;color:#111;white-space:pre-wrap;">{text_val}</div>'
            )

        html += f'<div style="{style}">{content}</div>'

    html += "</div></div>"
    return html

def parse_ai_output(text: str):
    """AI出力を6項目に分割"""
    sections = {
        "目標": "",
        "現状": "",
        "ビジネス課題": "",
        "調査目的": "",
        "問い": "",
        "仮説": "",
    }
    for key in sections.keys():
        pattern = rf"【{key}】(.*?)(?=【|$)"
        m = re.search(pattern, text, re.DOTALL)
        if m:
            sections[key] = m.group(1).strip()
    return sections



from pptx.enum.shapes import MSO_SHAPE_TYPE  # 既にインポート済みならこの行は重複していてもOK

from pptx.enum.shapes import MSO_SHAPE_TYPE  # 既にインポート済みならこの行は重複していてもOK
from pptx.dml.color import RGBColor         # ← これも上にあれば重複OK

def set_text_to_named_shape(slide, shape_name: str, text: str) -> bool:
    """
    スライド内の図形（グループ内も含む）から name=shape_name を探し、
    テキストを書き込む。見つかれば True、見つからなければ False を返す。
    - オートシェイプ／プレースホルダー：.text に書き込む
    - テーブル：全セルに同じテキストを書き込む（暫定）
    - 書き込んだテキストの文字色は黒（RGB 0,0,0）に設定する
    """

    def _set_font_black_textframe(text_frame):
        """text_frame 内の全 run のフォント色を黒にする"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font is not None:
                        run.font.color.rgb = RGBColor(0, 0, 0)
        except Exception:
            # フォーマット構造が想定外でも落ちないようにする
            pass

    def _search(shapes):
        for shp in shapes:
            # グループ内なら再帰
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                if _search(shp.shapes):
                    return True

            # 名前一致チェック
            if shp.name == shape_name:
                # テキスト枠があるタイプ
                if getattr(shp, "has_text_frame", False):
                    shp.text = text
                    _set_font_black_textframe(shp.text_frame)
                    return True

                # テーブルの場合
                if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        for row in shp.table.rows:
                            for cell in row.cells:
                                cell.text = text
                                _set_font_black_textframe(cell.text_frame)
                        return True
                    except Exception:
                        pass

        return False

    return _search(slide.shapes)




# ★ 調査仕様の項目（ラベルと session_state のキー）
SPEC_ITEMS = [
    ("調査手法", "spec_method"),
    ("抽出方法", "spec_sampling"),
    ("調査地域", "spec_region"),
    ("対象者条件", "spec_target"),
    ("サンプルサイズ", "spec_sample_size"),
    ("調査ボリューム", "spec_volume"),
    ("提示物", "spec_stimulus"),
    ("集計・分析仕様", "spec_analysis"),
    ("自由回答データの処理", "spec_openend"),
    ("業務範囲", "spec_scope"),
    ("納品物", "spec_deliverables"),
    ("インスペクションの方法", "spec_inspection"),
    ("謝礼の種類", "spec_incentive"),
    ("備考", "spec_notes"),
]

# ★ 調査仕様スライド（スライド6）の shape 名との対応
SPEC_LABEL_TO_SHAPE = {
    "調査手法": "Edit_SYUHO",
    "抽出方法": "Edit_Sampling",
    "調査地域": "Edit_Area",
    "対象者条件": "Edit_Joken",
    "サンプルサイズ": "Edit_Samplesize",
    "調査ボリューム": "Edit_Qvolume",
    "提示物": "Edit_review",
    "集計・分析仕様": "Edit_Analitics",
    "自由回答データの処理": "Edit_OAcdg",
    "業務範囲": "Edit_Hani",
    "納品物": "Edit_Nohin",
    "インスペクションの方法": "Edit_Inspection",
    "謝礼の種類": "Edit_Syarei",
    "備考": "Edit_Biko",
}

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt

def apply_text_format(shape, font_name="Arial", font_size=12, color=RGBColor(0, 0, 0)):
    """
    shape.text_frame の paragraph/run に書式を統一的に適用する
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT  # 左寄せ
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

import re

def parse_subquestions(ai_text: str):
    """
    『問いの分解』モードのAI出力をパースして、
    [
      {"subq": "...", "axis": "...", "items": "..."},
      ...
    ]
    のリストに変換する
    """
    if not ai_text:
        return []

    # 「- サブクエスチョン...」でブロックごとに分割
    blocks = re.split(r"\n(?=-\s*サブクエスチョン)", ai_text.strip())
    results = []

    for blk in blocks:
        # サブクエスチョン本体
        m_q = re.search(r"-\s*サブクエスチョン[0-9０-９]*[:：]\s*(.+)", blk)
        if not m_q:
            continue

        # 分析軸
        m_axis = re.search(r"分析軸[:：]\s*(.+)", blk)
        # 評価項目
        m_items = re.search(r"評価項目[:：]\s*(.+)", blk)

        results.append(
            {
                "subq": m_q.group(1).strip(),
                "axis": m_axis.group(1).strip() if m_axis else "",
                "items": m_items.group(1).strip() if m_items else "",
            }
        )

    return results



from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_text_style(shape):
    """
    指定した図形内テキストの書式を統一するヘルパー
    - フォント：Arial
    - サイズ：12pt
    - 色：黒
    - 配置：左揃え
    """
    if not getattr(shape, "has_text_frame", False):
        return

    try:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                if run.font is None:
                    continue
                run.font.name = "Arial"
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0, 0, 0)
    except Exception:
        # 万一フォーマット構造が想定外でも、ここでは落とさない
        pass


def generate_problem_reframe_premise():
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    funnel_text = st.session_state.get("funnel_text", "")

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    prompt = f"""
あなたは市場調査の企画責任者です。
以下の入力を踏まえ、「真の課題」にたどり着くための前提整理として、次の5観点をそれぞれ具体的に考察してください。
この出力はユーザーが編集する前提の一次ドラフトです。

【出力形式】
次のキーを持つ JSON オブジェクト「だけ」を出力してください。
{{
  "c1_next_action": "...",
  "c2_exec_summary": "...",
  "c4_business_brand": "..."
}}

【制約】
- c1_next_action には、「調査を依頼したクライアント担当者が調査結果を受けて何を実行するか」を記述すること
- c2_exec_summary には、報告先（事業責任者・部門長・経営層）が
  「この調査結果を見て何を判断したいのか」「どの選択肢で迷っているのか」が
  明確に分かる形で記述すること。
  単なる事実確認や現状把握ではなく、意思決定に直結する論点に限定すること。
- c4_business_brand には、短期的な施策課題ではなく、
  売上・シェア・ブランド価値・顧客構造など、
  事業またはブランドの中長期的な健全性に関わる論点として記述すること。
  個別施策の良し悪しではなく、構造的な問題として表現すること。
- 各項目は60〜120字程度
- 固有名詞・前提条件・意思決定者・意思決定タイミングなど、具体情報を優先する
- 不明な場合は「不明」と書いたうえで、推定ではなく不足情報として書く
- ###、**、コードブロック記号は使わない

【入力データ】
▼オリエン統合コンテキスト（アップロード抽出＋手入力）
{ori_texts[:4000]}

▼オリエン内容の整理（抜粋）
{orien_outline_text[:2000]}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

▼マーケティングファネル（トリガー／障壁）
{funnel_text}
""".strip()

    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査の企画責任者です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.4,
            max_tokens=800,
        )
        ai_text = response.choices[0].message.content.strip()
        if ai_text.startswith("```"):
            ai_text = ai_text.strip("`")
            ai_text = ai_text.replace("json", "", 1).strip()

        obj = json.loads(ai_text)

        st.session_state["reframe_c1_next_action"] = obj.get("c1_next_action", "")
        st.session_state["reframe_c2_exec_summary"] = obj.get("c2_exec_summary", "")
        #st.session_state["reframe_c3_process_gap"] = obj.get("c3_process_gap", "")
        st.session_state["reframe_c4_business_brand"] = obj.get("c4_business_brand", "")
        #st.session_state["reframe_c5_org_mission"] = obj.get("c5_org_mission", "")
        # c6 はユーザー任意入力枠。premise生成時は空で初期化（既入力があるなら上書きしない方が良い）
        if not (st.session_state.get("reframe_c6_user_notes") or "").strip():
            st.session_state["reframe_c6_user_notes"] = ""


        return True, ""
    except Exception as e:
        return False, f"課題変換（前提整理）の生成中にエラーが発生しました: {e}"



def generate_problem_reframe():
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    # ①の編集結果を取得
    c1 = (st.session_state.get("reframe_c1_next_action") or "").strip()
    c2 = (st.session_state.get("reframe_c2_exec_summary") or "").strip()
    #c3 = (st.session_state.get("reframe_c3_process_gap") or "").strip()
    c4 = (st.session_state.get("reframe_c4_business_brand") or "").strip()
    #c5 = (st.session_state.get("reframe_c5_org_mission") or "").strip()
    c6 = (st.session_state.get("reframe_c6_user_notes") or "").strip()

    if not all([c1, c2, c4, c6]):
        return False, "①の5観点が未入力の項目があります。先に前提整理を埋めてください。"

    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    funnel_text = st.session_state.get("funnel_text", "")

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    prompt = f"""
あなたは市場調査の企画責任者です。
以下の入力と、ユーザーが編集した「5観点の前提整理」を踏まえて、
依頼窓口が表現している課題を「調査で検証可能な真の課題」へ課題変換してください。

【ユーザー編集済み：5観点の前提整理】
1) ネクストアクション：{c1}
2) 報告先が知りたいこと：{c2}
3) 事業・ブランド課題：{c4}
4) ユーザー任意追記（補足・別視点）：{c6 if c6 else "（未記入）"}

【出力形式】
次のキーを持つ JSON オブジェクト「だけ」を出力してください。
{{
      "client_stated_problem": "...",
  "true_problem": "...",
  "logic_mapping": "...",
  "rationale": "...",
  "implications_for_research": "...",
  "kickoff_impact": {{
    "ビジネス課題への反映案": "...",
    "問いへの影響": "...",
    "仮説への影響": "..."
  }}
}}


【制約】
- true_problem は80〜120字程度。市場調査で仮説検証可能な表現に限定する。
- rationale は、5観点と整合した理由を2〜4文で説明する。
- ###、**、コードブロック記号は使わない。

【入力データ】
▼オリエン統合コンテキスト（アップロード抽出＋手入力）
{ori_texts[:4000]}

▼オリエン内容の整理（抜粋）
{orien_outline_text[:2000]}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

▼マーケティングファネル（トリガー／障壁）
{funnel_text}
""".strip()

    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査の企画責任者です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.4,
            max_tokens=900,
        )
        ai_text = response.choices[0].message.content.strip()
        if ai_text.startswith("```"):
            ai_text = ai_text.strip("`")
            ai_text = ai_text.replace("json", "", 1).strip()

        obj = json.loads(ai_text)

        # ②の出力をセッションに保持（別ボックスで表示）
        st.session_state["reframe_client_stated_problem"] = obj.get("client_stated_problem", "")
        st.session_state["reframe_true_problem"] = obj.get("true_problem", "")
        st.session_state["reframe_rationale"] = obj.get("rationale", "")
        st.session_state["reframe_implications_for_research"] = obj.get("implications_for_research", "")

        ki = obj.get("kickoff_impact", {}) or {}
        st.session_state["reframe_kickoff_impact_problem"] = ki.get("ビジネス課題への反映案", "")
        st.session_state["reframe_kickoff_impact_question"] = ki.get("問いへの影響", "")
        st.session_state["reframe_kickoff_impact_hypothesis"] = ki.get("仮説への影響", "")
        st.session_state["reframe_logic_map"] = obj.get("logic_mapping", "")

        return True, ""
    except Exception as e:
        return False, f"課題変換（採用課題）の生成中にエラーが発生しました: {e}"





# =========================================================
# 調査目的マトリクス（キックオフノート用）
# =========================================================
PURPOSE_MATRIX = {
    "市場・競合把握": "市場規模、成長性、競合構造などの理解を目的とした調査です。",
    "ニーズボリューム把握": "消費者ニーズの量的分布を明らかにし、優先ターゲットを特定します。",
    "実態・意識把握": "消費者の行動実態や意識構造を明らかにする調査です。",
    "ニーズ探索": "潜在的な消費者ニーズやウォンツを発掘・探索します。",
    "アイデアスクリーニング": "複数のアイデア案を評価・選抜するための初期テストを行います。",
    "コンセプト受容性把握": "商品・サービスコンセプトの受容度、共感度、理解度を測定します。",
    "スぺック評価把握": "製品スペック（機能・特徴）の重要度や評価ポイントを明らかにします。",
    "価格弾力性把握": "価格設定に対する需要反応（価格弾力性）を推定します。",
    "需要予測": "市場シェアや販売量の見込みを予測する調査です。",
    "訴求ポイント把握": "広告・コミュニケーションで強調すべきメッセージを明確化します。",
    "浸透状況把握": "ブランド・製品の市場浸透率や認知度、利用率を測定します。",
    "サービス使用評価": "実際のサービス利用体験を通じた満足度・課題を抽出します。",
    "プロモ効果測定": "キャンペーンやプロモーションの効果を定量的に評価します。",
    "ユーザー評価": "既存ユーザーからの製品・サービス評価を把握します。",
     "（該当なし／その他）": "下の自由記述欄に今回の調査目的を具体的に記載してください。",
}
# PURPOSE_MATRIX 定義の直後に追加
if st.session_state.get("kickoff_selected_purpose") not in PURPOSE_MATRIX:
    st.session_state["kickoff_selected_purpose"] = next(iter(PURPOSE_MATRIX))




import uuid
from datetime import datetime

def _now_id(prefix: str = "rev") -> str:
    return f"{prefix}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"

REV_STORE_KEYS = ["proposal_revisions", "revisions"]  # どちらでも動くようにする

def _get_revision_store_key() -> str:
    for k in REV_STORE_KEYS:
        if k in st.session_state:
            return k
    # 何もなければ proposal_revisions を作る
    st.session_state["proposal_revisions"] = []
    return "proposal_revisions"

def get_revisions() -> list[dict]:
    k = _get_revision_store_key()
    return st.session_state.get(k, []) or []

def find_revision(rev_id: str) -> dict | None:
    for r in get_revisions():
        if r.get("rev_id") == rev_id:
            return r
    return None

def set_active_revision(rev_id: str | None) -> None:
    st.session_state["active_rev_id"] = rev_id

def get_active_revision() -> dict | None:
    rid = st.session_state.get("active_rev_id")
    if not rid:
        return None
    return find_revision(rid)

def ensure_revision_store() -> None:
    if "proposal_revisions" not in st.session_state or st.session_state["proposal_revisions"] is None:
        st.session_state["proposal_revisions"] = []

def append_revision(
    stage: str,
    purpose_key: str,
    kickoff: dict,
    subq_list: list | None = None,
    subquestions_raw: str = "",
    analysis_blocks: list | None = None,
    target_condition: str = "",
    survey_items: dict | None = None,
    spec: dict | None = None,
    label: str | None = None,
    parent_rev_id: str | None = None,
    notes: str = "",
) -> dict:
    ensure_revision_store()

    rev_id = _now_id("rev")
    created_at = datetime.now().isoformat(timespec="seconds")

    # ★フル生成の軸（課題ピボット）をRevisionに保存する
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = st.session_state.get("fullgen_axis_source", "pivot")

    # purpose_key は「比較表示・ラベル用の識別子」に寄せる（互換のため引数は維持）
    # もし呼び出し側で purpose_key を渡さない運用にするなら、ここで強制上書きしてもよい。
    effective_purpose_key = purpose_key or "PIVOT"

    # PURPOSE_MATRIX を排除するので、purpose_desc は軸テキストから作る
    # 長すぎると比較UIが見づらいので短縮
    def _shorten(s: str, n: int = 220) -> str:
        s = (s or "").strip()
        return s if len(s) <= n else s[:n] + "…"

    purpose_desc = _shorten(axis_text, 220) if axis_text else ""

    if not label:
        label = f"{len(st.session_state['proposal_revisions'])+1}. {stage.upper()} / {effective_purpose_key}"

    rev = {
        "rev_id": rev_id,
        "label": label,
        "created_at": created_at,
        "stage": stage,

        # 既存UI互換のため残す
        "purpose_key": effective_purpose_key,
        "purpose_desc": purpose_desc,

        # PURPOSE_MATRIX由来の自由記述も不要なら空固定でOK
        # ただし既存UI/ロジックが参照しているならキー自体は残す
        "purpose_free": "",

        "axis_source": axis_source,
        "axis_text": axis_text,

        # ★新設：軸を明示的に保存（再現性のため）
        "fullgen_axis": {
            "source": axis_source,
            "text": axis_text,
        },

        "parent_rev_id": parent_rev_id,
        "notes": notes,
        "orien": {
            "orien_outline_text": st.session_state.get("orien_outline_text", ""),
            "orien_outline_ai_draft": st.session_state.get("orien_outline_ai_draft", ""),
            "orien_outline_manual": st.session_state.get("orien_outline_manual", ""),
            "uploaded_docs": st.session_state.get("uploaded_docs", []) or [],
        },
        "kickoff": {
            "目標": kickoff.get("目標", ""),
            "現状": kickoff.get("現状", ""),
            "ビジネス課題": kickoff.get("ビジネス課題", ""),
            "調査目的": kickoff.get("調査目的", ""),
            "問い": kickoff.get("問い", ""),
            "仮説": kickoff.get("仮説", ""),
            "ポイント": kickoff.get("ポイント", ""),
        },
        "subq_list": subq_list or [],
        "subquestions_raw": subquestions_raw or "",
        "analysis_blocks": analysis_blocks or [],
        "target_condition": target_condition or "",
        "survey_item_rows": st.session_state.get("survey_item_rows", []) or [],
        "survey_items": survey_items or {},
        "spec": spec or {},

        # ★追加：課題変換（前処理）を Revision に必ず紐づけ
        "problem_reframe": {
            "c1_next_action": st.session_state.get("reframe_c1_next_action", ""),
            "c2_exec_summary": st.session_state.get("reframe_c2_exec_summary", ""),
            "c4_business_brand": st.session_state.get("reframe_c4_business_brand", ""),
            "c6_user_notes": st.session_state.get("reframe_c6_user_notes", ""),
            "output": st.session_state.get("problem_reframe_output", {}) or {},
        },
        "true_problem_text": st.session_state.get("true_problem_text", "") or "",
    }

    st.session_state["proposal_revisions"].append(rev)
    set_active_revision(rev_id)

    st.session_state["revisions_version"] = st.session_state.get("revisions_version", 0) + 1
    return rev


def save_session_keys_to_active_revision() -> None:
    rev = get_active_revision()
    if not rev:
        return
    rev.setdefault("kickoff", {})

    # 企画書本文（下書き）を保存
    rev["purpose_free"] = st.session_state.get("kickoff_purpose_free_editor", "") or ""


    # rev["orien"] = {
    #     "orien_outline_text": st.session_state.get("orien_outline_text", ""),
    #     "orien_outline_ai_draft": st.session_state.get("orien_outline_ai_draft", ""),
    #     "orien_outline_manual": st.session_state.get("orien_outline_manual", ""),
    #     "uploaded_docs": st.session_state.get("uploaded_docs", []) or [],
    # }

    # 1) kickoff
    rev["kickoff"]["目標"] = st.session_state.get("ai_目標", "")
    rev["kickoff"]["現状"] = st.session_state.get("ai_現状", "")
    rev["kickoff"]["ビジネス課題"] = st.session_state.get("ai_ビジネス課題", "")
    rev["kickoff"]["調査目的"] = st.session_state.get("ai_調査目的", "")
    rev["kickoff"]["問い"] = st.session_state.get("ai_問い", "")
    rev["kickoff"]["仮説"] = st.session_state.get("ai_仮説", "")
    if "ai_ポイント" in st.session_state:
        rev["kickoff"]["ポイント"] = st.session_state.get("ai_ポイント", "")

    # 2) 問いの分解
    rev["subq_list"] = st.session_state.get("subq_list", []) or []
    rev["subquestions_raw"] = st.session_state.get("ai_subquestions", "") or ""


    # --------------------
    # 3) 分析アプローチ（保存のみ：session_stateは破壊しない）
    # --------------------
    base_blocks = st.session_state.get("analysis_blocks", []) or []

    # ★追加：subq_list（正）を取得しておく
    ss_subq_list = st.session_state.get("subq_list", []) or []
    ss_subq_texts = []
    for sq in ss_subq_list:
        d = dict(sq or {})
        txt = (d.get("subq") or d.get("question") or "").strip()
        ss_subq_texts.append(txt)

    blocks = []
    for i, b0 in enumerate(base_blocks, 1):
        # UI入力があれば優先、なければ既存ブロック
        subq = st.session_state.get(f"analysis_subq_{i}", b0.get("subq", ""))
        axis = st.session_state.get(f"analysis_axis_{i}", b0.get("axis", ""))
        items = st.session_state.get(f"analysis_items_{i}", b0.get("items", ""))
        approach = st.session_state.get(f"analysis_approach_{i}", b0.get("approach", ""))
        hypothesis = st.session_state.get(f"analysis_hypothesis_{i}", b0.get("hypothesis", ""))

        # ★重要：subq_list があれば、analysis側のsubqを追随させる（ズレ防止）
        # ただし、analysis_subq_i に明示入力があるならそれを尊重したい場合は条件を変えてOK
        if i <= len(ss_subq_texts) and ss_subq_texts[i - 1]:
            subq = ss_subq_texts[i - 1]

        # 完全空は落とす
        if not any([str(subq).strip(), str(axis).strip(), str(items).strip(), str(approach).strip(), str(hypothesis).strip()]):
            continue

        bid = (b0.get("id") or "").strip() or f"SQ{i}"
        blocks.append({
            "id": bid,
            "subq": subq,
            "axis": axis,
            "items": items,
            "approach": approach,
            "hypothesis": hypothesis,
        })

    rev["analysis_blocks"] = blocks




    # 4) 対象者条件案
    rev["target_condition"] = st.session_state.get("ai_target_condition", "") or ""

    # 5') 調査項目（analysis_blocks連動の表）
    rev["survey_item_rows"] = st.session_state.get("survey_item_rows", []) or []
    st.session_state["survey_item_rows"] = rev["survey_item_rows"]

    # 5) 調査項目案
    si = st.session_state.get("ai_survey_items", {}) or {}
    for ver in ["10問", "20問", "30問", "40問"]:
        ui_key = f"survey_items_{ver}"
        if ui_key in st.session_state:
            si[ver] = st.session_state.get(ui_key, "")
    rev["survey_items"] = si
    st.session_state["ai_survey_items"] = si

    # 6) 調査仕様案
    rev["spec"] = {label: st.session_state.get(ss_key, "") for (label, ss_key) in SPEC_ITEMS}

    # 7) 課題変換（前処理）
    rev["problem_reframe"] = {
        "c1_next_action": st.session_state.get("reframe_c1_next_action", ""),
        "c2_exec_summary": st.session_state.get("reframe_c2_exec_summary", ""),
        "c4_business_brand": st.session_state.get("reframe_c4_business_brand", ""),
        "c6_user_notes": st.session_state.get("reframe_c6_user_notes", ""),
        "output": st.session_state.get("problem_reframe_output", {}) or {},
    }
    rev["true_problem_text"] = st.session_state.get("true_problem_text", "") or ""


    st.session_state["__dbg_before_upsert"] = {
        "rev_id": rev.get("rev_id"),
        "len_rev_analysis_blocks": len((rev.get("analysis_blocks") or [])),
        "len_ss_analysis_blocks": len((st.session_state.get("analysis_blocks") or [])),
    }

    # ★最後に1回だけ保存
    upsert_revision(rev)

    rev3 = find_revision(rev.get("rev_id"))
    st.session_state["__dbg_after_upsert"] = {
        "found": bool(rev3),
        "len_found_analysis_blocks": len(((rev3 or {}).get("analysis_blocks") or [])),
    }



def upsert_revision(updated: dict) -> None:
    k = _get_revision_store_key()
    revs = st.session_state.get(k, []) or []

    rid = updated.get("rev_id")
    if not rid:
        raise ValueError("upsert_revision: rev_id がありません。")

    new_revs = []
    replaced = False
    for r in revs:
        if r.get("rev_id") == rid:
            new_revs.append(updated)  # 置換
            replaced = True
        else:
            new_revs.append(r)

    if not replaced:
        new_revs.append(updated)

    st.session_state[k] = new_revs




def apply_revision_to_session(rev: dict) -> None:
    if not rev:
        return

    # --------------------
    # Streamlit安全代入ヘルパー
    # （ウィジェット生成後の上書きによるクラッシュ防止）
    # --------------------
    def _safe_set(key: str, value):
        # すでに同名keyのウィジェットがこのrunで生成されている場合は触らない
        if key in st.session_state:
            return
        st.session_state[key] = value


    # ウィジェットkeyは触らず editor側に反映し、保存用にも同期
    st.session_state["kickoff_purpose_free_editor"] = rev.get("purpose_free", "") or ""
    st.session_state["kickoff_purpose_free"] = st.session_state["kickoff_purpose_free_editor"]

    # 0) orien（前提）を復元
    orien = rev.get("orien", {}) or {}

    # uploaded_docs は revision 側が空でも上書きしない（現セッション優先）
    if "uploaded_docs" in orien and orien.get("uploaded_docs"):
        st.session_state["uploaded_docs"] = orien.get("uploaded_docs") or []

    # ai_draft/manual/text は「値があるときだけ」上書きする
    val = (orien.get("orien_outline_ai_draft") or "").strip()
    if val:
        st.session_state["orien_outline_ai_draft"] = val

    val = (orien.get("orien_outline_manual") or "").strip()
    if val:
        st.session_state["orien_outline_manual"] = val

    val = (orien.get("orien_outline_text") or "").strip()
    if val:
        st.session_state["orien_outline_text"] = val



    # --------------------
    # 1) kickoff（ai_*）
    # --------------------
    k = rev.get("kickoff", {}) or {}
    st.session_state["ai_目標"] = k.get("目標", "")
    st.session_state["ai_現状"] = k.get("現状", "")
    st.session_state["ai_ビジネス課題"] = k.get("ビジネス課題", "")
    st.session_state["ai_調査目的"] = k.get("調査目的", "")
    st.session_state["ai_問い"] = k.get("問い", "")
    st.session_state["ai_仮説"] = k.get("仮説", "")
    if "ai_ポイント" in st.session_state or k.get("ポイント"):
        st.session_state["ai_ポイント"] = k.get("ポイント", "")

    # --------------------
    # 2) 問いの分解
    # --------------------
    _subq_list = rev.get("subq_list", []) or []
    norm_list = []
    for sq in _subq_list:
        d = dict(sq or {})
        txt = (d.get("subq") or d.get("question") or "").strip()
        d["subq"] = txt
        d["question"] = txt
        norm_list.append(d)

    st.session_state["subq_list"] = norm_list
    st.session_state["ai_subquestions"] = rev.get("subquestions_raw", "") or ""

    # 表示キャッシュは切替時にクリア（残像対策）
    st.session_state.pop("subq_structured_view", None)
    st.session_state.pop("EDIT1_subQ", None)

    # --------------------
    # 3) 分析アプローチ
    #   - ここが本丸：subq_list（正）を analysis_blocks / analysis_subq_i に反映させる
    # --------------------
    MAX_I = 9

    blocks = (rev.get("analysis_blocks") or []) or []
    blocks2 = []

    # 既存blocksがあるなら subq を subq_list に追随させる
    if blocks:
        for idx, b in enumerate(blocks):
            b2 = dict(b or {})
            if idx < len(norm_list):
                b2["subq"] = norm_list[idx].get("subq", "") or ""
            blocks2.append(b2)
    else:
        # 既存blocksが無い場合でも、analysis_subq_i が空にならないように
        # ここでは blocks2 は空のまま（生成は別ボタン）でOK。UIキーだけ入れる。
        blocks2 = []

    # 正規化済み blocks を session に保存
    st.session_state["analysis_blocks"] = blocks2

    # UIが参照している analysis_*_{i} を上書き（残像を消す）
    for i in range(1, MAX_I + 1):
        if i <= len(blocks2):
            b = blocks2[i - 1] or {}
            st.session_state[f"analysis_subq_{i}"] = b.get("subq", "") or ""
            st.session_state[f"analysis_axis_{i}"] = b.get("axis", "") or ""
            st.session_state[f"analysis_items_{i}"] = b.get("items", "") or ""
            st.session_state[f"analysis_approach_{i}"] = b.get("approach", "") or ""
            st.session_state[f"analysis_hypothesis_{i}"] = b.get("hypothesis", "") or ""
        else:
            st.session_state[f"analysis_subq_{i}"] = (norm_list[i - 1]["subq"] if i <= len(norm_list) else "")
            st.session_state[f"analysis_axis_{i}"] = ""
            st.session_state[f"analysis_items_{i}"] = ""
            st.session_state[f"analysis_approach_{i}"] = ""
            st.session_state[f"analysis_hypothesis_{i}"] = ""

    # --------------------
    # 4) 対象者条件案
    # --------------------
    st.session_state["ai_target_condition"] = rev.get("target_condition", "") or ""
    # editor側のキーも合わせて上書き（UI残像対策）
    st.session_state["ai_target_condition_editor"] = st.session_state["ai_target_condition"]
    st.session_state.pop("target_condition_textarea", None)  # 旧キャッシュ使ってる場合の保険

    # 5') 調査項目（analysis_blocks連動の表）
    st.session_state["survey_item_rows"] = rev.get("survey_item_rows", []) or []

    # --------------------
    # 5) 調査項目案（10/20/30/40）
    # --------------------
    si = rev.get("survey_items", {}) or {}
    st.session_state["ai_survey_items"] = si
    for ver in ["10問", "20問", "30問", "40問"]:
        st.session_state[f"survey_items_{ver}"] = si.get(ver, "")

    # --------------------
    # 6) 調査仕様案
    # --------------------
    spec_obj = rev.get("spec", {}) or {}
    # SPEC_ITEMS が定義済み前提（あなたのコード通り）
    for label, ss_key in SPEC_ITEMS:
        val = spec_obj.get(label, "")
        if not isinstance(val, str):
            val = json.dumps(val, ensure_ascii=False) if val else ""
        st.session_state[ss_key] = val

    # 7) 課題変換（前処理）
    pr = rev.get("problem_reframe", {}) or {}

    _safe_set("reframe_c1_next_action", pr.get("c1_next_action", ""))
    _safe_set("reframe_c2_exec_summary", pr.get("c2_exec_summary", ""))
    _safe_set("reframe_c4_business_brand", pr.get("c4_business_brand", ""))
    _safe_set("reframe_c6_user_notes", pr.get("c6_user_notes", ""))

    _safe_set("problem_reframe_output", pr.get("output", {}) or {})
    _safe_set("true_problem_text", rev.get("true_problem_text", "") or "")
    _safe_set(
        "problem_reframe_generated",
        bool(pr.get("output", {}) or {})
    )



def request_apply_revision(rev_id: str):
    st.session_state["pending_apply_rev_id"] = rev_id



def append_default_revision_from_current_ai(stage: str = "default") -> dict:
    """
    現在の session_state の ai_*（キックオフ）を「デフォルトRevision」として履歴に追加する。
    - 既に default が存在する場合は追加しない（重複防止）
    """
    ensure_revision_store()

    # 既に default を作っているなら二重追加しない
    for r in get_revisions():
        if r.get("stage") == "default":
            return r

    purpose_key = st.session_state.get("kickoff_selected_purpose") or list(PURPOSE_MATRIX.keys())[0]

    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
        "ポイント": st.session_state.get("ai_ポイント", ""),
    }

    rev = append_revision(
        stage=stage,  # "default"
        purpose_key=purpose_key,
        kickoff=kickoff,
        subq_list=st.session_state.get("subq_list", []) or [],
        label=f"0. DEFAULT / {purpose_key}",
        parent_rev_id=None,
        notes="初回『企画書下書きを生成』で作成されたベース案",
    )
    return rev





def sync_active_revision_to_session_keys() -> None:
    """
    active_rev_id が指す Revision を取り出し、
    proposal_draft UI が参照する session_state キーへ展開する。
    """
    rev = get_active_revision()
    if not rev:
        return
    apply_revision_to_session(rev)


def run_full_generation_and_append_revision(axis_text: str, axis_source: str = "pivot") -> tuple[bool, str]:
    """
    フル生成（①〜⑦）→ session_state に成果物を格納 → kickoff を回収 → Revision追加

    改修：PURPOSE_MATRIX を廃止し、課題ピボット（6観点連結テキスト）をフル生成の軸として扱う。
    """

    # ★残像対策：フル生成開始時に下流成果物をクリア
    st.session_state["ai_subquestions"] = ""
    st.session_state["subq_list"] = []
    st.session_state["analysis_blocks"] = []


    # ---------------------------------------------------------
    # 軸テキストの妥当性チェック
    # ---------------------------------------------------------
    axis_text = (axis_text or "").strip()
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    # ---------------------------------------------------------
    # 下流生成関数が参照できるように session_state に保存
    # generate_kickoff_draft / generate_subquestions_draft 等が
    # st.session_state['fullgen_axis_text'] を読む前提
    # ---------------------------------------------------------
    st.session_state["fullgen_axis_source"] = axis_source
    st.session_state["fullgen_axis_text"] = axis_text

    # ① オリエン内容の整理
    ok, msg = ensure_orien_outline()
    if not ok:
        return False, msg

    # ② キックオフノート
    ok, msg = generate_kickoff_draft()
    if not ok:
        return False, msg

    # ③ 問いの分解（subq_list を作る）
    ok, msg = generate_subquestions_draft()
    if not ok:
        return False, msg

    # ★追加：生成したのに parse 結果が空のケースを検知
    if not (st.session_state.get("subq_list") or []):
        raw = (st.session_state.get("ai_subquestions") or "")[:1200]
        return False, f"サブクエスチョンの解析結果が空でした。出力形式が崩れている可能性があります。\n---\n{raw}"


    # ④ 分析アプローチ
    ok, msg = generate_analysis_approach_draft()
    if not ok:
        return False, msg

    # ⑤ 対象者条件（先に作る：⑤’の参考に入れる）
    ok, msg = generate_target_condition_draft()
    if not ok:
        return False, msg

    # ⑤' 調査項目（analysis_blocks連動）
    ok, msg = generate_survey_items_linked_draft()
    if not ok:
        return False, msg

    # ⑤（互換）10/20/30/40 を rows から再構成
    rows = st.session_state.get("survey_item_rows", []) or []
    st.session_state["ai_survey_items"] = build_survey_versions_from_rows(rows)
    st.session_state["ai_survey_items_raw"] = ""  # 旧rawを使わないなら空でOK

    # ⑦ 調査仕様案
    ok, msg = generate_spec_draft()
    if not ok:
        return False, msg

    st.session_state["proposal_draft_generated"] = True

    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
        "ポイント": st.session_state.get("ai_ポイント", ""),
    }

    rev_index = len(get_revisions()) + 1
    rev_label = f"{rev_index}. FULL / axis:{axis_source}"

    # spec のキー名衝突を回避
    spec_dict = {spec_label: st.session_state.get(ss_key, "") for (spec_label, ss_key) in SPEC_ITEMS}

    # purpose_key は互換用に固定値を入れておくのが安全（UIや既存ロジック対策）
    # 例："PIVOT" / "AXIS" / "pivot" など
    rev = append_revision(
        stage="full",
        purpose_key="PIVOT",
        kickoff=kickoff,
        subq_list=st.session_state.get("subq_list", []) or [],
        subquestions_raw=st.session_state.get("ai_subquestions", "") or "",
        analysis_blocks=st.session_state.get("analysis_blocks", []) or [],
        target_condition=st.session_state.get("ai_target_condition", "") or "",
        survey_items=st.session_state.get("ai_survey_items", {}) or {},
        spec=spec_dict,
        label=rev_label,
    )

    request_apply_revision(rev["rev_id"])
    return True, ""




def run_phaseA_generation_and_append_revision(axis_text: str, axis_source: str = "pivot") -> tuple[bool, str]:
    """
    PhaseA: KON（キックオフノート）〜 サブクエスチョンまでを生成して Revision に追加する。
    - 生成・比較タブの高速比較用（分析アプローチ以降はここでは作らない）
    """

    # ★残像対策：PhaseA開始時に下流成果物をクリア（比較の混入を防止）
    st.session_state["ai_subquestions"] = ""
    st.session_state["subq_list"] = []
    st.session_state["analysis_blocks"] = []            # PhaseAでは作らない
    st.session_state["ai_target_condition"] = ""        # PhaseAでは作らない
    st.session_state["survey_item_rows"] = []           # PhaseAでは作らない
    st.session_state["ai_survey_items"] = {}            # PhaseAでは作らない
    st.session_state["ai_survey_items_raw"] = ""        # PhaseAでは作らない
    st.session_state["proposal_draft_generated"] = False  # PhaseA完了=フル生成完了ではない

    # ---------------------------------------------------------
    # 軸テキストの妥当性チェック
    # ---------------------------------------------------------
    axis_text = (axis_text or "").strip()
    if not axis_text:
        return False, "生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    # ---------------------------------------------------------
    # 下流生成関数が参照できるように session_state に保存
    # generate_kickoff_draft / generate_subquestions_draft が参照
    # ---------------------------------------------------------
    st.session_state["fullgen_axis_source"] = axis_source
    st.session_state["fullgen_axis_text"] = axis_text

    # ① オリエン内容の整理
    ok, msg = ensure_orien_outline()
    if not ok:
        return False, msg

    # ② キックオフノート
    ok, msg = generate_kickoff_draft()
    if not ok:
        return False, msg

    # ③ 問いの分解（subq_list を作る）
    ok, msg = generate_subquestions_draft()
    if not ok:
        return False, msg

    # ★追加：生成したのに parse 結果が空のケースを検知
    if not (st.session_state.get("subq_list") or []):
        raw = (st.session_state.get("ai_subquestions") or "")[:1200]
        return False, f"サブクエスチョンの解析結果が空でした。出力形式が崩れている可能性があります。\n---\n{raw}"

    # ---------------------------------------------------------
    # Revision 保存（PhaseA）
    # ---------------------------------------------------------
    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
        "ポイント": st.session_state.get("ai_ポイント", ""),
    }

    rev_index = len(get_revisions()) + 1
    rev_label = f"{rev_index}. PHASEA / axis:{axis_source}"

    # spec は空でも良いが、既存ロジック互換のためキーだけ揃える
    spec_dict = {spec_label: "" for (spec_label, _ss_key) in SPEC_ITEMS}

    rev = append_revision(
        stage="phaseA",          # ★比較タブ側でも拾えるようにしておく
        purpose_key="PIVOT",     # ★互換用に固定（fullと同じ）
        kickoff=kickoff,
        subq_list=st.session_state.get("subq_list", []) or [],
        subquestions_raw=st.session_state.get("ai_subquestions", "") or "",
        analysis_blocks=[],      # PhaseAでは未生成
        target_condition="",     # PhaseAでは未生成
        survey_items={},         # PhaseAでは未生成
        spec=spec_dict,
        label=rev_label,
        notes="PhaseA（KON〜SQ）生成。分析アプローチ以降は編集タブで実行。",
    )

    request_apply_revision(rev["rev_id"])
    return True, ""

# def _render_kon_sq_compact_view(rev: dict | None, title: str, key_prefix: str):
#     st.markdown(f"**{title}：KONの問い → サブクエスチョン（構造表示）**")
#     if not rev:
#         st.info("Revisionが未選択、または取得できませんでした。")
#         return

#     rev_id = rev.get("rev_id") or "no_rev"
#     kickoff = (rev.get("kickoff") or {})
#     kon_q = (kickoff.get("問い") or "").strip()

#     subq_list = (rev.get("subq_list") or []) or []

#     lines = []
#     lines.append("【KON：問い】")
#     lines.append(kon_q if kon_q else "（未設定）")
#     lines.append("")
#     lines.append("【SQ：サブクエスチョン】")

#     if not subq_list:
#         lines.append("（サブクエスチョンなし）")
#     else:
#         for i, sq in enumerate(subq_list, 1):
#             sq_text = (sq.get("subq") or sq.get("question") or "").strip()
#             lines.append(f"SQ{i}: {sq_text}")

#     text = "\n".join(lines).strip()

#     widget_key = f"{key_prefix}__{rev_id}__kon_sq_view"
#     st.session_state[widget_key] = text  # 表示専用なので毎回上書きでOK

#     st.text_area(
#         "KON〜SQ（構造表示）",
#         key=widget_key,
#         height=320,
#         disabled=True,  # ←まずは表示統一だけ。保存拡張は後で
#     )





import time, uuid

def create_new_revision():
    rid = f"rev_{time.strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"
    return {
        "rev_id": rid,
        "label": f"Revision {time.strftime('%Y/%m/%d %H:%M:%S')}",
        "created_at": time.time(),
        "kickoff": {"目標":"","現状":"","ビジネス課題":"","調査目的":"","問い":"","仮説":"","ポイント":""},
        "subq_list": [],
        "subquestions_raw": "",
        "analysis_blocks": [],
        "target_condition": "",
        "survey_item_rows": [],
        "survey_items": {},
        "spec": {},
        "problem_reframe": {},
        "true_problem_text": "",
        "orien": {},
        "purpose_free": "",
    }











# =========================================================
# 企画書下書き：各ステップを AI で一括生成するヘルパー
# =========================================================
import json
import pandas as pd


def ensure_orien_outline():
    docs_text = "\n".join(st.session_state.get("uploaded_docs", []) or []).strip()
    manual_text = (st.session_state.get("orien_outline_manual") or "").strip()

    # アップロードも手入力も無ければ生成不可
    if not docs_text and not manual_text:
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    # AIには、資料抽出テキストを主に渡す（手入力は任意で補助）
    # ※「手入力最優先」を崩さないため、手入力をAIに渡すかは方針次第。
    #   ここでは補助として渡す（ただし後段統合では手入力優先）
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    prompt = f"""
あなたは市場調査の専門家です。
以下のオリエン資料から以下のことをまとめてください。
特に言及がなければ項目ごとに「なし」と記載してください。

【出力形式】
・企業名：
・ブランド名：
・カテゴリー（市場）名：
・議事録の要約（500文字程度）：
・分析手法に関する要望：
・調査仕様に関する要望
    調査エリア：
    スクリーニング調査有無：
    対象者条件：
    質問数：
    サンプルサイズ：
    調査画面で画像や動画の提示：
    ウェイトバック集計の有無：
    自由回答のコーディング処理の有無：
    調査票作成（クライアントがやるか当社がやるか）：
    報告書は必要か：
・スケジュールに関する要望
    企画提案予定日：
    調査票や画像に関する提供可能日：
    希望する納期：
    請求日/月：
    クライアントの重要な会議日：
    その他スケジュールに関する要望：
・費用に関する要望
    見積金額上限：
    複数パターンの見積を希望しているか：
・会議参加者のお名前・役職・役割
・調査とは直接関係ないが雑談したこと：
・その他調査に関する特記事項（広告がいつから投下されるかなど）：

オリエン資料：
{ori_texts[:4000]}
"""

    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
            max_tokens=900,
        )
        ai_result = response.choices[0].message.content.strip()

        # ★必ずAI整理結果として保持（UIの整理欄はこれを見る）
        # AI整理結果は永続キーに保存
        st.session_state["data_orien_outline_ai_draft"] = ai_result

        # UIが存在するタイミングならUIにも反映（任意だが推奨）
        st.session_state["ui_orien_outline_ai_draft"] = ai_result

        # 正本（手入力優先）
        manual_text = (st.session_state.get("data_orien_outline_manual") or "").strip()
        st.session_state["orien_outline_text"] = manual_text if manual_text else ai_result


        return True, ""
    except Exception as e:
        return False, f"オリエン内容の整理生成中にエラーが発生しました: {e}"


def generate_kickoff_draft():
    """キックオフノート（①〜⑥）を AI で生成（フル生成の軸は課題ピボット6観点を使用）"""
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    funnel_text = st.session_state.get("funnel_text", "")

    # ---------------------------------------------------------
    # ★フル生成の軸：課題ピボット（6観点）を使用（PURPOSE_MATRIX依存を廃止）
    # ---------------------------------------------------------
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = st.session_state.get("fullgen_axis_source", "pivot")

    if not axis_text:
        return (
            False,
            "フル生成の軸テキスト（課題ピボット6観点）が見つかりません。先に『課題ピボット』で生成・編集してください。",
        )

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    true_problem = (st.session_state.get("true_problem_text") or "").strip()
    problem_reframe_obj = st.session_state.get("problem_reframe_output", {}) or {}

    prompt = f"""
あなたは市場調査設計の専門家です。
以下のオリエン資料、ブランド診断結果、および課題ピボット（6観点）をもとに、
調査設計の初期段階で用いる「キックオフノート」を作成してください。

【出力形式】
【目標】
【現状】
【ビジネス課題】
【調査目的】
【問い】
【仮説】
【ポイント】


【条件】
- 各項目は80〜120字以内
- オリエン資料にある固有名詞や文脈を十分に生かしてください。
- 【目標】や【現状】は経営課題や社会問題など、調査では解決できない抽象課題は避けてください。
  あくまで「消費者・市場・ブランド・広告・顧客体験」など、市場調査で仮説検証できる範囲に課題を限定してください。
- 【問い】はオリエンシートやブランド診断を踏まえた現在の対象ブランドの"リサーチクエスチョン"のことです。
- 【ポイント】にはなぜキックオフノートの各項目にそう記載したのか、特に注意すべき点や補足説明を簡潔に記載してください。
- 【ビジネス課題】は、課題変換（前処理）の採用課題（true_problem）がある場合は参照してもよいが、
  まずは「課題ピボット（6観点）」を最優先の前提として要約・反映してください。
- ###、** などの記号は使わないでください。

【入力データ】
▼オリエン内容の整理（抜粋）
{orien_outline_text[:2000]}

▼課題ピボット（フル生成の軸：最優先で反映する／任意追記含む）
（source: {axis_source}）
{axis_text}

▼課題変換（前処理）の結果（採用課題：補助情報。必要に応じて参照）
{true_problem if true_problem else "（未実施）"}

▼課題変換（前処理）の補足（必要に応じて参照）
{json.dumps(problem_reframe_obj, ensure_ascii=False)}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

▼マーケティングファネル
{funnel_text}
"""
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.6,
            max_tokens=900,
        )
        result = response.choices[0].message.content
        sections = parse_ai_output(result)

        for key in sections:
            st.session_state[f"ai_{key}"] = sections[key]

        return True, ""
    except Exception as e:
        return False, f"キックオフノート生成中にエラーが発生しました: {e}"

def generate_subquestions_draft():
    """問いの分解（サブクエスチョン）を AI で JSON 生成（MAX_I=9固定）"""
    ori_texts = get_orien_context_text()
    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    main_question = st.session_state.get("ai_問い", "")

    # ★フル生成の軸（課題ピボット）
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()

    if not main_question.strip():
        return False, "キックオフノート⑤『問い』がまだ生成されていません。"
    if not ori_texts.strip():
        return False, "オリエン資料が未入力です（アップロード or 手入力が必要です）。"
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    MAX_I = 9

    prompt = f"""
あなたは市場調査設計の専門家です。
以下の情報をもとに、キックオフノート⑤『問い』（リサーチクエスチョン）を深掘りするための
サブクエスチョン案を作成してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（6観点）」は、依頼課題を調査で検証可能な採用課題（真の課題）へ変換したものです。
- サブクエスチョンは、必ずこの軸に整合すること（論点がずれないこと）。

▼軸情報（source: {axis_source}）
{axis_text}

【キックオフノート⑤ 問い】
{main_question}

▼オリエン統合コンテキスト（アップロード抽出＋手入力）
{ori_texts[:4000]}

▼オリエン内容の整理（抜粋）
{orien_outline_text[:2000]}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

【出力要件】
- 出力は「JSON配列のみ」。前後に説明文、コードブロック、見出しは一切出さない。
- 要素数は最大{MAX_I}件。
- 各要素は次のキーを必ず持つこと：
  - id: "SQ1" のようなID（連番）
  - subq: サブクエスチョン本文（1文）
  - axis: 分析軸（セグメント）案（簡潔に）
  - items: 評価項目案（簡潔に）

【JSON出力例（この形のまま返す）】
[
  {{
    "id": "SQ1",
    "subq": "・・・？",
    "axis": "・・・",
    "items": "・・・"
  }}
]
""".strip()

    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.6,
            max_tokens=2000,
        )

        ai_text = (response.choices[0].message.content or "").strip()

        # ```json 対策
        if ai_text.startswith("```"):
            ai_text = ai_text.strip("`")
            ai_text = ai_text.replace("json", "", 1).strip()

        # ★raw保存（失敗解析のため必須）
        st.session_state["ai_subquestions"] = ai_text

        # JSONとして解釈
        try:
            items = json.loads(ai_text)
        except json.JSONDecodeError:
            return False, "サブクエスチョンがJSON形式で返ってきませんでした。プロンプト出力形式が崩れている可能性があります。"

        if not isinstance(items, list):
            return False, "サブクエスチョンのJSONが配列ではありません。"

        # 空要素の除外 → MAX_I適用
        clean = []
        for i, it in enumerate(items, 1):
            if not isinstance(it, dict):
                continue
            subq = (it.get("subq") or "").strip()
            if not subq:
                continue

            clean.append({
                "id": (it.get("id") or f"SQ{i}").strip(),
                "subq": subq,
                "axis": (it.get("axis") or "").strip(),
                "items": (it.get("items") or "").strip(),
            })
            if len(clean) >= MAX_I:
                break

        if not clean:
            return False, "サブクエスチョンの解析結果が空でした（JSONは取れたが有効なsubqがありません）。"

        st.session_state["subq_list"] = clean
        return True, ""

    except Exception as e:
        return False, f"問いの分解（サブクエスチョン）生成中にエラーが発生しました: {e}"


def generate_analysis_approach_draft():
    """サブクエスチョン別の分析アプローチ案を AI で JSON 生成（active revision を唯一の正とする）"""

    # =========================================================
    # 0) active revision を確定（生成の正）
    # =========================================================
    rid = st.session_state.get("active_rev_id")
    if not rid:
        return False, "active_rev_id が未設定です。"

    rev = find_revision(rid)
    if not rev:
        return False, f"active_rev_id={rid} のRevisionが見つかりません。"

    # =========================================================
    # 1) 必須入力チェック
    # =========================================================
    ori_texts = get_orien_context_text()
    if not (ori_texts or "").strip():
        return False, "オリエン資料が未入力です（アップロード or 手入力が必要です）。"

    # ★追加：フル生成の軸（課題ピボット6観点）
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    orien_outline_text = st.session_state.get("orien_outline_text", "") or ""
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")

    # =========================================================
    # 2) kickoff / subq_list の取得（rev を正として同期）
    # =========================================================
    kickoff = (rev.get("kickoff") or {}) or {}
    kickoff_text = json.dumps(
        {
            "目標": kickoff.get("目標", ""),
            "現状": kickoff.get("現状", ""),
            "ビジネス課題": kickoff.get("ビジネス課題", ""),
            "調査目的": kickoff.get("調査目的", ""),
            "問い": kickoff.get("問い", ""),
            "仮説": kickoff.get("仮説", ""),
        },
        ensure_ascii=False,
        indent=2,
    )

    # ★重要：subq_list は rev を唯一の正とする
    _subq_list = (rev.get("subq_list") or []) or []
    if not _subq_list:
        return False, "サブクエスチョンがまだ生成されていません（active revision の subq_list が空です）。"

    # 正規化：subq/question を必ず揃える（過去の揺れ・None混入対策）
    norm_list = []
    for sq in _subq_list:
        d = dict(sq or {})
        txt = (d.get("subq") or d.get("question") or "").strip()
        d["subq"] = txt
        d["question"] = txt
        norm_list.append(d)

    # 生成直前に session_state も同期しておく（UI/他処理の参照ズレ防止）
    st.session_state["subq_list"] = norm_list
    st.session_state["ai_subquestions"] = rev.get("subquestions_raw", "") or ""

    subq_text_lines = [f"SQ{i}: {sq.get('subq', '')}" for i, sq in enumerate(norm_list, 1)]
    subq_text = "\n".join(subq_text_lines)

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    # =========================================================
    # 3) prompt
    # =========================================================
    prompt = f"""
あなたは市場調査設計の専門家です。
以下のサブクエスチョンそれぞれについて、次の6項目の観点から分析アプローチの下書きを作成してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（6観点）」は、依頼課題を調査で検証可能な採用課題（真の課題）へ変換したものです。
- axis/items/approach/hypothesis は、必ずこの軸に整合すること（論点がずれないこと）。
- 軸の補足条件・制約・懸念（6観点目の任意追記を含む）を前提として、必要な分析軸・評価項目を具体化すること。

▼軸情報（source: {axis_source}）
{axis_text}

【対象となる6項目】
- id: "SQ1" のようなID
- subq: サブクエスチョン本文
- axis: 分析軸（セグメント）
- items: 評価項目
- approach: 主な分析アプローチ（どのような切り口で分析するか）
- hypothesis: 検証する仮説

▼オリエン統合コンテキスト（アップロード抽出＋手入力）
{(ori_texts or "")[:4000]}

▼オリエン内容の整理（抜粋）
{(orien_outline_text or "")[:2000]}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

▼キックオフノート（参考）
{kickoff_text}

【サブクエスチョン一覧】
{subq_text}

【出力形式】
- 必ず JSON 配列のみを出力してください（余計な文章やコードブロックは書かないこと）

[
  {{
    "id": "SQ1",
    "subq": "...",
    "axis": "...",
    "items": "...",
    "approach": "...",
    "hypothesis": "..."
  }}
]

- 配列の要素数は、入力されたサブクエスチョンの数と同じにしてください。
- axis: 分析軸（セグメント）の案が複数ある場合は最も優先度の高いもの1つを提示してください。
- axis: 分析軸（セグメント）の案には15歳未満の対象属性を含めないこと（市場調査の綱領にて15歳未満にはアンケートを依頼することができないため）
- axis: 分析軸（セグメント）の案には80歳以上はアンケートに回答できない可能性が高いため含めないこと
- items: 評価項目案の後に（）で具体的な項目を記載してください。
- hypothesis: 語尾に「〜の可能性が高い（ある）」を用いないでください。
"""

    # =========================================================
    # 4) call + parse
    # =========================================================
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.6,
            max_tokens=2000,
        )
        ai_text = (response.choices[0].message.content or "").strip()

        # コードフェンス除去（最低限）
        if ai_text.startswith("```"):
            ai_text = ai_text.strip().strip("`")
            ai_text = ai_text.replace("json", "", 1).strip()

        blocks = json.loads(ai_text)
        if not isinstance(blocks, list):
            raise ValueError("JSON配列ではありません。")

        blocks = blocks[:9]  # MAX_I=9

        for i, b in enumerate(blocks, 1):
            if isinstance(b, dict):
                b["id"] = f"SQ{i}"
                # subq も空なら norm_list から埋める（安全策）
                if not (b.get("subq") or "").strip() and i <= len(norm_list):
                    b["subq"] = norm_list[i - 1].get("subq", "")

        # session_state に反映
        st.session_state["analysis_blocks"] = blocks

        # ★保存：active rev に analysis_blocks を確実に書き戻す（後でズレない）
        rev2 = dict(rev)
        rev2["analysis_blocks"] = blocks
        upsert_revision(rev2)

        return True, ""

    except Exception as e:
        return False, f"分析アプローチ案の生成中にエラーが発生しました: {e}"




def generate_target_condition_draft():
    """対象者条件案を AI で生成"""
    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    subquestions = st.session_state.get("ai_subquestions", "")

    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
    }

    # ★統合コンテキスト
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    # ★追加：フル生成の軸（課題ピボット6観点）
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    # ★推奨：kickoffを読みやすく
    kickoff_text = json.dumps(kickoff, ensure_ascii=False, indent=2)

    prompt = f"""
あなたは市場調査設計の専門家です。
以下の情報をもとに、この調査の「対象者条件」を検討してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（6観点）」は、依頼課題を調査で検証可能な採用課題（真の課題）へ変換したものです。
- 対象者条件は、この軸に整合するように設定してください（対象者がズレると検証できません）。
- 6観点目の任意追記（制約・懸念・補足）がある場合は必ず反映してください。

▼軸情報（source: {axis_source}）
{axis_text}

【出力形式】
- 対象者イメージ：
- エリア：
- 年齢・性別条件：
- 行動・意識・その他属性の条件：
- 除外条件：


【オリエン内容の整理（抜粋）】
{orien_outline_text[:2000]}

【キックオフノート（参考）】
{kickoff_text}

【問いの分解（AI生成サブクエスチョン）】
{subquestions}

【条件】
- 市場調査綱領にて、15歳未満にはアンケートを依頼することができないので対象者条件に含めないこと
- 80歳以上はアンケートに回答できない可能性が高いので対象者条件に含めないこと
- 対象者イメージは冒頭に簡潔に記載してください。
- “なんとなく広く”ではなく、上記の軸と問いに対して検証力が最大化するように絞り込んでください。
- 「●●来場経験者であること」ではなく「●●来場者」、「購買意思決定に関与する層」ではなく「購買意思決定者」など、簡潔な表現を用いてください。
- 除外条件に"調査に協力する意思がない層"、"過去の調査で回答の質に問題があった層"などアンケート条件とならないものは含めないでください。
- ###、** などの記号は使わないでください。
"""
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.6,
            max_tokens=650,
        )
        ai_text = response.choices[0].message.content.strip()
        st.session_state["ai_target_condition"] = ai_text
        return True, ""
    except Exception as e:
        return False, f"対象者条件案の生成中にエラーが発生しました: {e}"


def normalize_analysis_blocks(blocks: list[dict]) -> list[dict]:
    """analysis_blocks の id 欠落を吸収し、常に SQ1.. を付与する。"""
    if not isinstance(blocks, list):
        return []
    out = []
    for i, blk in enumerate(blocks, 1):
        if not isinstance(blk, dict):
            continue
        if not (blk.get("id") or "").strip():
            blk["id"] = f"SQ{i}"
        out.append(blk)
    return out



def generate_survey_items_linked_draft() -> tuple[bool, str]:
    """
    ⑤' 調査項目案を analysis_blocks（SQ別）起点で生成し、
    survey_item_rows（統合テーブル）として session_state に保持する（SQ分割生成）。
    """
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料が未入力です（アップロード or 手入力が必要です）。"

    orien_outline_text = st.session_state.get("orien_outline_text", "")
    if not (orien_outline_text or "").strip():
        return False, "オリエン内容の整理がまだ生成されていません。"

    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()
    if not axis_text:
        return False, "フル生成の軸が空です。『課題ピボット』→『生成・比較へ反映（確定）』を確認してください。"

    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
    }
    target_condition = st.session_state.get("ai_target_condition", "")

    blocks = normalize_analysis_blocks(st.session_state.get("analysis_blocks", []) or [])
    if not blocks:
        return False, "分析アプローチ案（analysis_blocks）がまだありません。先に③（分析アプローチ）まで生成してください。"

    # SQ id を必ず埋める（SQ最大9想定）
    blocks_min = []
    for i, b in enumerate(blocks, 1):
        bid = (b.get("id") or f"SQ{i}").strip()
        blocks_min.append({
            "id": bid,
            "subq": b.get("subq", ""),
            "axis": b.get("axis", ""), #分析軸
            "items": b.get("items", ""), #評価項目
            "approach": b.get("approach", ""),
            "hypothesis": b.get("hypothesis", ""),
        })
    

    kickoff_text = json.dumps(kickoff, ensure_ascii=False, indent=2)

    # UIから調整可能にしてもよいが、まず固定でOK（3±2）
    per_sq_target = int(st.session_state.get("per_sq_target_items", 3) or 3)
    per_sq_target = max(1, min(8, per_sq_target))

    all_rows: list[dict] = []
    try:
        for sq in blocks_min:
            # SQごと生成（打ち切り回避）
            sq_rows = _generate_items_for_one_sq(
                client=client,
                model=DEPLOYMENT,
                axis_source=axis_source,
                axis_text=axis_text,
                ori_texts=ori_texts,
                orien_outline_text=orien_outline_text,
                kickoff_text=kickoff_text,
                target_condition=target_condition,
                sq_block=sq,
                per_sq_target=per_sq_target,
                temperature=0.4,
                max_tokens=1200,
            )
            all_rows.extend(sq_rows)

        all_rows = _dedupe_rows(all_rows)

        if not all_rows:
            return False, "調査項目の生成結果が空でした（JSONの中身を確認してください）。"

        st.session_state["survey_item_rows"] = all_rows
        return True, ""

    except Exception as e:
        return False, f"調査項目（analysis_blocks連動）の生成中にエラーが発生しました: {e}"




import json
import re

def _strip_json_fence(s: str) -> str:
    s = (s or "").strip()
    if s.startswith("```"):
        s = s.strip("`").strip()
        # 先頭が json と書かれるケースを除去
        if s.lower().startswith("json"):
            s = s[4:].strip()
    return s

def _normalize_one_row(r: dict) -> dict | None:
    if not isinstance(r, dict):
        return None

    sq_id = (r.get("sq_id") or "").strip()
    if not sq_id:
        return None

    pr = r.get("priority")
    try:
        pr = int(pr)
    except Exception:
        pr = 3
    pr = max(1, min(5, pr))

    # まず out（返却するdict）を組み立てる
    out = {
        "sq_id": sq_id,
        "sq_subq": (r.get("sq_subq") or "").strip(),
        "items": (r.get("items") or "").strip(),
        "approach": (r.get("approach") or "").strip(),
        "var_name": (r.get("var_name") or "").strip(),
        "item_text": (r.get("item_text") or "").strip(),
        "recommended_type": (r.get("recommended_type") or "").strip(),
        "recommended_scale": (r.get("recommended_scale") or "").strip(),
        "priority": pr,
        # 入力側に入っていれば受け取り、なければ空
        "table_role": (r.get("table_role") or "").strip(),
    }

    # table_role を保証（欠損・表記ゆれを吸収）
    tr = str(out.get("table_role") or "").strip()
    if tr in ["表頭", "表側"]:
        out["table_role"] = tr
    elif tr in ["row", "side", "表側（行）", "表側(行)"]:
        out["table_role"] = "表側"
    elif tr in ["col", "head", "表頭（列）", "表頭(列)"]:
        out["table_role"] = "表頭"
    else:
        out["table_role"] = "表側"  # デフォルト

    return out


def _dedupe_rows(rows: list[dict]) -> list[dict]:
    """
    重複排除の最小実装：sq_id + var_name + item_text で重複判定。
    var_name が空の場合は item_text をより重視。
    """
    seen = set()
    out = []
    for r in rows:
        key = (
            (r.get("sq_id") or "").strip(),
            (r.get("var_name") or "").strip().lower(),
            (r.get("item_text") or "").strip(),
        )
        if key in seen:
            continue
        seen.add(key)
        out.append(r)
    return out

def _build_sq_prompt(
    *,
    axis_source: str,
    axis_text: str,
    ori_texts: str,
    orien_outline_text: str,
    kickoff_text: str,
    target_condition: str,
    sq_block: dict,
    per_sq_target: int,
) -> str:
    # SQ単位に絞ることで出力安定＆打ち切り回避
    return f"""
あなたは市場調査設計の専門家です。
以下の「分析アプローチ（対象SQ）」に基づいて、調査票の“調査項目（変数）”を設計してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（軸）」に論点を必ず整合させてください。
- items / approach / hypothesis を“測定可能な調査項目”に落としてください。
- 調査項目は「分析で使える変数」になるように書いてください（抽象語で終わらない）。

▼軸情報（source: {axis_source}）
{axis_text}

▼オリエン統合コンテキスト（抜粋）
{(ori_texts or "")[:2500]}

▼オリエン内容の整理（抜粋）
{(orien_outline_text or "")[:1500]}

▼キックオフノート（参考）
{kickoff_text}

▼対象者条件（参考）
{target_condition}

▼分析アプローチ（対象SQ：JSON）
{json.dumps(sq_block, ensure_ascii=False, indent=2)}

【出力形式（厳守）】
- 必ず JSON 配列 “だけ” を出力してください（余計な文章、コードブロック禁止）。
- 1要素＝1調査項目（変数案）。
- 各要素は以下キーを必須とします（空でもキーは出す）。

[
  {{
    "sq_id": "{sq_block.get('id','')}",
    "sq_subq": "{sq_block.get('subq','')}",
    "items": "{sq_block.get('items','')}",
    "approach": "{sq_block.get('approach','')}",
    "var_name": "",
    "item_text": "",
    "recommended_type": "SA/MA/尺度/数値/自由回答",
    "recommended_scale": "",
    "table_role": "表頭/表側",
    "priority": 1
  }}
]

【生成ルール】
- このSQについて {per_sq_target}±2 個の調査項目を提案してください。
- priority は 1（最重要）〜5（補助）で付与してください。
- item_text は1行で、全角60字以内目安。
- sq_id / sq_subq / items / approach は上記の対象SQ情報を維持してください（勝手に別SQへ変更しない）。
- table_role の定義：
  表頭：属性・セグメント・分類軸（例：年代、利用有無、業種、頻度区分など）
  表側：評価・態度・行動・量・スコア（例：満足、購入意向、理由、利用実態など）
  必ず "表頭" または "表側" のどちらかを出力すること。
""".strip()

def _generate_items_for_one_sq(
    *,
    client,
    model: str,
    axis_source: str,
    axis_text: str,
    ori_texts: str,
    orien_outline_text: str,
    kickoff_text: str,
    target_condition: str,
    sq_block: dict,
    per_sq_target: int = 3,
    temperature: float = 0.4,
    max_tokens: int = 1200,
) -> list[dict]:
    prompt = _build_sq_prompt(
        axis_source=axis_source,
        axis_text=axis_text,
        ori_texts=ori_texts,
        orien_outline_text=orien_outline_text,
        kickoff_text=kickoff_text,
        target_condition=target_condition,
        sq_block=sq_block,
        per_sq_target=per_sq_target,
    )

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
            {"role": "user", "content": prompt},
        ],
        temperature=temperature,
        max_tokens=max_tokens,
    )

    ai_text = _strip_json_fence((response.choices[0].message.content or "").strip())
    rows = json.loads(ai_text)
    if not isinstance(rows, list):
        raise ValueError("JSON配列ではありません。")

    norm_rows = []
    for r in rows:
        nr = _normalize_one_row(r)
        if not nr:
            continue

        # 念のため：sq_id を対象SQに寄せる（モデル揺れ対策）
        nr["sq_id"] = (sq_block.get("id") or nr["sq_id"]).strip()
        nr["sq_subq"] = nr["sq_subq"] or (sq_block.get("subq") or "").strip()
        nr["items"] = nr["items"] or (sq_block.get("items") or "").strip()
        nr["approach"] = nr["approach"] or (sq_block.get("approach") or "").strip()

        if not nr["item_text"]:
            continue

        norm_rows.append(nr)

    return norm_rows













def build_survey_versions_from_rows(rows: list[dict]) -> dict:
    """
    survey_item_rows から 10/20/30/40 問バージョンを再構成する。
    priority（小さいほど重要）→SQ順の安定ソートで上位から切り出し。
    """
    if not rows:
        return {"10問": "", "20問": "", "30問": "", "40問": ""}

    # 安定ソート（priority→sq_id→var_name）
    def key_fn(r: dict):
        return (
            int(r.get("priority", 3)),
            str(r.get("sq_id", "")),
            str(r.get("var_name", "")),
            str(r.get("item_text", "")),
        )

    sorted_rows = sorted(rows, key=key_fn)

    def to_text(n: int) -> str:
        picked = sorted_rows[:n]
        lines = []
        for r in picked:
            sq = r.get("sq_id", "")
            item = r.get("item_text", "")
            if not item:
                continue
            # 紐づけが見えるように SQ を先頭に付与
            lines.append(f"・[{sq}] {item}")
        return "\n".join(lines)

    return {
        "10問": to_text(10),
        "20問": to_text(20),
        "30問": to_text(30),
        "40問": to_text(40),
    }






def generate_survey_items_draft():
    """調査項目案（10/20/30/40問バージョン）を AI で生成"""
    orien_outline_text = st.session_state.get("orien_outline_text", "")
    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")

    kickoff = {
        "目標": st.session_state.get("ai_目標", ""),
        "現状": st.session_state.get("ai_現状", ""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
        "調査目的": st.session_state.get("ai_調査目的", ""),
        "問い": st.session_state.get("ai_問い", ""),
        "仮説": st.session_state.get("ai_仮説", ""),
    }
    subquestions = st.session_state.get("ai_subquestions", "")
    target_condition = st.session_state.get("ai_target_condition", "")

    if not orien_outline_text.strip():
        return False, "オリエン内容の整理がまだ生成されていません。"

    # ★追加：オリエン統合コンテキスト（具体情報）
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    # ★追加：フル生成の軸（課題ピボット6観点）
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    prompt = f"""
あなたは市場調査設計の専門家です。
以下の情報をもとに、この調査で実施すべき調査項目案を提案してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（6観点）」は、依頼課題を調査で検証可能な採用課題（真の課題）へ変換したものです。
- 調査項目は、この軸と整合し、検証可能性が最大化するように構成してください。
- 6観点目の任意追記（制約・懸念・補足）がある場合は必ず反映してください。

▼軸情報（source: {axis_source}）
{axis_text}

【出力条件】
- 選択肢は不要（設問文のみ）
- 設問文は質問文形式でなく、調査項目名として簡潔に表現する
- 各バージョンで「ちょうど」下記の問数になるようにする（10問／20問／30問／40問）
- 各設問は1行で簡潔に（目安：全角60文字以内）
- 見出しと設問リスト以外の説明文は出力しない

【出力形式】（この見出し形式を厳守）
# 10問バージョン
・ ...
（10問まで）

# 20問バージョン
・ ...
（20問まで）

# 30問バージョン
・ ...
（30問まで）

# 40問バージョン
・ ...
（40問まで）

【オリエン統合コンテキスト（アップロード抽出＋手入力）】
{ori_texts[:3500]}

【オリエン内容の整理（抜粋）】
{orien_outline_text[:2000]}

【ブランド診断：カテゴリー構造】
{cat_text}

【ブランド診断：消費行動特性】
{beh_text}

【キックオフノート】
{kickoff}

【問いの要因分解】
{subquestions}

【対象者条件】
{target_condition}
"""
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.6,
            max_tokens=3200,
        )
        ai_text = response.choices[0].message.content.strip()
        st.session_state["ai_survey_items_raw"] = ai_text

        versions = {}
        for ver in ["10問", "20問", "30問", "40問"]:
            pattern = rf"#\s*{ver}バージョン(.*?)(?=#\s*\d+問バージョン|$)"
            m = re.search(pattern, ai_text, re.DOTALL)
            versions[ver] = m.group(1).strip() if m else ""
        st.session_state["ai_survey_items"] = versions
        return True, ""
    except Exception as e:
        return False, f"調査項目案の生成中にエラーが発生しました: {e}"

def generate_spec_draft():
    """調査仕様案を AI で JSON 生成"""
    orien_outline_text = st.session_state.get("orien_outline_text", "")
    target_condition = st.session_state.get("ai_target_condition", "")
    ai_survey_items = st.session_state.get("ai_survey_items", {}) or {}

    # ★追加：オリエン統合コンテキスト（具体情報）
    ori_texts = get_orien_context_text()
    if not ori_texts.strip():
        return False, "オリエン資料（アップロード）または手入力内容がありません。"

    # ★追加：フル生成の軸（課題ピボット6観点）
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    axis_source = (st.session_state.get("fullgen_axis_source") or "pivot").strip()
    if not axis_text:
        return False, "フル生成の軸（課題ピボット6観点）が空です。『課題ピボット』で確定してから実行してください。"

    # 調査項目案のうち、優先的に 20問 → 30問 → 40問 → 10問 の順で採用
    survey_items_selected = ""
    for ver in ["20問", "30問", "40問", "10問"]:
        if ai_survey_items.get(ver):
            survey_items_selected = ai_survey_items[ver]
            break

    # ★追加：調査項目案が空なら止める（仕様がテンプレ化しやすい）
    if not str(survey_items_selected).strip():
        return False, "調査項目案が未生成です。先に『調査項目案』を生成してください。"

    cat_df = st.session_state.get("df_category_structure")
    beh_df = st.session_state.get("df_behavior_traits")
    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

    if not orien_outline_text.strip():
        return False, "オリエン内容の整理がまだ生成されていません。"

    prompt = f"""
あなたは市場調査設計の専門家です。
以下の情報をもとに、この調査の「調査仕様案」を項目ごとに整理してください。

【最優先の前提（フル生成の軸）】
- 以下の「課題ピボット（6観点）」は、依頼課題を調査で検証可能な採用課題（真の課題）へ変換したものです。
- 調査仕様はこの軸と整合し、検証可能性・実査設計の妥当性が最大化するように記述してください。
- 6観点目の任意追記（制約・懸念・補足）がある場合は必ず反映してください。

▼軸情報（source: {axis_source}）
{axis_text}

【入力情報】
▼オリエン統合コンテキスト（アップロード抽出＋手入力）
{ori_texts[:3500]}

▼オリエン内容の整理（抜粋）
{orien_outline_text[:2000]}

▼対象者条件
{target_condition}

▼調査項目案（採用版）
{survey_items_selected}

▼参考情報：カテゴリー構造
{cat_text}

▼参考情報：消費行動特性
{beh_text}

【出力する項目】
- 調査手法
- 抽出方法
- 調査地域
- 対象者条件
- サンプルサイズ
- 調査ボリューム
- 提示物
- 集計・分析仕様
- 自由回答データの処理
- 業務範囲
- 納品物
- インスペクションの方法
- 謝礼の種類
- 備考

【出力形式】
次のキーを持つ JSON オブジェクト「だけ」を出力してください。

{{
  "調査手法": "...",
  "抽出方法": "...",
  "調査地域": "...",
  "対象者条件": "...",
  "サンプルサイズ": "...",
  "調査ボリューム": "...",
  "提示物": "...",
  "集計・分析仕様": "...",
  "自由回答データの処理": "...",
  "業務範囲": "...",
  "納品物": "...",
  "インスペクションの方法": "...",
  "謝礼の種類": "...",
  "備考": "..."
}}

【ルール】
- 調査手法は特に明記がなければ「インターネット調査」を基本としてください。
- 抽出方法は特に明記がなければ「割付抽出」を基本としてください。
- 調査ボリュームは「スクリーニング調査」「本調査」を2行に分けて記載してください。
- 自由回答データの処理は、入力テキストに記載がなければ「なし」を基本としてください。
- インスペクションの方法は、入力テキストに記載がなければ「性別・年齢（2歳以上）のアンマッチの場合は、対象除外とする。」を基本としてください。
- 謝礼の種類は、入力テキストに記載がなければ「ポイント謝礼」を基本としてください。
- 不明な項目は、前提（軸）に照らして現実的な案を補い、備考に「仮置き」と明記してください。
"""
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.5,
            max_tokens=1000,
        )
        ai_text = response.choices[0].message.content.strip()
        if ai_text.startswith("```"):
            ai_text = ai_text.strip("`")
            ai_text = ai_text.replace("json", "", 1).strip()

        spec_obj = json.loads(ai_text)
        for label, key in SPEC_ITEMS:
            st.session_state[key] = spec_obj.get(label, "")

        return True, ""
    except Exception as e:
        return False, f"調査仕様案の生成中にエラーが発生しました: {e}"

# =========================================================
# プロジェクト保存／読み込みヘルパー
# =========================================================

def build_project_from_session() -> dict:
    """
    現在の st.session_state からプロジェクト構造を組み立てて返す。
    JSON にして保存し、あとで apply_project_to_session() で復元できる前提。
    """
    # DataFrame は JSON で扱いやすいようにレコード配列に変換
    df_cat = st.session_state.get("df_category_structure")
    df_beh = st.session_state.get("df_behavior_traits")

    cat_records = df_cat.to_dict(orient="records") if isinstance(df_cat, pd.DataFrame) else []
    beh_records = df_beh.to_dict(orient="records") if isinstance(df_beh, pd.DataFrame) else []

    project_name = st.session_state.get("project_name", "")

    proj = {
        "version": 1,
        "meta": {
            "project_name": project_name,
        },
        "orien": {
            "outline_text": st.session_state.get("orien_outline_text", ""),
            "uploaded_raw_texts": st.session_state.get("uploaded_docs", []),
        },
        "brand_diagnosis": {
            "target_category": st.session_state.get("target_category", ""),
            "target_brand": st.session_state.get("target_brand", ""),
            "df_category_structure": cat_records,
            "df_behavior_traits": beh_records,
            "funnel_text": st.session_state.get("funnel_text", ""),
        },
        "proposal": {
            "kickoff_selected_purpose": st.session_state.get("kickoff_selected_purpose"),
            "kickoff_purpose_free": st.session_state.get("kickoff_purpose_free", ""),
            "kickoff": {
                "目標": st.session_state.get("ai_目標", ""),
                "現状": st.session_state.get("ai_現状", ""),
                "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
                "調査目的": st.session_state.get("ai_調査目的", ""),
                "問い": st.session_state.get("ai_問い", ""),
                "仮説": st.session_state.get("ai_仮説", ""),
            },
            "subquestions_raw": st.session_state.get("ai_subquestions", ""),
            "subquestions_structured": st.session_state.get("subq_list", []),
            "analysis_blocks": st.session_state.get("analysis_blocks", []),
            "target_condition": st.session_state.get("ai_target_condition", ""),
            "survey_items": st.session_state.get("ai_survey_items", {}),
            "spec": {label: st.session_state.get(key, "") for (label, key) in SPEC_ITEMS},
            "revisions": st.session_state.get("proposal_revisions", []),
            "active_rev_id": st.session_state.get("active_rev_id"),
            "problem_reframe_output": st.session_state.get("problem_reframe_output", {}) or {},
            "true_problem_text": st.session_state.get("true_problem_text", "") or "",


        },
    }

    return proj


def apply_project_to_session(project: dict) -> None:
    """
    保存済みプロジェクトJSONを st.session_state に展開して復元する。
    復元後に st.session_state["selected_mode"] = "proposal_draft" などをセットして rerun すれば、
    企画書下書きの画面がそのまま再掲される。
    """
    meta = project.get("meta", {})
    orien = project.get("orien", {})
    brand = project.get("brand_diagnosis", {})
    proposal = project.get("proposal", {})

    
    st.session_state["proposal_revisions"] = proposal.get("revisions", []) or []
    st.session_state["active_rev_id"] = proposal.get("active_rev_id")

    # active が無い場合は最新をactiveにする
    if not st.session_state["active_rev_id"] and st.session_state["proposal_revisions"]:
        st.session_state["active_rev_id"] = st.session_state["proposal_revisions"][-1].get("rev_id")

    # active revision を互換キーへ展開
    sync_active_revision_to_session_keys()


    # プロジェクト名
    st.session_state["project_name"] = meta.get("project_name", "")

    # オリエン関連
    st.session_state["orien_outline_text"] = orien.get("outline_text", "")
    st.session_state["orien_outline_editor"] = orien.get("outline_text", "")
    st.session_state["uploaded_docs"] = orien.get("uploaded_raw_texts", [])

    # ブランド診断
    st.session_state["target_category"] = brand.get("target_category", "")
    st.session_state["target_brand"] = brand.get("target_brand", "")
    st.session_state["funnel_text"] = brand.get("funnel_text", "")

    cat_records = brand.get("df_category_structure", [])
    beh_records = brand.get("df_behavior_traits", [])

    if cat_records:
        st.session_state["df_category_structure"] = pd.DataFrame(cat_records)
    else:
        st.session_state.pop("df_category_structure", None)

    if beh_records:
        st.session_state["df_behavior_traits"] = pd.DataFrame(beh_records)
    else:
        st.session_state.pop("df_behavior_traits", None)

    st.session_state["kickoff_selected_purpose"] = proposal.get("kickoff_selected_purpose")
    st.session_state["kickoff_purpose_free_editor"] = st.session_state.get("kickoff_purpose_free", "")



    # 企画書下書き：キックオフ
    kickoff = proposal.get("kickoff", {})
    st.session_state["ai_目標"] = kickoff.get("目標", "")
    st.session_state["ai_現状"] = kickoff.get("現状", "")
    st.session_state["ai_ビジネス課題"] = kickoff.get("ビジネス課題", "")
    st.session_state["ai_調査目的"] = kickoff.get("調査目的", "")
    st.session_state["ai_問い"] = kickoff.get("問い", "")
    st.session_state["ai_仮説"] = kickoff.get("仮説", "")

    # サブクエスチョン
    st.session_state["ai_subquestions"] = proposal.get("subquestions_raw", "")
    st.session_state["subq_list"] = proposal.get("subquestions_structured", [])

    # 分析アプローチ
    st.session_state["analysis_blocks"] = proposal.get("analysis_blocks", [])

    # 対象者条件
    st.session_state["ai_target_condition"] = proposal.get("target_condition", "")

    # 調査項目案
    st.session_state["ai_survey_items"] = proposal.get("survey_items", {})

    # 調査仕様
    spec_obj = proposal.get("spec", {})
    for label, key in SPEC_ITEMS:
        st.session_state[key] = spec_obj.get(label, "")

    # 企画書下書きモードを有効化して再掲できるようにする
    st.session_state["proposal_draft_generated"] = True
    st.session_state["selected_mode"] = "proposal_draft"

    st.session_state["problem_reframe_output"] = proposal.get("problem_reframe_output", {}) or {}
    st.session_state["true_problem_text"] = proposal.get("true_problem_text", "") or ""
    st.session_state["problem_reframe_generated"] = bool(st.session_state["problem_reframe_output"])



def reset_proposal_state():
    keys = [
        "ai_目標", "ai_現状", "ai_ビジネス課題", "ai_調査目的",
        "ai_問い", "ai_仮説",
        "ai_subquestions", "subq_list",
        "analysis_blocks",
        "ai_target_condition",
        "ai_survey_items",
        "proposal_draft_generated",
    ]
    for k in keys:
        st.session_state.pop(k, None)



def build_question_tree_text(session_state: dict) -> str:
    """
    以前の「問いの分解」画面と同等の構造ビュー文字列(tree_text)を生成する。
    目的 → メインクエスチョン → サブクエスチョン をテキストで表現。
    """
    main_question_text = (session_state.get("ai_問い", "") or "").strip()
    purpose = (session_state.get("ai_調査目的", "") or "").strip()
    subq_list = session_state.get("subq_list", []) or []

    def split_main_questions(text: str):
        if not text:
            return []
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        questions = []
        buf = ""
        for line in lines:
            m = re.match(r'^(?:\d+[\.\)]|Q\d+[:：])\s*(.+)', line)
            if m:
                if buf:
                    questions.append(buf.strip())
                buf = m.group(1)
            else:
                if buf:
                    buf += " " + line
                else:
                    buf = line
        if buf:
            questions.append(buf.strip())
        if not questions and text.strip():
            questions = [text.strip()]
        return questions

    main_questions = split_main_questions(main_question_text)

    tree_lines = []

    # 目的
    tree_lines.append("目的（キックオフノート）")
    tree_lines.append(f"  ┗ {purpose}" if purpose else "  ┗ （未設定）")
    tree_lines.append("")

    # 問い
    tree_lines.append("問い（メインクエスチョン）")

    if not main_questions:
        # メインQが抽出できない場合は全文を1つとして扱う
        if main_question_text:
            tree_lines.append(f"  ┗ {main_question_text}")
            if subq_list:
                for i, sq in enumerate(subq_list, 1):
                    tree_lines.append(f"       ┗ SQ{i}: {sq.get('subq', '')}")
            else:
                tree_lines.append("       ┗ （まだサブクエスチョンが生成されていません）")
        else:
            tree_lines.append("  ┗ （未設定）")
    else:
        # メインQが複数ある場合：subq_list を割当
        grouped = {mq: [] for mq in main_questions}

        if subq_list:
            for idx, sq in enumerate(subq_list):
                explicit_parent = sq.get("main_question") or sq.get("main")
                if explicit_parent and explicit_parent in grouped:
                    grouped[explicit_parent].append(sq)
                else:
                    target_mq = main_questions[idx % len(main_questions)]
                    grouped[target_mq].append(sq)

        for mq in main_questions:
            tree_lines.append(f"  ┗ {mq}")
            sq_items = grouped.get(mq, [])
            if not sq_items:
                tree_lines.append("       ┗ （まだサブクエスチョンが紐付いていません）")
            else:
                for i, sq in enumerate(sq_items, 1):
                    tree_lines.append(f"       ┗ SQ{i}: {sq.get('subq', '')}")
            tree_lines.append("")

    tree_text = "\n".join(tree_lines).strip()
    return tree_text


def ensure_question_tree_cached(session_state: dict) -> str:
    """
    STEP4向け：edited_texts["EDIT1_subQ"] を必ず埋める。
    - 既に入っていればそれを採用
    - 無ければ UI表示用 subq_structured_view を採用
    - それも無ければ ai_問い / subq_list から再生成
    返り値は最終的に採用したテキスト。
    """
    if session_state.get("edited_texts") is None:
        session_state["edited_texts"] = {}
    edited = session_state["edited_texts"]

    # 1) 既にキャッシュがあれば最優先
    val = (edited.get("EDIT1_subQ") or "").strip()
    if val:
        return val

    # 2) proposal_draft画面の text_area(key="subq_structured_view") があれば採用
    val = (session_state.get("subq_structured_view") or "").strip()
    if val:
        edited["EDIT1_subQ"] = val
        return val

    # 3) なければ再計算
    val = build_question_tree_text(session_state).strip()
    if val:
        edited["EDIT1_subQ"] = val
    return val

def ensure_target_condition_cached(session_state: dict) -> str:
    """
    STEP4向け：edited_texts["EDIT1_taisyosya"] を必ず埋める。
    優先順位：
      1) edited_texts["EDIT1_taisyosya"]（既に手で入れている場合）
      2) UI編集値 target_condition_textarea
      3) ai_target_condition
    """
    if session_state.get("edited_texts") is None:
        session_state["edited_texts"] = {}
    edited = session_state["edited_texts"]

    # 1) 既存キャッシュ
    val = (edited.get("EDIT1_taisyosya") or "").strip()
    if val:
        return val

    # 2) UI編集値（旧コードと同じキー）
    val = (session_state.get("target_condition_textarea") or "").strip()
    if val:
        edited["EDIT1_taisyosya"] = val
        return val

    # 3) AI生成結果
    val = (session_state.get("ai_target_condition") or "").strip()
    if val:
        edited["EDIT1_taisyosya"] = val
        return val

    return ""


def ensure_survey_items_cached(session_state: dict) -> dict:
    """
    STEP4向け：調査項目案（10/20/30/40）を edited_texts に必ず格納して返す。

    参照優先順位（verごと）：
      1) edited_texts[shape_name]（既に手動反映済みなど）
      2) session_state["survey_items_{ver}"]（proposal_draft 画面の編集値）
      3) session_state["ai_survey_items"][ver]（AI生成結果）
    Returns:
      {
        "10問": "...",
        "20問": "...",
        "30問": "...",
        "40問": "..."
      }
    """
    if session_state.get("edited_texts") is None:
        session_state["edited_texts"] = {}
    edited = session_state["edited_texts"]

    ai_items = session_state.get("ai_survey_items", {}) or {}

    ver_to_shape = {
        "10問": "EDIT1_Qimg",
        "20問": "EDIT2_Qimg",
        "30問": "EDIT3_Qimg",
        "40問": "EDIT4_Qimg",
    }

    out = {}

    for ver, shape in ver_to_shape.items():
        # 1) 既存キャッシュ
        val = (edited.get(shape) or "").strip()
        if val:
            out[ver] = val
            continue

        # 2) proposal_draft側の編集値（key名はあなたの現行UIに合わせる）
        ui_key = f"survey_items_{ver}"
        val = (session_state.get(ui_key) or "").strip()
        if val:
            edited[shape] = val
            out[ver] = val
            continue

        # 3) AI生成結果
        val = (ai_items.get(ver) or "").strip()
        if val:
            edited[shape] = val
            out[ver] = val
            continue

        out[ver] = ""

    return out




import json
from datetime import datetime

def build_ppt_update_payload(session_state: dict) -> dict:
    """
    Streamlitの session_state（＝UIのTXT入力結果）から、
    PPT反映用の正規化JSON payload を生成する。

    Returns:
      {
        "meta": {...},
        "items": [
          {
            "slide_index": 0,
            "slide_no": 1,
            "shape_name": "Edit_client",
            "text": "...",
            "source": {"type": "session_state", "key": "Edit_client"}
          },
          ...
        ]
      }
    """

    def add_item(items, slide_index: int, shape_name: str, text: str, source_key: str | None = None):
        """空文字はスキップ。"""
        if text is None:
            return
        if isinstance(text, str) and text.strip() == "":
            return

        items.append({
            "slide_index": int(slide_index),
            "slide_no": int(slide_index) + 1,   # 1始まり表示用
            "shape_name": str(shape_name),
            "text": str(text),
            "source": ({
                "type": "session_state",
                "key": source_key
            } if source_key else {"type": "computed"})
        })

    items: list[dict] = []

    # ------------------------------------------------------------
    # スライド1（表紙）: slide_index=0
    # ------------------------------------------------------------
    add_item(items, 0, "Edit_client", session_state.get("Edit_client", ""), "Edit_client")
    add_item(items, 0, "Edit_title",  session_state.get("Edit_title",  ""), "Edit_title")
    add_item(items, 0, "Edit_date",   session_state.get("Edit_date",   ""), "Edit_date")

    # ------------------------------------------------------------
    # スライド2（キックオフノート）: slide_index=1
    # ------------------------------------------------------------
    kickoff_map = {
        "EDIT_TO_BE":     "ai_目標",
        "EDIT_AS_IS":     "ai_現状",
        "EDIT_PROBLEM":   "ai_ビジネス課題",
        "EDIT_PURPOSE":   "ai_調査目的",
        "EDIT_QUESTION":  "ai_問い",
        "EDIT_HYPOTHESIS":"ai_仮説",
    }
    for shape, key in kickoff_map.items():
        add_item(items, 1, shape, session_state.get(key, ""), key)


    # ------------------------------------------------------------
    # スライド3（問いの分解ツリー）: slide_index=2
    # 以前の画面と同等の tree_text を生成し、EDIT1_subQ に反映する
    # ------------------------------------------------------------
    tree_text = ensure_question_tree_cached(session_state)
    add_item(items, 2, "EDIT1_subQ", tree_text, "EDIT1_subQ")


    # ------------------------------------------------------------
    # 対象者条件（旧：スライド4）: slide_index=12
    # shape: EDIT1_taisyosya
    # ------------------------------------------------------------
    target_text = ensure_target_condition_cached(session_state)
    add_item(items, 12, "EDIT1_taisyosya", target_text, "ai_target_condition")



    # ------------------------------------------------------------
    # スライド4-12（分析アプローチ）：slide_index=3..11（最大9件 = i=1..9）
    # shape: EDIT1_subQ{i}_1..5
    # 値：session_state["analysis_*_{i}"] から取得
    # ------------------------------------------------------------
    MAX_I = 9
    for i in range(1, MAX_I + 1):
        slide_index = 3 + (i - 1)  # 3..11
        add_item(items, slide_index, f"EDIT1_subQ{i}_1", session_state.get(f"analysis_subq_{i}", ""), f"analysis_subq_{i}")
        add_item(items, slide_index, f"EDIT1_subQ{i}_2", session_state.get(f"analysis_axis_{i}", ""), f"analysis_axis_{i}")
        add_item(items, slide_index, f"EDIT1_subQ{i}_3", session_state.get(f"analysis_items_{i}", ""), f"analysis_items_{i}")
        add_item(items, slide_index, f"EDIT1_subQ{i}_4", session_state.get(f"analysis_approach_{i}", ""), f"analysis_approach_{i}")
        add_item(items, slide_index, f"EDIT1_subQ{i}_5", session_state.get(f"analysis_hypothesis_{i}", ""), f"analysis_hypothesis_{i}")

    # ------------------------------------------------------------
    # 調査項目案（10/20/30/40）: slide_index=13
    # shape:
    #   10問→EDIT1_Qimg
    #   20問→EDIT2_Qimg
    #   30問→EDIT3_Qimg
    #   40問→EDIT4_Qimg
    # ------------------------------------------------------------
    survey_map = ensure_survey_items_cached(session_state)

    add_item(items, 13, "EDIT1_Qimg", survey_map.get("10問", ""), "survey_items_10問")
    add_item(items, 13, "EDIT2_Qimg", survey_map.get("20問", ""), "survey_items_20問")
    add_item(items, 13, "EDIT3_Qimg", survey_map.get("30問", ""), "survey_items_30問")
    add_item(items, 13, "EDIT4_Qimg", survey_map.get("40問", ""), "survey_items_40問")




    # ------------------------------------------------------------
    # スライド14（調査仕様案）: slide_index=14
    # SPEC_ITEMS / SPEC_LABEL_TO_SHAPE は既存コードの定義を利用する前提
    # ------------------------------------------------------------
    try:
        for label, ss_key in SPEC_ITEMS:
            shape_name = SPEC_LABEL_TO_SHAPE.get(label)
            if shape_name:
                add_item(items, 14, shape_name, session_state.get(ss_key, ""), ss_key)
    except NameError:
        # まだSPEC_ITEMS等を読み込んでいない場合は何もしない
        pass

    # ------------------------------------------------------------
    # スライド16（概算見積）: slide_index=16（ご指定どおり）
    # shape: EDIT_amount1..5
    # 値：session_state["estimate_summary{idx}"] を優先
    # ------------------------------------------------------------
    for idx in range(1, 6):
        key = f"estimate_summary{idx}"
        add_item(items, 16, f"EDIT_amount{idx}", session_state.get(key, ""), key)

    payload = {
        "meta": {
            "generated_at": datetime.now().isoformat(timespec="seconds"),
            "mode": session_state.get("selected_mode"),
            "pptx_path": session_state.get("pptx_path"),
            "items_count": len(items),
        },
        "items": items
    }
    return payload


def payload_to_pretty_json(payload: dict) -> str:
    return json.dumps(payload, ensure_ascii=False, indent=2)


from pptx import Presentation
from pathlib import Path
from datetime import datetime

def reflect_payload_to_pptx(
    pptx_path: str,
    payload: dict,
    session_state: dict,
    save_dir: Path,
    out_filename_prefix: str = "proposal_reflected",
    apply_default_format: bool = True,
):
    """
    STEP4: payload(JSON) を PPT に反映して保存する。

    Args:
      pptx_path: 反映対象のpptxファイルパス（テンプレ or 作業中ppt）
      payload: build_ppt_update_payload() が返す dict
      session_state: st.session_state を渡す（edited_texts 等の更新用）
      save_dir: get_session_dir() の戻りなど。保存先ディレクトリ
      out_filename_prefix: 出力ファイルprefix
      apply_default_format: Trueなら apply_text_format を適用（shapeによって例外あり）

    Returns:
      (out_path: Path, report: dict)
    """

    prs = Presentation(pptx_path)

    items = payload.get("items", []) or []

    # --- レポート ---
    report = {
        "pptx_in": str(pptx_path),
        "pptx_out": None,
        "total_items": len(items),
        "applied": 0,
        "skipped_empty": 0,
        "slide_oob": 0,
        "shape_not_found": 0,
        "errors": [],
    }

    # edited_texts キャッシュ（プレビュー用）を更新
    if "edited_texts" not in session_state or session_state["edited_texts"] is None:
        session_state["edited_texts"] = {}
    edited_texts = session_state["edited_texts"]

    for it in items:
        try:
            slide_index = int(it.get("slide_index"))
            shape_name = str(it.get("shape_name", ""))
            text_val = it.get("text", "")

            # 空はスキップ（STEP3でも基本弾いているが二重ガード）
            if text_val is None or (isinstance(text_val, str) and text_val.strip() == ""):
                report["skipped_empty"] += 1
                continue

            # スライド範囲外
            if slide_index < 0 or slide_index >= len(prs.slides):
                report["slide_oob"] += 1
                continue

            slide = prs.slides[slide_index]

            # 既存関数で書き込み（グループ内探索含む）
            ok = set_text_to_named_shape(slide, shape_name, str(text_val))

            if not ok:
                report["shape_not_found"] += 1
                continue

            # --- 書式適用（必要に応じて分岐）---
            # set_text_to_named_shape は「黒」にしてくれるが、
            # フォント名/サイズ/左寄せ統一をしたい場合は apply_text_format を追加適用
            if apply_default_format:
                # shapeを（グループ外なら）直接取得して apply_text_format
                shp = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name), None)
                if shp and getattr(shp, "has_text_frame", False):
                    # 概算見積（slide_index=16）のEDIT_amount*は10pt固定など、既存仕様に合わせる
                    if slide_index == 16 and shape_name.startswith("EDIT_amount"):
                        apply_text_format(shp, font_size=10)
                    else:
                        apply_text_format(shp)  # Arial 12pt 黒 左寄せ（あなたの定義どおり）
                # グループ内shapeは next() で取れない可能性があるので、取れない場合はスキップ（落とさない）

            # プレビュー用キャッシュ
            edited_texts[shape_name] = str(text_val)
            report["applied"] += 1

        except Exception as e:
            report["errors"].append({"item": it, "error": str(e)})

    # --- 保存 ---
    save_dir.mkdir(parents=True, exist_ok=True)
    out_path = save_dir / f"{out_filename_prefix}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(out_path)

    report["pptx_out"] = str(out_path)
    return out_path, report


def run_step4_apply_current_ui_to_ppt(session_state: dict):
    """
    Streamlit上で呼び出すための薄いラッパー。
    - STEP3: payload生成
    - STEP4: 反映 + 保存
    - session_state["pptx_path"] 更新
    """

    pptx_path = session_state.get("pptx_path")
    if not pptx_path:
        raise ValueError("pptx_path が未設定です。先にPPTテンプレートをアップロードしてください。")

    # STEP3（既に用意済み）
    payload = build_ppt_update_payload(session_state)

    # 保存先
    save_dir = get_session_dir()

    out_path, report = reflect_payload_to_pptx(
        pptx_path=str(pptx_path),
        payload=payload,
        session_state=session_state,
        save_dir=save_dir,
        out_filename_prefix="proposal_reflected_all",
        apply_default_format=True,
    )

    # 次工程のためにカレントPPTを更新
    session_state["pptx_path"] = str(out_path)

    # ★ダウンロード用に bytes をセッションに保持（これが重要）
    with open(out_path, "rb") as f:
        session_state["final_pptx_bytes"] = f.read()
    session_state["final_pptx_name"] = out_path.name
    session_state["final_pptx_key"] = f"dl_{out_path.name}"  # ★毎回変わるキー

    return out_path, report



def render_ppt_download_button():
    pptx_bytes = st.session_state.get("final_pptx_bytes")
    filename = st.session_state.get("final_pptx_name", "reflected.pptx")
    dl_key = st.session_state.get("final_pptx_key", "dl_default")

    if not pptx_bytes:
        #st.info("まだPPTが生成されていません。")
        return

    st.download_button(
        label="⬇️ 反映済みPPTをダウンロード",
        data=pptx_bytes,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key=dl_key,  # ★これが効きます
        use_container_width=True,
    )

def delete_revision(rev_id: str) -> tuple[bool, str]:
    ensure_revision_store()

    revs = get_revisions()
    if not revs:
        return False, "Revisionがありません。"

    target = find_revision(rev_id)
    if not target:
        return False, "指定したRevisionが見つかりません。"

    if target.get("stage") == "default":
        return False, "DEFAULT Revision は削除できない設定です。"

    # 実削除
    st.session_state["proposal_revisions"] = [r for r in revs if r.get("rev_id") != rev_id]

    # pending apply が削除対象なら消す
    if st.session_state.get("pending_apply_rev_id") == rev_id:
        st.session_state.pop("pending_apply_rev_id", None)

    # activeが削除対象なら付け替え（最新 or default）
    if st.session_state.get("active_rev_id") == rev_id:
        new_revs = get_revisions()
        new_active = new_revs[-1]["rev_id"] if new_revs else None
        st.session_state["active_rev_id"] = new_active
        if new_active:
            apply_revision_to_session(find_revision(new_active))

    # 比較UIは初期化し直し
    st.session_state["compare_initialized"] = False
    st.session_state.pop("compare_left_label", None)
    st.session_state.pop("compare_right_label", None)

    # 編集タブの selector とその周辺を必ずリセット（★追加）
    st.session_state.pop("active_revision_selector", None)
    st.session_state.pop("active_revision_selector_index", None)  # 使っていれば
    st.session_state.pop("active_revision_selector_label", None)  # 使っていれば

    # 削除したrevが選ばれていた可能性があるので、Streamlit widget状態を強制更新
    st.session_state.pop("active_revision_selector", None)
    st.session_state.pop("active_revision_selector__cache", None)  # 保険

    # ★追加：Revision一覧が変わったことを示す（編集タブのselectbox再生成用）
    st.session_state["revisions_version"] = st.session_state.get("revisions_version", 0) + 1

    # ★強制掃除：active_revision_selector 系を全部消す（残像対策の決定版）
    for k in list(st.session_state.keys()):
        if k.startswith("active_revision_selector__"):
            st.session_state.pop(k, None)

    # ★追加：Revisionがゼロになった場合は編集UIのトリガーも落とす
    if not st.session_state.get("proposal_revisions"):
        reset_proposal_state()
        st.session_state["proposal_draft_generated"] = False
        st.session_state.pop("active_rev_id", None)
        st.session_state.pop("edited_texts", None)


    return True, ""


MODE_HINTS = {
    "オリエン内容の整理": "資料（PDF/PPTX/TXT/DOCX/XLSX/ZIP）をアップロードしてください。自動整理が走ったら、右側の手入力欄で必要な追記・修正をします。",
    "proposal_draft": "まずは「課題ピボット」→「生成・比較」→「編集・PPT反映」の順で進めると迷いません。",
}

def switch_mode(next_mode: str):
    if st.session_state.get("selected_mode") == "proposal_draft":
        st.session_state.pop("__proposal_draft_hydrated", None)

    st.session_state["selected_mode"] = next_mode
    #st.session_state["__ui_mode_hint"] = MODE_HINTS.get(next_mode, "")
    if next_mode == "problem_reframe":
        hydrate_reframe_ui_from_data_if_empty()
    st.rerun()






def sync_orien_from_ui():
    # UI -> 永続データ
    st.session_state["data_orien_outline_ai_draft"] = st.session_state.get("ui_orien_outline_ai_draft", "")
    st.session_state["data_orien_outline_manual"] = st.session_state.get("ui_orien_outline_manual", "")

    # 既存の後段参照キーもここで統一して作る（手入力最優先）
    manual = (st.session_state["data_orien_outline_manual"] or "").strip()
    ai = (st.session_state["data_orien_outline_ai_draft"] or "").strip()
    st.session_state["orien_outline_text"] = manual if manual else ai


import pandas as pd
from pathlib import Path
import streamlit as st


@st.cache_data(show_spinner=False)
def load_case_db_from_upload(uploaded_file) -> pd.DataFrame:
    enc_candidates = ["utf-8-sig", "cp932", "utf-8"]
    last_err = None

    for enc in enc_candidates:
        try:
            uploaded_file.seek(0)
            df = pd.read_csv(uploaded_file, encoding=enc)
            break
        except Exception as e:
            last_err = e
            df = None

    if df is None:
        raise last_err

    # NaNを潰す（検索・embedding前提）
    for c in df.columns:
        if df[c].dtype == "object":
            df[c] = df[c].fillna("")

    # 日付列があればdatetime化（任意）
    col_date = "提案日（企画書の日付）"
    if col_date in df.columns:
        df[col_date] = pd.to_datetime(df[col_date], errors="coerce")

    return df


def filter_case_db(df: pd.DataFrame, *, industry: str, method: str, region: str, keyword: str) -> pd.DataFrame:
    out = df

    def _apply_eq(col: str, val: str) -> None:
        nonlocal out
        if val and val != "（指定なし）" and col in out.columns:
            out = out[out[col].astype(str) == val]

    _apply_eq("業界", industry)
    _apply_eq("調査手法", method)
    _apply_eq("調査地域", region)

    # キーワードは “複数列を横断してcontains”
    if keyword.strip():
        key = keyword.strip()
        cols = [
            "企画書タイトル", "クライアント", "クライアントのビジネス課題",
            "リサーチ目的", "リサーチの問い", "調査対象", "対象者条件", "集計分析仕様"
        ]
        cols = [c for c in cols if c in out.columns]
        if cols:
            mask = False
            for c in cols:
                mask = mask | out[c].astype(str).str.contains(key, case=False, na=False)
            out = out[mask]

    # 並び順：新しい順（提案日がある場合）
    col_date = "提案日（企画書の日付）"
    if col_date in out.columns:
        out = out.sort_values(col_date, ascending=False, na_position="last")

    return out



def render_case_review_tab(df):
    import streamlit as st

    st.session_state.setdefault("case_refs", [])

    if df is None or getattr(df, "empty", True):
        st.info("過去事例DBが読み込まれていません。")
        return

    # -----------------------------
    # フィルタUI
    # -----------------------------
    col_f1, col_f2, col_f3, col_f4 = st.columns([1, 1, 1, 2], gap="small")

    def uniq(col):
        if col in df.columns:
            return sorted([str(x) for x in df[col].unique() if str(x).strip()])
        return []

    with col_f1:
        industry = st.selectbox("業界", ["（指定なし）"] + uniq("業界"), key="case_filter_industry")
    with col_f2:
        method = st.selectbox("調査手法", ["（指定なし）"] + uniq("調査手法"), key="case_filter_method")
    with col_f3:
        region = st.selectbox("調査地域", ["（指定なし）"] + uniq("調査地域"), key="case_filter_region")
    with col_f4:
        keyword = st.text_input(
            "キーワード（タイトル／課題／目的／問い／対象などを横断検索）",
            value="",
            key="case_filter_keyword",
        )

    # -----------------------------
    # フィルタ処理
    # -----------------------------
    filtered = df.copy()

    def apply_eq(col, val):
        nonlocal filtered
        if val and val != "（指定なし）" and col in filtered.columns:
            filtered = filtered[filtered[col].astype(str) == val]

    apply_eq("業界", industry)
    apply_eq("調査手法", method)
    apply_eq("調査地域", region)

    if keyword.strip():
        cols = [
            "企画書タイトル",
            "クライアント",
            "クライアントのビジネス課題",
            "リサーチ目的",
            "リサーチの問い",
            "調査対象",
            "対象者条件",
            "集計分析仕様",
        ]
        cols = [c for c in cols if c in filtered.columns]
        if cols:
            mask = False
            for c in cols:
                mask = mask | filtered[c].astype(str).str.contains(keyword, case=False, na=False)
            filtered = filtered[mask]

    # 並び順（新しい順）
    date_col = "提案日（企画書の日付）"
    if date_col in filtered.columns:
        filtered = filtered.sort_values(date_col, ascending=False, na_position="last")

    st.caption(f"ヒット件数：{len(filtered)}")
    if filtered.empty:
        st.info("条件に合う事例がありません。")
        return

    # -----------------------------
    # 一覧表示
    # -----------------------------
    view_cols = [
        c for c in [
            "提案日（企画書の日付）",
            "業界",
            "クライアント",
            "企画書タイトル",
            "調査手法",
            "調査対象",
            "サンプルサイズ",
        ]
        if c in filtered.columns
    ]

    st.dataframe(filtered[view_cols].head(200), use_container_width=True, hide_index=True)

    # -----------------------------
    # 詳細表示
    # -----------------------------
    options = filtered.index.tolist()
    sel = st.selectbox(
        "詳細表示する事例を選択",
        options=options,
        format_func=lambda i: filtered.at[i, "企画書タイトル"]
        if "企画書タイトル" in filtered.columns else str(i),
        key="case_detail_selector",
    )

    row = filtered.loc[sel].to_dict()

    st.markdown("#### 事例詳細")

    important_cols = [
        "企画書タイトル",
        "提案日（企画書の日付）",
        "業界",
        "クライアント",
        "クライアントのビジネス課題",
        "リサーチ目的",
        "リサーチの問い",
        "リサーチ仮説",
        "調査対象",
        "対象者条件",
        "調査手法",
        "集計分析仕様",
        "納品物",
    ]

    for k in important_cols:
        v = row.get(k)
        if v is not None and str(v).strip():
            st.markdown(f"**{k}**")
            st.write(v)

    # -----------------------------
    # 参考リスト操作
    # -----------------------------
    col_a, col_b = st.columns([1, 3])

    with col_a:
        if st.button("この事例を参考に入れる", use_container_width=True, key="btn_add_case_ref"):
            key_id = (row.get("ファイル名") or row.get("企画書タイトル") or "").strip()
            existing = {
                (x.get("ファイル名") or x.get("企画書タイトル") or "").strip()
                for x in st.session_state["case_refs"]
            }
            if key_id and key_id in existing:
                st.info("既に参考リストに入っています。")
            else:
                st.session_state["case_refs"].append(row)
                st.success("参考リストに追加しました。")

    with col_b:
        with st.expander("参考リスト（このプロジェクト内）", expanded=False):
            refs = st.session_state.get("case_refs", [])
            st.caption(f"{len(refs)}件")
            for i, r in enumerate(refs, 1):
                st.markdown(f"{i}. **{r.get('企画書タイトル','')}**（{r.get('クライアント','')}）")

            if refs and st.button("参考リストをクリア", key="btn_clear_case_refs"):
                st.session_state["case_refs"] = []
                st.rerun()


def build_case_search_text(row: dict) -> str:
    cols = [
        "業界","クライアント","企画書タイトル",
        "クライアントのビジネス目標","クライアントのビジネス現状","クライアントのビジネス課題",
        "リサーチ目的","リサーチの問い","リサーチ仮説",
        "調査対象","対象者条件","調査手法","抽出方法","調査地域",
        "集計分析仕様","自由回答データの処理","納品物",
    ]
    parts = []
    for c in cols:
        v = (row.get(c) or "")
        if isinstance(v, float) and pd.isna(v):
            v = ""
        v = str(v).strip()
        if v:
            parts.append(f"{c}:{v}")
    return "\n".join(parts)


import numpy as np

@st.cache_data(show_spinner=False)
def load_case_vectors(path: Path) -> pd.DataFrame:
    df = pd.read_parquet(path)
    # vector列をnp配列化（計算を速くする）
    df["vector_np"] = df["vector"].apply(lambda x: np.asarray(x, dtype=np.float32))
    return df

def build_current_case_query_text() -> str:
    axis_text = (st.session_state.get("fullgen_axis_text") or "").strip()
    kickoff = {
        "目標": st.session_state.get("ai_目標",""),
        "現状": st.session_state.get("ai_現状",""),
        "ビジネス課題": st.session_state.get("ai_ビジネス課題",""),
        "調査目的": st.session_state.get("ai_調査目的",""),
        "問い": st.session_state.get("ai_問い",""),
        "仮説": st.session_state.get("ai_仮説",""),
    }
    subq = st.session_state.get("subq_list", []) or []
    return "\n".join([
        "【軸】", axis_text,
        "【キックオフ】", str(kickoff),
        "【サブクエスチョン】", str(subq),
    ]).strip()

def cosine_np(a: np.ndarray, b: np.ndarray) -> float:
    na = np.linalg.norm(a); nb = np.linalg.norm(b)
    if na == 0 or nb == 0:
        return 0.0
    return float(np.dot(a, b) / (na * nb))

def retrieve_similar_cases(df_vec: pd.DataFrame, query_vec: np.ndarray, top_k: int = 5) -> list[dict]:
    scores = []
    for i, r in df_vec.iterrows():
        s = cosine_np(query_vec, r["vector_np"])
        scores.append((s, i))
    scores.sort(reverse=True, key=lambda x: x[0])
    out = []
    for s, i in scores[:top_k]:
        row = df_vec.loc[i].drop(labels=["vector_np"]).to_dict()
        row["_score"] = s
        out.append(row)
    return out


def format_case_for_prompt(row: dict) -> str:
    keys = ["業界","クライアント","企画書タイトル","リサーチ目的","リサーチの問い","リサーチ仮説","調査対象","対象者条件","調査手法","集計分析仕様"]
    lines = []
    for k in keys:
        v = (row.get(k) or "").strip()
        if v:
            lines.append(f"- {k}: {v}")
    return "\n".join(lines)

def build_case_context_for_prompt(max_manual: int = 2, max_auto: int = 2) -> str:
    manual = st.session_state.get("case_refs", []) or []
    manual = manual[:max_manual]

    auto = st.session_state.get("auto_case_recs", []) or []
    auto = auto[:max_auto]

    if not manual and not auto:
        return ""

    parts = ["【参考：過去事例（要点）】", "※今回の軸に整合する範囲で“構造”を参考にし、文言のコピペは避けてください。"]
    if manual:
        parts.append("▼手動選択")
        for i, r in enumerate(manual, 1):
            parts.append(f"事例M{i}\n{format_case_for_prompt(r)}")
    if auto:
        parts.append("▼類似レコメンド")
        for i, r in enumerate(auto, 1):
            parts.append(f"事例A{i}（score={r.get('_score',0):.3f}）\n{format_case_for_prompt(r)}")
    return "\n\n".join(parts).strip()


def render_case_db_uploader_sidebar():
    import streamlit as st

    uploaded = st.file_uploader(
        "過去企画の要点CSVをアップロード",
        type=["csv"],
        key="case_db_upload_sidebar",  # ★ sidebar用にkeyを変える（重複回避）
    )

    if uploaded is None:
        st.caption("未読み込み：過去事例レビュータブで参照できます。")
        return None

    sig = (uploaded.name, getattr(uploaded, "size", None))
    if st.session_state.get("case_db_upload_sig") != sig:
        st.session_state["case_db_upload_sig"] = sig
        try:
            df = load_case_db_from_upload(uploaded)
            st.session_state["case_db_df"] = df
            st.success(f"読み込み完了（{len(df)}件）")
        except Exception as e:
            st.session_state.pop("case_db_df", None)
            st.error(f"CSVの読み込みに失敗しました: {e}")
            return None

    return st.session_state.get("case_db_df")



import pandas as pd
import streamlit as st

@st.cache_data(show_spinner=False)
def load_case_db_from_upload(uploaded_file) -> pd.DataFrame:
    enc_candidates = ["utf-8-sig", "cp932", "utf-8"]
    last_err = None

    for enc in enc_candidates:
        try:
            uploaded_file.seek(0)
            df = pd.read_csv(uploaded_file, encoding=enc)
            break
        except Exception as e:
            last_err = e
            df = None

    if df is None:
        raise last_err

    # 文字列NaNを潰す（検索で事故りやすいので）
    for c in df.columns:
        if df[c].dtype == "object":
            df[c] = df[c].fillna("")

    # 日付列（存在するなら）をdatetime化
    col_date = "提案日（企画書の日付）"
    if col_date in df.columns:
        df[col_date] = pd.to_datetime(df[col_date], errors="coerce")

    return df


def render_case_review_screen():
    import streamlit as st

    st.markdown("# 過去事例レビュー")

    # DBが未読み込みなら案内
    df = st.session_state.get("case_db_df")
    if df is None:
        st.info("左ペインで過去事例DB（CSV）をアップロードしてください。")
        return

    # # （任意）戻るボタン
    # col1, col2 = st.columns([1, 5])
    # with col1:
    #     if st.button("← 戻る", use_container_width=True):
    #         st.session_state["selected_mode"] = "proposal_draft"
    #         st.rerun()

    # 既存のレビューUI（uploader無し・df引数版）
    render_case_review_tab(df)




def maybe_autosave_active_revision() -> None:
    """Tab3の“毎回save”を1回だけ抑止するためのラッパー。"""
    if st.session_state.pop("__skip_autosave_once", False):
        return
    if get_active_revision() is None:
        return
    save_session_keys_to_active_revision()


# =========================================================
# Pending project apply (MUST run before any widgets)
# =========================================================
if "pending_project" in st.session_state:
    try:
        reset_proposal_state()
        apply_project_to_session(st.session_state["pending_project"])
    finally:
        st.session_state.pop("pending_project", None)
    st.rerun()

if st.session_state.get("pending_apply_rev_id"):
    rid = st.session_state.pop("pending_apply_rev_id")
    set_active_revision(rid)
    rev = get_active_revision()
    if rev:
        apply_revision_to_session(rev)
    st.rerun()



#render_dbg_sidebar()


# パワポフォーマット連携
import os
from pathlib import Path

import streamlit as st
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient


def get_setting(key: str) -> str:
    # 1) App Service のアプリ設定（環境変数）を最優先
    v = os.getenv(key)
    if v:
        return v

    # 2) ローカル実行用：secrets.toml がある場合だけ使う
    try:
        return st.secrets[key]
    except Exception:
        raise RuntimeError(
            f"Missing setting: {key}. "
            f"Set it as an App Service app setting (recommended) or in .streamlit/secrets.toml."
        )


@st.cache_resource
def get_blob_service_client():
    account_url = get_setting("STORAGE_ACCOUNT_URL")
    return BlobServiceClient(
        account_url=account_url,
        credential=DefaultAzureCredential(),
    )


def ensure_server_template_loaded():
    # すでにロード済みなら何もしない
    if st.session_state.get("template_loaded") and st.session_state.get("pptx_path"):
        return

    container = get_setting("TEMPLATE_CONTAINER")
    blob_name = get_setting("TEMPLATE_BLOB_NAME")

    bsc = get_blob_service_client()
    blob = bsc.get_blob_client(container=container, blob=blob_name)

    pptx_bytes = blob.download_blob().readall()

    session_dir = get_session_dir()  # 既存関数（Pathを返す前提）
    tpl_dir = session_dir / "pptx"
    tpl_dir.mkdir(parents=True, exist_ok=True)

    # 保存ファイル名は固定にするのが安全（日本語blob名でも事故らない）
    target = tpl_dir / "template.pptx"
    target.write_bytes(pptx_bytes)

    st.session_state["pptx_path"] = str(target)
    st.session_state["template_loaded"] = True

import io
import json
import zipfile
from datetime import datetime

import pandas as pd
import streamlit as st


def _make_survey_items_excel_bytes(session_state) -> bytes:
    rows = session_state.get("survey_item_rows", []) or []
    df = pd.DataFrame(rows)

    # 期待列が無いケースに備えて最低限の列を用意（あなたの列設計に寄せる）
    base_cols = [
        "sq_id", "sq_subq", "items", "approach",
        "var_name", "item_text", "recommended_type", "recommended_scale",
        "priority", "table_role", "is_selected"
    ]
    for c in base_cols:
        if c not in df.columns:
            df[c] = ""

    # bool正規化
    if "is_selected" in df.columns:
        df["is_selected"] = df["is_selected"].apply(
            lambda x: True if str(x).strip() in ["", "True", "true", "1", "yes", "Yes", "Y", "y"] else False
        )

    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="openpyxl") as writer:
        df.to_excel(writer, sheet_name="all_items", index=False)

        sel = df[df["is_selected"] == True].copy()
        sel.to_excel(writer, sheet_name="selected_only", index=False)

        # SQ別（見やすさ用）
        if "sq_id" in df.columns:
            df_sorted = df.sort_values(["sq_id", "priority"], ascending=[True, True])
            df_sorted.to_excel(writer, sheet_name="by_sq", index=False)

    return out.getvalue()


def _make_project_json_bytes(session_state) -> bytes:
    proj = build_project_from_session()  # ←あなたの既存関数
    s = json.dumps(proj, ensure_ascii=False, indent=2)
    return s.encode("utf-8")


def _make_output_zip_bytes(session_state) -> tuple[bytes, str]:
    """
    returns: (zip_bytes, zip_filename)
    """
    # ZIP名（プロジェクト名＋日時）
    project_name = (session_state.get("project_name") or "project").strip() or "project"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    zip_name = f"{project_name}_{ts}.zip"

    # 1) Excel
    excel_bytes = _make_survey_items_excel_bytes(session_state)

    # 2) JSON
    json_bytes = _make_project_json_bytes(session_state)

    # 3) PPT（反映済みPPTを生成）
    # run_step4_apply_current_ui_to_ppt が (out_path, report) を返す想定なので
    # out_path を読み込んで bytes化します
    out_path, report = run_step4_apply_current_ui_to_ppt(session_state)  # ←あなたの既存関数
    ppt_bytes = out_path.read_bytes()
    ppt_name = out_path.name  # 既存のファイル名をそのままZIPへ

    # ZIP化
    zbuf = io.BytesIO()
    with zipfile.ZipFile(zbuf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{project_name}/survey_items.xlsx", excel_bytes)
        zf.writestr(f"{project_name}/{project_name}.json", json_bytes)
        zf.writestr(f"{project_name}/{ppt_name}", ppt_bytes)

        # 任意：反映レポートも入れると便利
        rep = json.dumps(report, ensure_ascii=False, indent=2).encode("utf-8")
        zf.writestr(f"{project_name}/ppt_apply_report.json", rep)

    return zbuf.getvalue(), zip_name


import os
from azure.storage.blob import BlobServiceClient

def ensure_pptx_path() -> bool:
    """
    st.session_state['pptx_path'] が無ければ Blob から /tmp/template.pptx にDLしてセット。
    """
    if st.session_state.get("pptx_path"):
        return True

    conn_str = os.getenv("AZURE_STORAGE_CONNECTION_STRING", "")
    container = os.getenv("PPTX_TEMPLATE_CONTAINER", "")
    blob_name = os.getenv("PPTX_TEMPLATE_BLOB_NAME", "")  # 例: templates/template.pptx

    if not conn_str or not container or not blob_name:
        return False

    try:
        bsc = BlobServiceClient.from_connection_string(conn_str)
        bc = bsc.get_blob_client(container=container, blob=blob_name)
        data = bc.download_blob().readall()

        local_path = "/tmp/template.pptx"
        with open(local_path, "wb") as f:
            f.write(data)

        st.session_state["pptx_path"] = local_path
        return True

    except Exception:
        return False


# =========================================================
# pending apply（ウィジェット生成前にrev内容をsessionへ適用）
# =========================================================
# pending_id = st.session_state.pop("pending_apply_rev_id", None)
# if pending_id:
#     rev = find_revision(pending_id)
#     if rev:
#         apply_revision_to_session(rev)  # ←ここなら安全（まだウィジェットが作られていない）



# =========================
# レイアウト構成（左ペイン＋中央ペインのみ）
# =========================
left, center = st.columns([1, 3], gap="large")

if "selected_mode" not in st.session_state:
    st.session_state["selected_mode"] = "オリエン内容の整理"


# =========================
# 左ペイン
# =========================
with left:

    #ensure_server_template_loaded()


    if st.button("オリエン内容の整理", use_container_width=True):
        switch_mode("オリエン内容の整理")

    #if st.button("業界/ブランド診断", use_container_width=True):
    #    switch_mode("brand_diagnosis")

    if st.button("企画書下書き", use_container_width=True):
        switch_mode("proposal_draft")

    #if st.button("過去事例レビュー", use_container_width=True):
    #    st.session_state["selected_mode"] = "case_review"
    #    st.rerun()    

    st.divider()

    # 保存済みプロジェクトの読み込み
    uploaded_proj = st.file_uploader(
        "保存済みファイル読み込み",
        type=["json"],
        key="project_json_upload",
    )

    if uploaded_proj is not None:
        if st.button("このプロジェクトを読み込んで再編集を開始", use_container_width=True):
            try:
                raw = uploaded_proj.getvalue()  # read() ではなく getvalue()
                proj_loaded = json.loads(raw.decode("utf-8-sig"))

                # ここでは apply せず、「予約」だけして rerun
                st.session_state["pending_project"] = proj_loaded
                st.rerun()

            except Exception as e:
                st.error(f"プロジェクト読み込み中にエラーが発生しました: {e}")
   
    # # PPTテンプレートアップロード（セッション専用ディレクトリに保存）
    # uploaded_pptx = st.file_uploader(
    #     "企画書テンプレートをアップロードしてください（PPTX）",
    #     type=["pptx"],
    #     key="pptx_upload",
    # )

    # # アップロードの変化検知（同名でも内容/サイズが変われば再保存）
    # if "pptx_upload_sig" not in st.session_state:
    #     st.session_state["pptx_upload_sig"] = None

    # if uploaded_pptx is not None:
    #     sig = (uploaded_pptx.name, getattr(uploaded_pptx, "size", None))

    #     # 初回 or 差し替え時のみ保存し直す
    #     if st.session_state["pptx_upload_sig"] != sig:
    #         st.session_state["pptx_upload_sig"] = sig

    #         session_dir = get_session_dir()
    #         tpl_dir = session_dir / "pptx"
    #         tpl_dir.mkdir(parents=True, exist_ok=True)

    #         target = tpl_dir / uploaded_pptx.name
    #         with open(target, "wb") as f:
    #             f.write(uploaded_pptx.getbuffer())

    #         st.session_state["pptx_path"] = str(target)
    #         st.session_state["template_loaded"] = True

    #         # 既に反映済みPPTがある場合はクリア（取り違え防止）
    #         st.session_state.pop("final_pptx_bytes", None)
    #         st.session_state.pop("final_pptx_name", None)
    #         st.session_state.pop("final_pptx_key", None)

    #         st.success(f"{uploaded_pptx.name} を読み込みました（保存先: {target}）。")

    # 過去企画のアップロード（セッション専用ディレクトリに保存）
    # df = render_case_db_uploader_sidebar()
    # sig = st.session_state.get("case_db_upload_sig")
    # if sig:
    #     name, size = sig
    #     st.caption(f"読み込み済み: {name}")


# =========================
# 中央ペイン
# =========================
# ★モード切替時のガイド（1回だけ表示）
hint = st.session_state.pop("__ui_mode_hint", None)
# if hint:
#     render_character_guide("次にやること", hint, img_width=120, kind="info")


with center:
    mode = st.session_state.get("selected_mode")

    # ★初回アクセス時だけガイド表示（初動の心理的壁を下げる）
    if "__ui_first_visit" not in st.session_state:
        render_character_guide(
            "みなの知恵へ、ようこそ！",
            "現在このツールでは企画書作成をサポートしているよ。作成できる企画書のスライドは以下です。\n"
            "- キックオフノート（KON）\n"
            "- サブクエスチョン（SQ）\n"
            "- 分析アプローチ\n"
            "- 調査対象者案\n"
            "- 調査項目案\n\n\n"
            "次の手順で進めていこう！\n"
            "1) 左の「オリエン内容の整理」で資料をアップロード\n"
            "2) 内容を確認したら「企画書下書き」へ進んでください。\n"
            "3) 企画を一時保存してある場合は「保存済みファイル読み込み」から再開できます。\n",
            img_width=300,
            kind="info",
        )
        st.session_state["__ui_first_visit"] = True


    #WATCH_PREFIXES = ("pivot_", "reframe_", "proposal_", "kickoff_", "analysis_", "ai_", "survey_", "problem_", "orien_")
    WATCH_PREFIXES = (
        "pivot_", "reframe_", "proposal_", "kickoff_", "analysis_", "ai_", "survey_", "problem_", "orien_",
        "draft_", "edit_", "EDIT_", "ppt_", "pptx_", "slide_", "tab_", "gen_", "cmp_", "rev_"
    )


    def _snap(prefixes=WATCH_PREFIXES) -> dict:
        ss = st.session_state
        out = {}
        for k in list(ss.keys()):
            if k.startswith(prefixes):
                v = ss.get(k)
                if isinstance(v, str):
                    out[k] = f"str({len(v)})"
                elif isinstance(v, list):
                    out[k] = f"list({len(v)})"
                elif isinstance(v, dict):
                    out[k] = f"dict({len(v)})"
                else:
                    out[k] = type(v).__name__
        return out

    def _diff(before: dict, after: dict) -> dict:
        b = set(before.keys())
        a = set(after.keys())
        removed = sorted(list(b - a))
        added = sorted(list(a - b))
        changed = sorted([k for k in (b & a) if before.get(k) != after.get(k)])
        return {"removed": removed[:80], "added": added[:80], "changed": changed[:80]}

    # モード変化検出
    if "__dbg_prev_mode" not in st.session_state:
        st.session_state["__dbg_prev_mode"] = mode
        st.session_state["__dbg_prev_snap"] = _snap()

    if mode != st.session_state.get("__dbg_prev_mode"):
        before = st.session_state.get("__dbg_prev_snap", {})
        after = _snap()
        st.session_state["__dbg_mode_change"] = {
            "from": st.session_state.get("__dbg_prev_mode"),
            "to": mode,
            "diff": _diff(before, after),
        }
        st.session_state["__dbg_prev_mode"] = mode
        st.session_state["__dbg_prev_snap"] = after

    # dbg = st.session_state.get("__dbg_mode_change")
    # if dbg:
    #     st.warning(f"DEBUG mode change: {dbg['from']} -> {dbg['to']}")
    #     st.json(dbg["diff"])






    # ----------------------------------------
    # オリエン内容の整理（単独モード）
    # ----------------------------------------
    if mode == "オリエン内容の整理":
        st.markdown("## オリエン内容の整理")


        render_character_guide2(
            "ここではクライアントと話した内容を整理しよう",
            "- 下から オリエンのファイルを読み込んでください。自動で内容を整理してくれるよ。\n"
            "- 保存は不要です。\n"
            "- 完了したら 「企画書下書き」 に進んでください。",
            img_width=500,
            kind="info",
        )
        # st.caption(
        #     f"DEBUG orien_ai_draft len={len(st.session_state.get('data_orien_outline_ai_draft',''))} / "
        #     f"manual len={len(st.session_state.get('data_orien_outline_manual',''))} / "
        #     f"uploaded_docs={len(st.session_state.get('uploaded_docs',[]) or [])}"
        # )



        uploaded_files = st.file_uploader(
            "オリエン資料をアップロードしてください（PDF / PPTX / TXT / DOCX / XLSX / ZIP）",
            type=["pdf", "pptx", "txt", "docx", "xlsx", "zip"],
            accept_multiple_files=True,
            key="orien_upload",
        )
        process_orien_upload(uploaded_files)

        # アップロードが入ったら自動で整理（ただし手入力は上書きしない）
        if uploaded_files and (not st.session_state.get("orien_auto_generated")):
            with st.spinner("オリエン内容を所定フォームで整理しています..."):
                ok, msg = ensure_orien_outline()
            if ok:
                st.session_state["orien_auto_generated"] = True
                st.rerun()
            else:
                st.error(msg)

        # 整理結果（AI）と手入力欄を常に表示
        col1, col2 = st.columns(2, gap="large")

        with col1:
            # 表示直前に data -> ui を注入（UIが消えても data は残る）
            st.session_state["ui_orien_outline_ai_draft"] = st.session_state.get("data_orien_outline_ai_draft", "")

            st.text_area(
                "整理結果（所定フォーム）",
                key="ui_orien_outline_ai_draft",
                height=800,
                on_change=sync_orien_from_ui,
            )

        with col2:
            st.session_state["ui_orien_outline_manual"] = st.session_state.get("data_orien_outline_manual", "")

            st.text_area(
                "手入力（最優先：追記・修正）",
                key="ui_orien_outline_manual",
                height=800,
                placeholder="ここに補足・修正を入力すると、この内容が最優先で後工程に反映されます。",
                on_change=sync_orien_from_ui,
            )

        # 統合用の正本を同期（後段処理の参照先）
        sync_orien_from_ui()



    # ----------------------------------------
    # ブランド診断モード（中央ペインに集約）
    # ----------------------------------------
    elif mode == "brand_diagnosis":
        st.markdown("## 業界/ブランド診断")
        st.caption("オリエン資料から業界・ブランドを推測し、市場特性を整理します。")

        ori_texts = get_orien_context_text()

        # 初期値
        st.session_state.setdefault("target_category", "")
        st.session_state.setdefault("target_brand", "")

        # カテゴリー・ブランド推測
        if st.button("業界/ブランドを推測", use_container_width=True):
            if not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("カテゴリーとブランドを推測中..."):
                    prompt = f"""
あなたは市場調査の専門家です。
以下のオリエン資料から、今回の調査対象となるカテゴリー（市場）とブランド名を推定してください。

【出力形式】
カテゴリー（市場）:
ブランド:

資料:
{ori_texts[:4000]}
"""
                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.5,
                            max_tokens=200,
                        )
                        ai_result = response.choices[0].message.content
                        cat_match = re.search(r"カテゴリー（市場）[:：]\s*(.*)", ai_result)
                        brand_match = re.search(r"ブランド[:：]\s*(.*)", ai_result)

                        st.session_state["target_category"] = cat_match.group(1).strip() if cat_match else ""
                        st.session_state["target_brand"] = brand_match.group(1).strip() if brand_match else ""
                        st.success("カテゴリーとブランドを抽出しました。")
                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")

        # 手動編集欄
        st.text_input(
            "対象カテゴリー（市場）",
            key="target_category",
            placeholder="例：清涼飲料、化粧品、通信キャリアなど",
        )
        st.text_input(
            "対象ブランド",
            key="target_brand",
            placeholder="例：キッザニア、SUUMO、カローラ など",
        )

        st.markdown("---")
        st.markdown("### カテゴリー・ブランドについて検索")

        if st.button("カテゴリー・ブランドについて検索", use_container_width=True):
            cat = st.session_state.get("target_category", "")
            brand = st.session_state.get("target_brand", "")
            if not cat:
                st.warning("カテゴリーを入力してください。")
            else:
                with st.spinner("市場特性を検索中..."):
                    prompt = f"""
あなたは市場分析の専門家です。
次のカテゴリーとブランドに関する市場構造と消費行動特性を整理してください。

【カテゴリー】{cat}
【ブランド】{brand}

出力は以下の2表形式のMarkdownで記載してください。
# カテゴリーに関する検索項目
|項目|内容|
|市場タイプ|...|
...

# カテゴリーの消費行動特性
|項目|内容|
|検討期間|...|
...
"""
                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場分析の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=900,
                        )
                        result = response.choices[0].message.content

                        def extract_md_table(md_text, header):
                            if header in md_text:
                                section = md_text.split(header, 1)[1]
                                table_part = section.split("#")[0]
                                rows = [
                                    ln.strip()
                                    for ln in table_part.splitlines()
                                    if "|" in ln and not ln.startswith("|項目|----|")
                                ]
                                data = []
                                for ln in rows:
                                    cols = [c.strip() for c in ln.strip("|").split("|")]
                                    if len(cols) >= 2:
                                        data.append(cols[:2])
                                if data:
                                    df = pd.DataFrame(data[1:], columns=data[0])
                                    return df
                            return pd.DataFrame(columns=["項目", "内容"])

                        st.session_state["df_category_structure"] = extract_md_table(result, "# カテゴリーに関する検索項目")
                        st.session_state["df_behavior_traits"] = extract_md_table(result, "# カテゴリーの消費行動特性")

                        # マーケティングファネルも合わせて生成
                        with st.spinner("マーケティングファネルを生成中..."):
                            prompt_funnel = f"""
あなたはブランドマーケティングの専門家です。
以下のカテゴリー・ブランドについて、認知から再接点・ロイヤリティまでのマーケティングファネルを
トリガー／障壁の観点で階層構造として整理してください。

【カテゴリー】{cat}
【ブランド】{brand}
"""
                            response_funnel = client.chat.completions.create(
                                model=DEPLOYMENT,
                                messages=[
                                    {"role": "system", "content": "あなたはブランドマーケティングの専門家です。"},
                                    {"role": "user", "content": prompt_funnel},
                                ],
                                temperature=0.6,
                                max_tokens=1800,
                            )
                            st.session_state["funnel_text"] = response_funnel.choices[0].message.content

                        st.success("市場特性とマーケティングファネルを整理しました。")

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")

        # 結果表示
        if "df_category_structure" in st.session_state:
            st.markdown("### カテゴリーに関する検索項目")
            st.data_editor(
                st.session_state["df_category_structure"],
                hide_index=True,
                num_rows="fixed",
                use_container_width=True,
            )

        if "df_behavior_traits" in st.session_state:
            st.markdown("### カテゴリーの消費行動特性")
            st.data_editor(
                st.session_state["df_behavior_traits"],
                hide_index=True,
                num_rows="fixed",
                use_container_width=True,
            )

        if st.session_state.get("funnel_text"):
            st.markdown("### マーケティングファネル（トリガー／障壁）")
            st.code(st.session_state["funnel_text"], language="text")


    elif mode == "proposal_draft":

        ensure_revision_store()

        # tabsより前（proposal_draft の中）
        if not ensure_pptx_path():
            st.warning("PPTテンプレート（Blob）の取得に失敗しました。環境変数設定を確認してください。")

        # =========================================================
        # active_rev_id が変わったら、必ずそのrevを session_state に適用する
        # （apply前に save が走って巻き戻るのを防ぐ）
        # =========================================================
        active_now = st.session_state.get("active_rev_id")
        last_applied = st.session_state.get("__last_applied_rev_id")

        if active_now and active_now != last_applied:
            rev = find_revision(active_now)
            if rev:
                apply_revision_to_session(rev)
                st.session_state["__last_applied_rev_id"] = active_now
                # ★次の描画1回だけ自動保存を止める（巻き戻り防止）
                st.session_state["__skip_autosave_once"] = True


        # -------------------------------------------------
        # 0) active_rev_id を確実に用意
        # -------------------------------------------------
        if not st.session_state.get("active_rev_id"):
            revs = get_revisions() or []
            if revs:
                st.session_state["active_rev_id"] = revs[-1].get("rev_id")
            else:
                rev = create_new_revision()
                upsert_revision(rev)
                st.session_state["active_rev_id"] = rev["rev_id"]

        # -------------------------------------------------
        # 1) KON～SQ からの「採用して反映」予約を最優先で処理（UI生成より前）
        #    ※ pending_apply_rev_id はここで 1回だけ pop する
        # -------------------------------------------------
        # pending_id = st.session_state.pop("pending_apply_rev_id", None)
        # if pending_id:
        #     rev = find_revision(pending_id)
        #     if rev:
        #         st.session_state["active_rev_id"] = pending_id
        #         apply_revision_to_session(rev)

        #         # hydration 管理をリセット（後段の保険）
        #         st.session_state["__hydrated_rev_id"] = pending_id
        #     else:
        #         st.warning("採用したRevisionが見つかりませんでした（削除された可能性があります）。")

        # -------------------------------------------------
        # 2) hydration（active_rev_id が変わった時だけ apply）
        #    ＝初回 or active切替時に必ず session_state を同期
        # -------------------------------------------------
        active_id = st.session_state.get("active_rev_id")
        if st.session_state.get("__hydrated_rev_id") != active_id:
            rev = get_active_revision()
            if rev:
                apply_revision_to_session(rev)
            st.session_state["__hydrated_rev_id"] = active_id

        # -------------------------------------------------
        # 3) ここから下で UI を描き始める
        # -------------------------------------------------
        st.markdown("## 企画書下書き")
        st.session_state["proposal_draft_generated"] = bool(get_revisions())

        tab_pivot, tab_gen, tab_edit = st.tabs(["課題ピボット", "KON～SQ", "分析イメージ"])



        # =========================================================
        # TAB 1: 課題ピボット
        # =========================================================
        with tab_pivot:

            render_character_guide3(
                "課題のピボット",
                "- クライアントが言ったことを一度立ち止まって考えるステップだよ。\n"
                "- クライアントはなぜその課題をみんなに投げかけてきたのかな？\n"
                "- 「新規作成」ボタンを押して下さい。オリエン内容をもとに課題の背景を考察するよ。\n"
                "- ここでは眺めるだけでOK。もし自分なりの考えがあれば「クライアント課題（手書き）」に直接記入してください。\n"
                "- 確認が済んだらページ下にある 「確認完了」 を押してください。\n"
                "- 次は「KON～SQ」タブに進みます。",
                img_width=500,
                kind="info",
            )

            #st.markdown("### 課題のピボット")
            

            # 入力参照（読み取り専用）
            orien_outline_text = st.session_state.get("orien_outline_text", "")
            cat_df = st.session_state.get("df_category_structure")
            beh_df = st.session_state.get("df_behavior_traits")
            funnel_text = st.session_state.get("funnel_text", "")

            with st.expander("参照している前提", expanded=False):
                st.markdown("### オリエン内容の整理（抜粋）")
                st.code(orien_outline_text[:2000] if orien_outline_text else "（未生成）", language="text")

                st.markdown("### カテゴリー構造")
                if cat_df is not None and not cat_df.empty:
                    st.data_editor(cat_df, hide_index=True, num_rows="fixed", use_container_width=True)
                else:
                    st.caption("（未生成）")

                st.markdown("### 消費行動特性")
                if beh_df is not None and not beh_df.empty:
                    st.data_editor(beh_df, hide_index=True, num_rows="fixed", use_container_width=True)
                else:
                    st.caption("（未生成）")

                st.markdown("### マーケティングファネル")
                st.code(funnel_text[:2000] if funnel_text else "（未生成）", language="text")

            col_btn, col_note = st.columns([1, 3], gap="small")

            with col_btn:
                if st.button("新規作成", use_container_width=True, key="btn_premise_generate"):
                    with st.spinner("前提整理（5観点）を生成しています..."):
                        ok, msg = generate_problem_reframe_premise()

                    if ok:
                        st.success("①の前提整理ドラフトを生成しました。必要に応じて編集してください。")
                        st.rerun()
                    else:
                        st.error(msg)

            with col_note:
                st.caption("課題変換（前処理）は自動実行しません。必要なタイミングでボタン押下して生成してください。")

            st.divider()

            # =========================================================
            # ★3レイヤー（WHY / WHAT / HOW）で入力欄を構造化
            # =========================================================
            st.markdown("### 観点を3レイヤーで整理（WHY / WHAT / HOW）")
            #st.caption("上位→下位の順に埋めると、後続の『生成・比較』でブレにくくなります。")

            # WHY（上位）
            with st.container(border=True):
                st.markdown("#### WHY（上位）：事業視点 — 何が問題で、なぜ調査するのか")
                st.text_area(
                    "事業やブランドが抱える課題（根本課題）",
                    key="reframe_c4_business_brand",
                    height=110,
                )

            # WHAT（中位）
            with st.container(border=True):
                st.markdown("#### WHAT（中位）：意思決定視点 — 調査で何を明らかにし、何を判断するのか")
                st.text_area(
                    "報告先（組織長など）が知りたいこと（意思決定論点）",
                    key="reframe_c2_exec_summary",
                    height=110,
                )

            # HOW（下位）
            with st.container(border=True):
                st.markdown("#### HOW（下位）：実行視点 — 誰が何を担い、調査後にどう動くのか")
                st.text_area(
                    "調査結果を受けて次にすること（ネクストアクション）",
                    key="reframe_c1_next_action",
                    height=110,
                )

            # 任意追記（補助レイヤー）
            with st.container(border=True):
                st.markdown("#### クライアント課題（手書き）")
                st.text_area(
                    "任意の記載",
                    key="reframe_c6_user_notes",
                    height=140,
                    #placeholder="例：意思決定会議が2月中旬にある／競合Aの新商品が影響／調査対象外の制約条件／現場の肌感など",
                )

            st.markdown("")
            st.divider()

            # =========================================================
            # 「生成・比較」へ渡すための確定（反映）処理（既存のまま）
            # =========================================================

            pivot_items = [
                ("1) 調査結果を受けて次にすること（ネクストアクション）", st.session_state.get("reframe_c1_next_action", "")),
                ("2) 報告先（組織長など）が知りたいこと", st.session_state.get("reframe_c2_exec_summary", "")),
                ("3) 事業やブランドが抱える課題", st.session_state.get("reframe_c4_business_brand", "")),
                ("4) 任意の追記（補足・前提条件・懸念・別視点など）", st.session_state.get("reframe_c6_user_notes", "")),
            ]

            committed = bool(st.session_state.get("pivot_axis_texts_committed", False))
            if committed:
                st.success("課題ピボット（4観点）は『KON～SQ』に反映済みです。内容を変更した場合は、再度『反映（確定）』してください。")
            else:
                st.warning("課題ピボット（4観点）は、まだ『KON～SQ』に反映されていません。下のボタンで反映してください。")

            col_commit, col_preview = st.columns([1, 2], gap="small")

            with col_commit:
                if st.button("確認完了", use_container_width=True, key="btn_commit_pivot_axis"):
                    pivot_map = {
                        "c4_business_brand": (st.session_state.get("reframe_c4_business_brand", "") or "").strip(),
                        "c2_exec_summary": (st.session_state.get("reframe_c2_exec_summary", "") or "").strip(),
                        "c1_next_action": (st.session_state.get("reframe_c1_next_action", "") or "").strip(),
                        "c6_user_notes": (st.session_state.get("reframe_c6_user_notes", "") or "").strip(),
                    }

                    pivot_labels = {
                        "c4_business_brand": "事業やブランド視点",
                        "c2_exec_summary": "意思決定視点",
                        "c1_next_action": "実行視点",
                        "c6_user_notes": "クライアント課題（手書き）",
                    }

                    ordered_keys = [
                        "c4_business_brand",
                        "c2_exec_summary",
                        "c1_next_action",
                        "c6_user_notes",
                    ]

                    pivot_items = [(pivot_labels[k], pivot_map.get(k, "")) for k in ordered_keys]

                    non_empty_count = sum(1 for _, t in pivot_items if (t or "").strip())
                    if non_empty_count == 0:
                        st.error("6観点がすべて空です。少なくとも1つ以上入力してから反映してください。")
                    else:
                        st.session_state["pivot_axis_text_map"] = pivot_map
                        st.session_state["pivot_axis_labels"] = pivot_labels
                        st.session_state["pivot_axis_texts_committed"] = True

                        st.session_state["pivot_axis_texts"] = [
                            {"title": title, "text": (text or "").strip()} for title, text in pivot_items
                        ]

                        compiled = "\n\n".join(
                            [f"### {item['title']}\n{item['text']}" for item in st.session_state["pivot_axis_texts"] if item["text"]]
                        )
                        st.session_state["pivot_axis_text_compiled"] = compiled

                        st.success("『生成・比較』に反映しました。")
                        st.rerun()


        # =========================================================
        # TAB 2: 課題マトリクス選択 + PhaseA生成（KON〜SQ） + 左右比較
        # =========================================================
        with tab_gen:

            render_character_guide3(
                "KON〜SQ",
                "- キックオフノートからサブクエスチョンまでを生成して考察・比較するステップだよ。\n"
                "- ポイントは顧客課題を決めつけないこと。\n"
                "- はじめに中心となる課題を選んで「新規作成」を押してください。\n"
                "- 中心となる課題を変えて「新規作成」を押すと、別なKON～SQが生成されるよ。\n"
                "- ちなみに、SQとはサブクエスチョン（KONの「問い」に答えるために設定する下位項目）のことです。",
                img_width=500,
                kind="info",
            )

            st.markdown("### KON～サブクエスチョン")

            # =========================================================
            # ★課題ピボット（6観点）の受け取り（生成の軸）
            # =========================================================
            pivot_committed = bool(st.session_state.get("pivot_axis_texts_committed", False))
            pivot_map = st.session_state.get("pivot_axis_text_map", {}) or {}
            pivot_labels = st.session_state.get("pivot_axis_labels", {}) or {}

            if not pivot_committed or not pivot_map:
                st.warning("課題ピボット（軸）が未確定です。先に課題ピボット側で入力・確定してください。")
                selected_key = None
                selected_text = ""
            else:
                # 選択肢（キー）を固定順で並べる
                axis_keys = [
                    "c4_business_brand",
                    "c2_exec_summary",
                    "c1_next_action",
                    "c6_user_notes",
                ]

                def _fmt_axis(k: str) -> str:
                    # 表示を「日本語ラベル」にしたい場合はここ
                    return pivot_labels.get(k, k)

                st.selectbox(
                    "中心となる課題（課題ピボットの観点を選択）",
                    options=axis_keys,
                    key="pivot_axis_selected_key",
                    format_func=_fmt_axis,
                )

                selected_key = st.session_state.get("pivot_axis_selected_key")
                selected_text = (pivot_map.get(selected_key, "") or "").strip()

                if not selected_text:
                    st.warning("選択した観点のテキストが空です。課題ピボット側で入力してから確定してください。")

                with st.expander("選択中の軸（本文プレビュー）", expanded=False):
                    st.code(selected_text if selected_text else "（空）", language="text")

            # =========================================================
            # PhaseA 生成（Revision追加：KON〜SQまで）
            # =========================================================
            col_a, col_b = st.columns([1, 2], gap="small")
            with col_a:
                if st.button(
                    "新規作成",
                    use_container_width=True,
                    key="btn_phaseA_tab_gen",
                    disabled=(not pivot_committed or not selected_text),
                ):
                    with st.spinner("KON〜サブクエスチョンを生成しています..."):
                        ok, msg = run_phaseA_generation_and_append_revision(
                            axis_text=selected_text,
                            axis_source=f"pivot:{selected_key}",
                        )

                    if ok:
                        st.success("KON〜SQを追加しました。下で比較できます。")

                        revs_tmp = get_revisions() or []
                        if revs_tmp:
                            latest = revs_tmp[-1]
                            latest_id = latest.get("rev_id")

                            # ★追加：作成した最新Revisionをアクティブにする
                            if latest_id:
                                st.session_state["active_rev_id"] = latest_id
                                st.session_state.pop("__proposal_draft_hydrated", None)

                                # （左右表示順の制御は既にあなたが反映済み：左=最新、右=1つ前）
                                st.session_state["__force_compare_left_rev_id"] = latest_id

                                if len(revs_tmp) >= 2:
                                    prev_id = revs_tmp[-2].get("rev_id")
                                    if prev_id:
                                        st.session_state["__force_compare_right_rev_id"] = prev_id

                        st.rerun()  # ★成功時だけ rerun
                    else:
                        st.error(msg)
                        # ★失敗時は rerun しない（エラーを残す）

            with col_b:
                st.caption("複数の軸で “KON〜SQ” を素早く比較し、採用案を決めてから編集タブで詳細化します。")

            # =========================================================
            # Revision一覧（存在する場合）
            # =========================================================
            revs = get_revisions()
            if not revs:
                st.info("新規作成ボタンを押してください。")
            else:
                # ---------- 表示専用関数（ここで定義してOK） ----------

                def _render_kickoff_block_editable(rev: dict | None, title: str, key_prefix: str):
                    if not rev:
                        st.info("Revisionが未選択、または取得できませんでした。")
                        return

                    rev_id = rev.get("rev_id") or "no_rev"
                    k = (rev.get("kickoff") or {})

                    st.write(f"**{title}：{rev.get('label','')}**")
                    axis = rev.get("axis_source") or ""
                    pkey = rev.get("purpose_key") or ""
                    meta = f"axis:{axis}" if axis else f"purpose:{pkey}"
                    st.caption(f"{rev.get('created_at','')} / {meta}")

                    fields = ["目標", "現状", "ビジネス課題", "調査目的", "問い", "仮説"]

                    for field_label in fields:
                        widget_key = f"{key_prefix}__{rev_id}__kickoff__{field_label}"

                        # 初回だけ rev→UI に注入（毎回上書きしない）
                        if widget_key not in st.session_state:
                            st.session_state[widget_key] = k.get(field_label, "") or ""

                        st.text_area(
                            field_label,
                            height=120,
                            key=widget_key,
                            disabled=False,
                        )

                    # ★ここから下は「ループの外」で1回だけ
                    new_kickoff = dict(k)
                    for field_label in fields:
                        widget_key = f"{key_prefix}__{rev_id}__kickoff__{field_label}"
                        new_kickoff[field_label] = st.session_state.get(widget_key, "") or ""

                    # ★重要：最新のRevisionを取り直して、kickoffだけ差し替える（SQ巻き戻り防止）
                    rev_latest = find_revision(rev_id) or rev
                    rev2 = dict(rev_latest)
                    rev2["kickoff"] = new_kickoff
                    upsert_revision(rev2)



                def _render_subq_block_editable(rev: dict | None, title: str, key_prefix: str):
                    st.markdown(f"**{title}：問いの分解（サブクエスチョン）**")
                    if not rev:
                        st.info("Revisionが未選択、または取得できませんでした。")
                        return

                    rev_id = rev.get("rev_id") or "no_rev"
                    subq_list = (rev.get("subq_list") or []) or []

                    if not subq_list:
                        st.info("サブクエスチョンがありません。")
                        return

                    for i, sq in enumerate(subq_list, 1):
                        widget_key = f"{key_prefix}__{rev_id}__subq__{i}"
                        if widget_key not in st.session_state:
                            st.session_state[widget_key] = sq.get("subq") or sq.get("question") or ""

                        st.text_area(f"SQ{i}", key=widget_key, height=120, disabled=False)



                def _render_kon_sq_compact_view(rev: dict | None, title: str, key_prefix: str):
                    """
                    KONの「問い（メインクエスチョン）」とSQを、分析イメージタブと同じ階層表示で1つのtext_areaにまとめて表示（read-only）。
                    """
                    st.markdown(f"**{title}：KONの問い → サブクエスチョン（構造表示）**")

                    if not rev:
                        st.info("Revisionが未選択、または取得できませんでした。")
                        return

                    import re

                    rev_id = rev.get("rev_id") or "no_rev"

                    rev = find_revision(rev_id) or rev   # ★最新を取り直す

                    kickoff = (rev.get("kickoff") or {})
                    main_question_text = (kickoff.get("問い") or "").strip()
                    subq_list = (rev.get("subq_list") or []) or []

                    def split_main_questions(text: str):
                        """
                        分析イメージタブと同じ分割ロジック（番号/ Qx: などを想定）
                        """
                        if not text:
                            return []
                        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
                        questions, buf = [], ""
                        for line in lines:
                            m = re.match(r"^(?:\d+[\.\)]|Q\d+[:：])\s*(.+)", line)
                            if m:
                                if buf:
                                    questions.append(buf.strip())
                                buf = m.group(1)
                            else:
                                buf = (buf + " " + line).strip() if buf else line
                        if buf:
                            questions.append(buf.strip())
                        return questions or ([text.strip()] if text.strip() else [])

                    main_questions = split_main_questions(main_question_text)

                    # まず親（メインクエスチョン）ごとに箱を用意
                    if not main_questions:
                        grouped = {"(メインクエスチョン未設定)": []}
                    else:
                        grouped = {mq: [] for mq in main_questions}

                    # subq_list を親に紐付け
                    if subq_list:
                        for idx, sq in enumerate(subq_list):
                            explicit_parent = (sq.get("main") or sq.get("main_question") or "").strip()

                            if explicit_parent and explicit_parent in grouped:
                                grouped[explicit_parent].append(sq)
                            else:
                                # フォールバック：順番で均等に割り当て
                                if main_questions:
                                    parent = main_questions[idx % len(main_questions)]
                                else:
                                    parent = "(メインクエスチョン未設定)"
                                grouped.setdefault(parent, []).append(sq)

                    # 表示文字列を構造化
                    lines = []
                    lines.append("【問い（メインクエスチョン）→ サブクエスチョン】")
                    lines.append("")

                    if not grouped:
                        lines.append("（未設定）")
                    else:
                        for mq_idx, (mq, sqs) in enumerate(grouped.items(), 1):
                            lines.append(f"{mq_idx}. メインクエスチョン{mq_idx}：{mq}")
                            if not sqs:
                                lines.append("    ┗ （サブクエスチョン未設定）")
                            else:
                                for j, sq in enumerate(sqs, 1):
                                    sq_text = (sq.get("subq") or sq.get("question") or "").strip()
                                    lines.append(f"    ┗ {j}. サブクエスチョン{j}：{sq_text}")
                            lines.append("")

                    text = "\n".join(lines).strip()

                    widget_key = f"{key_prefix}__{rev_id}__kon_sq_view"
                    st.session_state[widget_key] = text  # 表示専用なので毎回上書き

                    st.text_area(
                        "KON〜SQ（構造表示）",
                        key=widget_key,
                        height=380,
                        disabled=True,
                    )




                # ---------- rev_id 方式：選択値は rev_id に統一 ----------
                rev_ids = [r["rev_id"] for r in revs if r.get("rev_id")]

                if not rev_ids:
                    st.error("Revisionの rev_id が取得できませんでした。append_revision の実装を確認してください。")
                else:
                    # rev_id -> rev を引ける辞書（format_func で使う）
                    rev_by_id = {r.get("rev_id"): r for r in revs if r.get("rev_id")}

                    def _display_rev_name(rev_id: str) -> str:
                        """
                        セレクトボックスの選択肢表記を「中心となる課題（課題ピボット）のKEY」基準にする。
                        axis_source が "pivot:<key>" の場合は pivot_labels を使って表示する。
                        """
                        r = rev_by_id.get(rev_id, {}) or {}
                        axis = (r.get("axis_source") or "").strip()

                        if axis.startswith("pivot:"):
                            pivot_key = axis.split("pivot:", 1)[1].strip()

                            # 表示を「KEYのみ」にしたい場合はこちら
                            # return pivot_key

                            # 表示を「日本語ラベル + (KEY)」にしたい場合はこちら（おすすめ）
                            return f"{pivot_labels.get(pivot_key, pivot_key)} ({pivot_key})"

                        # pivot由来でないRevisionは label にフォールバック
                        return r.get("label") or rev_id

                    def _find_rev_id_by_stage(stage: str) -> str | None:
                        for r in revs:
                            if r.get("stage") == stage:
                                return r.get("rev_id")
                        return None

                    default_rev_id = _find_rev_id_by_stage("default")

                    latest_phaseA_rev_id = None
                    for r in reversed(revs):
                        if r.get("stage") in ["phaseA", "ko_sq", "kickoff_subq"]:
                            latest_phaseA_rev_id = r.get("rev_id")
                            break
                    if latest_phaseA_rev_id is None:
                        latest_phaseA_rev_id = rev_ids[-1]

                    # 初回だけ（rev_idで初期化）
                    if "compare_left_rev_id" not in st.session_state:
                        st.session_state["compare_left_rev_id"] = default_rev_id or rev_ids[0]
                    if "compare_right_rev_id" not in st.session_state:
                        st.session_state["compare_right_rev_id"] = latest_phaseA_rev_id or rev_ids[-1]

                    # ★新規作成直後は左/右を強制（rev_id）
                    force_left_id = st.session_state.pop("__force_compare_left_rev_id", None)
                    if force_left_id and force_left_id in rev_ids:
                        st.session_state["compare_left_rev_id"] = force_left_id

                    force_right_id = st.session_state.pop("__force_compare_right_rev_id", None)
                    if force_right_id and force_right_id in rev_ids:
                        st.session_state["compare_right_rev_id"] = force_right_id

                    st.markdown("#### KON～SQ比較（左右）")

                    c1, c2 = st.columns(2, gap="large")

                    with c1:
                        st.selectbox(
                            "左に表示するKON～SQ（比較A）",
                            options=rev_ids,
                            key="compare_left_rev_id",
                            format_func=_display_rev_name,
                        )
                        left_rev = find_revision(st.session_state["compare_left_rev_id"])
                        left_subq = (left_rev or {}).get("subq_list", [])

                        if left_rev:
                            _render_kickoff_block_editable(left_rev, "左（比較A）", "cmp_left")
                        else:
                            st.info("左側に表示するRevisionを選択してください。")

                    with c2:
                        st.selectbox(
                            "右に表示するKON～SQ（比較B）",
                            options=rev_ids,
                            key="compare_right_rev_id",
                            format_func=_display_rev_name,
                        )
                        right_rev = find_revision(st.session_state["compare_right_rev_id"])
                        right_subq = (right_rev or {}).get("subq_list", [])

                        if right_rev:
                            _render_kickoff_block_editable(right_rev, "右（比較B）", "cmp_right")
                        else:
                            st.info("右側に表示するRevisionを選択してください。")

                    # --- 問いの分解（左右比較） ---
                    # --- KONの問い → SQ（左右比較：統合表示） ---
                    c1, c2 = st.columns(2, gap="large")

                    with c1:
                        _render_kon_sq_compact_view(left_rev, "左（比較A）", "cmp_left")

                        # 編集用（保存ボタンで確定）
                        with st.expander("SQの編集（左）", expanded=False):
                            _render_subq_block_editable(left_rev, "左（比較A）", "cmp_left")

                            if st.button("左のSQを保存", use_container_width=True, key="btn_save_left_sq"):
                                if left_rev and left_rev.get("rev_id"):
                                    rev_id = left_rev["rev_id"]

                                    # 念のため最新のRevisionを取り直す（保存競合防止）
                                    rev_latest = find_revision(rev_id) or left_rev
                                    subq_list = (rev_latest.get("subq_list") or []) or []

                                    new_list = []
                                    lines_for_raw = []

                                    for i, sq in enumerate(subq_list, 1):
                                        widget_key = f"cmp_left__{rev_id}__subq__{i}"
                                        txt = (st.session_state.get(widget_key, "") or "").strip()

                                        new_sq = dict(sq)
                                        # ★重要：subq だけでなく question も同じ値に揃える（戻り防止）
                                        new_sq["subq"] = txt
                                        new_sq["question"] = txt

                                        new_list.append(new_sq)
                                        if txt:
                                            lines_for_raw.append(f"SQ{i}: {txt}")

                                    rev2 = dict(rev_latest)
                                    rev2["subq_list"] = new_list
                                    rev2["subquestions_raw"] = "\n".join(lines_for_raw)

                                    # ★追加：analysis_blocks の subq も subq_list に追随させる（戻り防止の本丸）
                                    blocks = (rev_latest.get("analysis_blocks") or []) or []
                                    blocks2 = []
                                    for idx, b in enumerate(blocks):
                                        b2 = dict(b or {})
                                        if idx < len(new_list):
                                            b2["subq"] = new_list[idx].get("subq", "") or ""
                                        blocks2.append(b2)
                                    rev2["analysis_blocks"] = blocks2

                                    upsert_revision(rev2)

                                    # 比較レーン固定して rerun
                                    st.session_state["__force_compare_left_rev_id"] = st.session_state.get("compare_left_rev_id")
                                    st.session_state["__force_compare_right_rev_id"] = st.session_state.get("compare_right_rev_id")

                                    st.success("左（比較A）のSQを保存しました。")
                                    st.rerun()
                                else:
                                    st.warning("左（比較A）のRevisionが取得できません。")

                    with c2:
                        _render_kon_sq_compact_view(right_rev, "右（比較B）", "cmp_right")

                        # 編集用（保存ボタンで確定）
                        with st.expander("SQの編集（右）", expanded=False):
                            _render_subq_block_editable(right_rev, "右（比較B）", "cmp_right")

                            if st.button("右のSQを保存", use_container_width=True, key="btn_save_right_sq"):
                                if right_rev and right_rev.get("rev_id"):
                                    rev_id = right_rev["rev_id"]

                                    # 念のため最新のRevisionを取り直す（保存競合防止）
                                    rev_latest = find_revision(rev_id) or right_rev
                                    subq_list = (rev_latest.get("subq_list") or []) or []

                                    new_list = []
                                    lines_for_raw = []

                                    for i, sq in enumerate(subq_list, 1):
                                        widget_key = f"cmp_right__{rev_id}__subq__{i}"
                                        txt = (st.session_state.get(widget_key, "") or "").strip()

                                        new_sq = dict(sq)
                                        # ★重要：subq だけでなく question も同じ値に揃える（戻り防止）
                                        new_sq["subq"] = txt
                                        new_sq["question"] = txt

                                        new_list.append(new_sq)
                                        if txt:
                                            lines_for_raw.append(f"SQ{i}: {txt}")

                                    rev2 = dict(rev_latest)
                                    rev2["subq_list"] = new_list
                                    rev2["subquestions_raw"] = "\n".join(lines_for_raw)

                                    # ★追加：analysis_blocks の subq も subq_list に追随させる（戻り防止の本丸）
                                    blocks = (rev_latest.get("analysis_blocks") or []) or []
                                    blocks2 = []
                                    for idx, b in enumerate(blocks):
                                        b2 = dict(b or {})
                                        if idx < len(new_list):
                                            b2["subq"] = new_list[idx].get("subq", "") or ""
                                        blocks2.append(b2)
                                    rev2["analysis_blocks"] = blocks2

                                    upsert_revision(rev2)


                                    # 比較レーン固定して rerun
                                    st.session_state["__force_compare_left_rev_id"] = st.session_state.get("compare_left_rev_id")
                                    st.session_state["__force_compare_right_rev_id"] = st.session_state.get("compare_right_rev_id")

                                    st.success("右（比較B）のSQを保存しました。")
                                    st.rerun()
                                else:
                                    st.warning("右（比較B）のRevisionが取得できません。")

                    # =========================================================
                    # KON～SQ（左右比較）の採用 → 編集タブ（分析イメージ）へ反映
                    # =========================================================
                    st.markdown("---")
                    st.markdown("### この案を採用して、分析イメージタブへ反映")

                    col_apply_l, col_apply_r = st.columns(2, gap="small")

                    with col_apply_l:
                        if st.button("左（比較A）を採用して反映", use_container_width=True, key="btn_apply_left_to_active"):
                            if left_rev and left_rev.get("rev_id"):
                                rid = left_rev["rev_id"]

                                # ★採用前に、いま編集中のSQをrevに保存してから反映（戻り防止）
                                rev_latest = find_revision(rid) or left_rev
                                subq_list = (rev_latest.get("subq_list") or []) or []

                                new_list = []
                                lines_for_raw = []
                                for i, sq in enumerate(subq_list, 1):
                                    widget_key = f"cmp_left__{rid}__subq__{i}"
                                    txt = (st.session_state.get(widget_key, "") or "").strip()
                                    new_sq = dict(sq or {})
                                    new_sq["subq"] = txt
                                    new_sq["question"] = txt
                                    new_list.append(new_sq)
                                    if txt:
                                        lines_for_raw.append(f"SQ{i}: {txt}")

                                rev2 = dict(rev_latest)
                                rev2["subq_list"] = new_list
                                rev2["subquestions_raw"] = "\n".join(lines_for_raw)

                                blocks = (rev_latest.get("analysis_blocks") or []) or []
                                blocks2 = []
                                for idx, b in enumerate(blocks):
                                    b2 = dict(b or {})
                                    if idx < len(new_list):
                                        b2["subq"] = new_list[idx].get("subq", "") or ""
                                    blocks2.append(b2)
                                rev2["analysis_blocks"] = blocks2

                                upsert_revision(rev2)

                                # ★ここでは set_active_revision() を呼ばない（あなたの方針でOK）
                                st.session_state["active_rev_id"] = rid
                                st.session_state["pending_apply_rev_id"] = rid

                                st.session_state["__force_compare_left_rev_id"] = st.session_state.get("compare_left_rev_id")
                                st.session_state["__force_compare_right_rev_id"] = st.session_state.get("compare_right_rev_id")

                                st.rerun()

                            else:
                                st.warning("左（比較A）のRevisionが取得できません。")

                    with col_apply_r:
                        if st.button("（比較B）を採用して反映", use_container_width=True, key="btn_apply_right_to_active"):
                            if right_rev and right_rev.get("rev_id"):
                                rid = right_rev["rev_id"]

                                # ★採用前に、いま編集中のSQをrevに保存してから反映（戻り防止）
                                rev_latest = find_revision(rid) or right_rev
                                subq_list = (rev_latest.get("subq_list") or []) or []

                                new_list = []
                                lines_for_raw = []
                                for i, sq in enumerate(subq_list, 1):
                                    widget_key = f"cmp_right__{rid}__subq__{i}"
                                    txt = (st.session_state.get(widget_key, "") or "").strip()
                                    new_sq = dict(sq or {})
                                    new_sq["subq"] = txt
                                    new_sq["question"] = txt
                                    new_list.append(new_sq)
                                    if txt:
                                        lines_for_raw.append(f"SQ{i}: {txt}")

                                rev2 = dict(rev_latest)
                                rev2["subq_list"] = new_list
                                rev2["subquestions_raw"] = "\n".join(lines_for_raw)

                                blocks = (rev_latest.get("analysis_blocks") or []) or []
                                blocks2 = []
                                for idx, b in enumerate(blocks):
                                    b2 = dict(b or {})
                                    if idx < len(new_list):
                                        b2["subq"] = new_list[idx].get("subq", "") or ""
                                    blocks2.append(b2)
                                rev2["analysis_blocks"] = blocks2

                                upsert_revision(rev2)

                                # ★ここでは set_active_revision() を呼ばない（あなたの方針でOK）
                                st.session_state["active_rev_id"] = rid
                                st.session_state["pending_apply_rev_id"] = rid

                                st.session_state["__force_compare_left_rev_id"] = st.session_state.get("compare_left_rev_id")
                                st.session_state["__force_compare_right_rev_id"] = st.session_state.get("compare_right_rev_id")

                                st.rerun()

                            else:
                                st.warning("右（比較B）のRevisionが取得できません。")



                # --- ここから下は既存の削除UIへ（あなたのコードをそのまま続けてOK） ---

                # =========================================================
                # Revision 削除 UI（生成・比較タブ）
                # =========================================================
                st.markdown("---")
                st.markdown("### Revisionの削除")

                # default は削除不可にする
                deletable_revs = [r for r in revs if r.get("stage") != "default"]

                if not deletable_revs:
                    st.caption("削除可能なRevisionはありません。")
                else:
                    del_options = {r["label"]: r["rev_id"] for r in deletable_revs}
                    del_labels = list(del_options.keys())

                    del_label = st.selectbox(
                        "削除するRevisionを選択",
                        options=del_labels,
                        key="delete_revision_label",
                    )
                    del_rev_id = del_options[del_label]

                    st.warning("この操作は元に戻せません。")

                    confirm = st.checkbox(
                        "本当にこのRevisionを削除する",
                        key="delete_revision_confirm",
                    )

                    if st.button(
                        "選択したRevisionを削除",
                        use_container_width=True,
                        disabled=not confirm,
                    ):
                        ok, msg = delete_revision(del_rev_id)
                        if ok:
                            st.success("Revisionを削除しました。")
                            st.rerun()
                        else:
                            st.error(msg)


        # =========================================================
        # TAB 3: アクティブRevision選択 + 企画内容レビュー（編集UI / PhaseB詳細化）
        # =========================================================
        with tab_edit:

            render_character_guide3(
                "分析イメージ～対象者条件～調査項目の検討",
                "- ここはさっき作った“KON～SQ”を元に「分析イメージ」、「対象者条件」、「調査項目」を検討するステップだよ。\n"
                "- まず検討を進めたい「KON～SQ」を選んでください。\n"
                "- その後は「新規作成」を押しながらそれぞれのそれぞれの工程の内容を確認してね。\n"
                "- ページの最下部に「PPTに保存」、「一時保存」のボタンがあるからそこから保存してね。\n",
                img_width=500,
                kind="info",
            )

            st.markdown("### 編集・PPT反映")

            # --- tab_edit 冒頭：revs取得 ---
            revs = get_revisions()

            # active_rev_id が消えている／存在しない rev_id を指している場合の保険
            active_id = st.session_state.get("active_rev_id")
            if active_id and not find_revision(active_id):
                st.session_state["active_rev_id"] = None
                active_id = None

            # まだアクティブが無い場合は最新をアクティブ化（初回導線）
            if not active_id and revs:
                st.session_state["active_rev_id"] = revs[-1].get("rev_id")
                sync_active_revision_to_session_keys()

            if not revs:
                st.info("まだRevisionがありません。『生成・比較』タブで新規作成（KON〜SQ）を実行してください。")
                st.stop()
            else:
                # =========================================================
                # active選択（編集・PPT反映に使うRevision）
                #   - selectboxの値は rev_id（内部キーを安定化）
                #   - 表示は KON～SQ と同じ pivot_labels + (pivot_key)
                # =========================================================

                # KON～SQタブと同じ pivot 情報（無ければ空でOK）
                pivot_labels = st.session_state.get("pivot_axis_labels", {}) or {}

                # rev_id リスト＆参照辞書
                rev_ids = [r.get("rev_id") for r in revs if r.get("rev_id")]
                rev_by_id = {r.get("rev_id"): r for r in revs if r.get("rev_id")}

                def _display_rev_name(rev_id: str) -> str:
                    r = rev_by_id.get(rev_id, {}) or {}
                    axis = (r.get("axis_source") or "").strip()

                    # axis_source="pivot:<key>" の場合は KON～SQ と同表記
                    if axis.startswith("pivot:"):
                        pivot_key = axis.split("pivot:", 1)[1].strip()
                        return f"{pivot_labels.get(pivot_key, pivot_key)} ({pivot_key})"

                    # pivot由来でない場合は label にフォールバック
                    return r.get("label") or rev_id

                # 現在の active を index に反映（存在しなければ最後）
                active_id = st.session_state.get("active_rev_id")
                if active_id not in rev_ids:
                    active_id = (rev_ids[-1] if rev_ids else None)

                # selector key（revsの増減で変わる）
                rev_sig = "__".join(rev_ids)  # rev_id だけでOK
                selector_key = f"active_revision_selector__{rev_sig}"

                selected_rev_id = st.selectbox(
                    "編集・PPT反映に使うRevision（アクティブ）",
                    options=rev_ids,
                    index=(rev_ids.index(active_id) if active_id in rev_ids else len(rev_ids) - 1),
                    format_func=_display_rev_name,
                    key=selector_key,
                )

                if selected_rev_id != st.session_state.get("active_rev_id"):
                    set_active_revision(selected_rev_id)

                    # ★必ず tabs より前の apply 処理に流す（ウィジェット生成前に同期するため）
                    st.session_state["pending_apply_rev_id"] = selected_rev_id

                    # ★切替直後の “毎回save” を1回だけ止める（巻き戻り防止）
                    st.session_state["__skip_autosave_once"] = True

                    st.rerun()




            # =========================================================
            # PhaseA（KON〜SQ）が同期されているかの確認（proposal_draft_generated に依存しない）
            # =========================================================
            has_kickoff = any(
                (st.session_state.get(k) or "").strip()
                for k in ["ai_目標", "ai_現状", "ai_ビジネス課題", "ai_調査目的", "ai_問い", "ai_仮説"]
            )
            subq_list = st.session_state.get("subq_list", []) or []
            has_sq = len(subq_list) > 0

            if not has_kickoff and not has_sq:
                st.info("アクティブRevisionの内容がまだ同期されていません。『生成・比較』タブで作成したRevisionを選択してから、ここに戻ってください。")
                st.stop()

            # =========================================================
            # 1. キックオフノート
            # =========================================================
            st.markdown("### 1. キックオフノート（①〜⑥）")
            st.text_area("① 目標（to be）", key="ai_目標", height=80)
            st.text_area("② 現状（as is）", key="ai_現状", height=80)
            st.text_area("③ ビジネス課題", key="ai_ビジネス課題", height=80)
            st.text_area("④ 調査目的", key="ai_調査目的", height=80)
            st.text_area("⑤ 問い", key="ai_問い", height=80)
            st.text_area("⑥ 仮説", key="ai_仮説", height=80)

            # if get_active_revision() is not None:
            #     save_session_keys_to_active_revision()
            maybe_autosave_active_revision()

            st.markdown("---")

            # =========================================================
            # 2. 問いの分解（サブクエスチョン）
            # =========================================================
            st.markdown("### 2. 問いの分解（サブクエスチョン）")

            import math
            import re

            main_question_text = (st.session_state.get("ai_問い", "") or "").strip()
            subq_list = st.session_state.get("subq_list", []) or []

            def split_main_questions(text: str):
                if not text:
                    return []
                lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
                questions, buf = [], ""
                for line in lines:
                    m = re.match(r"^(?:\d+[\.\)]|Q\d+[:：])\s*(.+)", line)
                    if m:
                        if buf:
                            questions.append(buf.strip())
                        buf = m.group(1)
                    else:
                        buf = (buf + " " + line).strip() if buf else line
                if buf:
                    questions.append(buf.strip())
                return questions or ([text.strip()] if text.strip() else [])

            main_questions = split_main_questions(main_question_text)

            # --- 1) 親（メインQ）を用意 ---
            if not main_questions:
                grouped = {"(メインクエスチョン未設定)": []}
            else:
                grouped = {mq: [] for mq in main_questions}

            # --- 2) subq を親へ紐付け ---
            # 2-1) 明示的に main / main_question が入っていればそれを優先
            remaining = []
            for sq in subq_list:
                d = dict(sq or {})
                sq_txt = (d.get("subq") or d.get("question") or "").strip()
                d["subq"] = sq_txt
                d["question"] = sq_txt

                parent = (d.get("main") or d.get("main_question") or "").strip()
                if parent and parent in grouped:
                    grouped[parent].append(d)
                else:
                    remaining.append(d)

            # 2-2) 明示親が無い分は「順序で“塊”割当」して構造化表示にする
            # （例：メインQ3つ、SQ9つ → 3つずつ割当）
            if remaining:
                if main_questions:
                    n_main = len(main_questions)
                    chunk = max(1, math.ceil(len(remaining) / n_main))
                    idx = 0
                    for mq in main_questions:
                        grouped[mq].extend(remaining[idx: idx + chunk])
                        idx += chunk
                        if idx >= len(remaining):
                            break
                else:
                    grouped["(メインクエスチョン未設定)"].extend(remaining)

            # --- 3) 表示テキスト（ツリー）を作成 ---
            lines = []
            lines.append("【問い（メインクエスチョン）→ サブクエスチョン（構造表示）】")
            lines.append("")

            for mq_idx, (mq, sqs) in enumerate(grouped.items(), 1):
                lines.append(f"{mq_idx}. メインクエスチョン{mq_idx}：{mq}")
                if not sqs:
                    lines.append("    ┗ （サブクエスチョン未設定）")
                else:
                    for j, sq in enumerate(sqs, 1):
                        lines.append(f"    ┗ {j}. サブクエスチョン{j}：{sq.get('subq','')}")
                lines.append("")

            structured_text = "\n".join(lines).strip()

            # ★重要：value= を使わず、key を唯一のソースにする（表示が古くなるのを防ぐ）
            active_id = st.session_state.get("active_rev_id") or "no_rev"
            subq_view_key = f"subq_structured_view__{active_id}"

            st.session_state[subq_view_key] = structured_text  # 毎回上書き（表示用）
            st.text_area(
                "サブクエスチョン一覧（構造化）",
                key=subq_view_key,
                height=260,
                disabled=True,
            )

            st.markdown("---")

            # =========================================================
            # 3. 分析アプローチ（PhaseBで生成）
            # =========================================================
            st.markdown("### 3. 分析アプローチ")

            active_id = st.session_state.get("active_rev_id") or "no_rev"

            colb1, colb2 = st.columns([1, 3], gap="small")

            # rev_dbg = get_active_revision() or {}
            # ss_sq1 = ((st.session_state.get("subq_list") or [{}])[0] or {}).get("subq")
            # rev_sq1 = (((rev_dbg.get("subq_list") or [{}])[0]) or {}).get("subq")

            # st.write("DEBUG ss subq_list[0].subq:", ss_sq1)
            # st.write("DEBUG rev subq_list[0].subq:", rev_sq1)
            # st.write("DEBUG ss ai_subquestions head:", (st.session_state.get("ai_subquestions") or "")[:80])
            # st.write("DEBUG rev subquestions_raw head:", (rev_dbg.get("subquestions_raw") or "")[:80])



            with colb1:
                if st.button("新規作成", use_container_width=True, key=f"btn_gen_analysis_phaseB__{active_id}"):
                    with st.spinner("分析アプローチを生成しています..."):
                        ok, msg = generate_analysis_approach_draft()

                    if ok:
                        # 生成後、必ず保存（rev側にも確定）
                        if get_active_revision() is not None:
                            save_session_keys_to_active_revision()
                        st.success("分析アプローチを生成しました。")
                        st.rerun()
                    else:
                        st.error(msg or "分析アプローチ生成に失敗しました。")

            with colb2:
                st.caption("生成・比較タブではKON〜SQまで。ここで選択したRevisionを詳細化します。")

            analysis_blocks = st.session_state.get("analysis_blocks", []) or []

            if not analysis_blocks:
                st.info("分析アプローチ案がまだありません。上のボタンで生成してください。")
            else:
                # ★revごとの widget key を使う（これが本丸）
                for i, blk in enumerate(analysis_blocks, 1):
                    st.markdown(f"#### サブクエスチョン {i}")

                    k_subq = f"analysis_subq_{i}__{active_id}"
                    k_axis = f"analysis_axis_{i}__{active_id}"
                    k_items = f"analysis_items_{i}__{active_id}"
                    k_app  = f"analysis_approach_{i}__{active_id}"
                    k_hypo = f"analysis_hypothesis_{i}__{active_id}"

                    # 初回だけブロック値をUIへ注入（毎回上書きしない）
                    if k_subq not in st.session_state: st.session_state[k_subq] = blk.get("subq", "") or ""
                    if k_axis not in st.session_state: st.session_state[k_axis] = blk.get("axis", "") or ""
                    if k_items not in st.session_state: st.session_state[k_items] = blk.get("items", "") or ""
                    if k_app  not in st.session_state: st.session_state[k_app]  = blk.get("approach", "") or ""
                    if k_hypo not in st.session_state: st.session_state[k_hypo] = blk.get("hypothesis", "") or ""

                    # value=は使わない（keyのみ）
                    st.text_area("サブクエスチョン", height=60, key=k_subq)
                    st.text_area("分析軸（セグメント）", height=60, key=k_axis)
                    st.text_area("評価項目", height=60, key=k_items)
                    st.text_area("主な分析アプローチ", height=80, key=k_app)
                    st.text_area("検証する仮説", height=80, key=k_hypo)

                    # UI→analysis_blocksへ反映（このループ内で常に同期）
                    blk["subq"] = st.session_state.get(k_subq, "") or ""
                    blk["axis"] = st.session_state.get(k_axis, "") or ""
                    blk["items"] = st.session_state.get(k_items, "") or ""
                    blk["approach"] = st.session_state.get(k_app, "") or ""
                    blk["hypothesis"] = st.session_state.get(k_hypo, "") or ""

                    st.markdown("---")

                # ★最終的に session_state の analysis_blocks を更新して確定
                st.session_state["analysis_blocks"] = analysis_blocks




            # =========================================================
            # 4. 対象者条件（AI生成ボタンを追加）
            # =========================================================
            st.markdown("### 4. 対象者条件案")

            coltc1, coltc2 = st.columns([1, 3], gap="small")
            with coltc1:
                if st.button("新規作成", use_container_width=True, key="btn_gen_target_condition_phaseB"):
                    with st.spinner("対象者条件案を生成しています..."):
                        ok, msg = generate_target_condition_draft()
                    if ok:
                        st.session_state["__dbg_tc_timeline"] = []

                        # 生成直後（session側）
                        st.session_state["__dbg_tc_timeline"].append(
                            ("after_gen",
                            (st.session_state.get("ai_target_condition") or "")[:20],
                            len((st.session_state.get("ai_target_condition") or "").strip()))
                        )

                        # save 呼び出し直前
                        st.session_state["__dbg_tc_timeline"].append(
                            ("before_save_call",
                            (st.session_state.get("ai_target_condition") or "")[:20],
                            len((st.session_state.get("ai_target_condition") or "").strip()))
                        )
                        
                        gen_text = st.session_state.get("ai_target_condition", "") or ""
                        st.session_state["ai_target_condition_editor"] = gen_text
                        st.session_state["ai_target_condition"] = gen_text

                        if get_active_revision() is not None:
                            save_session_keys_to_active_revision()

                        # save 呼び出し直後（session側が消されていないか）
                        st.session_state["__dbg_tc_timeline"].append(
                            ("after_save_call",
                            (st.session_state.get("ai_target_condition") or "")[:20],
                            len((st.session_state.get("ai_target_condition") or "").strip()))
                        )

                        # 保存結果（rev側に入っているか）
                        rev2 = get_active_revision()
                        st.session_state["__dbg_tc_rev_len_after_save"] = len(((rev2 or {}).get("target_condition") or "").strip())

                        st.success("対象者条件案を生成しました（debug）")
                        st.rerun()

                    else:
                        st.error(msg or "対象者条件案の生成に失敗しました。")

            with coltc2:
                st.caption("軸（課題ピボット）・KON・サブクエスチョンに整合する対象者条件を提案します。")


            # if st.session_state.get("__dbg_tc_timeline"):
            #     st.warning(f"DEBUG tc timeline: {st.session_state['__dbg_tc_timeline']}")
            #     st.warning(f"DEBUG tc rev_len_after_save: {st.session_state.get('__dbg_tc_rev_len_after_save')}")


            # ① editor を唯一のUIソースにする（value= は使わない）
            st.text_area(
                "対象者条件案（編集可）",
                height=320,
                key="ai_target_condition_editor",
            )

            # ② 保存用キーへ同期
            editor_val = st.session_state.get("ai_target_condition_editor", "") or ""
            if editor_val.strip():
                st.session_state["ai_target_condition"] = editor_val
            else:
                # editorが空でも、既に ai_target_condition があるなら保持
                st.session_state["ai_target_condition"] = st.session_state.get("ai_target_condition", "") or ""


            # if get_active_revision() is not None:
            #     save_session_keys_to_active_revision()

            st.markdown("---")

            # =========================================================
            # 5. 調査項目案（分析アプローチ連動 / 10/20/30/40は作らない）
            # =========================================================
            st.markdown("### 5. 調査項目案")

            analysis_blocks_norm = normalize_analysis_blocks(st.session_state.get("analysis_blocks", []) or [])
            if not analysis_blocks_norm:
                st.info("分析アプローチがありません。先に③（分析アプローチ）を生成してください。")
            else:
                colx1, colx2 = st.columns([1, 3], gap="small")
                with colx1:
                    if st.button("新規作成", use_container_width=True, key="btn_gen_linked_items"):
                        with st.spinner("分析アプローチに紐づく調査項目を生成しています..."):
                            ok, msg = generate_survey_items_linked_draft()
                        if ok:
                            if get_active_revision() is not None:
                                save_session_keys_to_active_revision()
                            st.success("調査項目を生成しました。")
                            st.rerun()
                        else:
                            st.error(msg)

                with colx2:
                    st.caption("下の表で “分析アプローチ→調査項目” の内容（形式/尺度/表頭・表側/採用）を編集できます。")

                rows = st.session_state.get("survey_item_rows", []) or []
                if not rows:
                    st.info("連動調査項目がまだありません。上のボタンで生成してください。")
                else:
                    import pandas as pd

                    df = pd.DataFrame(rows)

                    base_cols = [
                        "sq_id", "sq_subq", "items", "approach",
                        "var_name", "item_text", "recommended_type", "recommended_scale",
                        "priority", "table_role", "is_selected"
                    ]
                    for c in base_cols:
                        if c not in df.columns:
                            df[c] = ""

                    df["is_selected"] = df["is_selected"].apply(
                        lambda x: True if str(x).strip() in ["", "True", "true", "1", "yes", "Yes", "Y", "y"] else False
                    )

                    st.markdown("#### 分析アプローチ（表頭） × 調査項目（行）")

                    edit_cols = ["is_selected", "sq_id", "item_text", "recommended_type", "recommended_scale", "table_role"]

                    colm1, colm2, colm3 = st.columns([1, 1, 2], gap="small")
                    with colm1:
                        view_mode = st.radio(
                            "表示",
                            options=["採用のみ", "全件"],
                            horizontal=True,
                            index=0,
                            key=f"survey_item_view_mode__{st.session_state.get('active_rev_id','no_rev')}",
                        )
                    with colm2:
                        if st.button("採用のみを別途出力（保存）", use_container_width=True, key="btn_export_selected_items"):
                            selected_df = df[df["is_selected"] == True].copy()
                            st.session_state["survey_item_rows_selected"] = selected_df.to_dict(orient="records")
                            if get_active_revision() is not None:
                                save_session_keys_to_active_revision()
                            st.success(f"採用のみを保存しました（{len(selected_df)}件）。")
                    with colm3:
                        st.caption("採用=チェックあり。まず全件生成→不要を外す運用を推奨します。")

                    display_df = df[df["is_selected"] == True].copy() if view_mode == "採用のみ" else df.copy()
                    editor_df = display_df[edit_cols].copy()

                    edited = st.data_editor(
                        editor_df,
                        hide_index=True,
                        num_rows="dynamic",
                        use_container_width=True,
                        column_config={
                            "is_selected": st.column_config.CheckboxColumn("採用", help="チェックあり＝採用（出力対象）"),
                            "sq_id": st.column_config.TextColumn("sq_id", disabled=True),
                            "item_text": st.column_config.TextColumn("調査項目", width="large"),
                            "recommended_type": st.column_config.SelectboxColumn(
                                "形式",
                                options=["SA", "MA", "尺度", "数値", "自由回答"],
                                help="推奨設問形式",
                            ),
                            "recommended_scale": st.column_config.TextColumn(
                                "尺度（必要なら）",
                                help="例：5件法（1=不満〜5=満足）など。SA/MA/自由回答なら空でOK",
                            ),
                            "table_role": st.column_config.SelectboxColumn(
                                "表頭/表側",
                                options=["表頭", "表側"],
                                help="表頭=分類軸（属性・セグメント）、表側=評価/行動など測る項目",
                            ),
                        },
                        key=f"survey_item_rows_editor__{st.session_state.get('active_rev_id','no_rev')}",
                    )

                    master = df[base_cols].copy()

                    def _row_key(d):
                        return (str(d.get("sq_id", "")).strip(), str(d.get("item_text", "")).strip())

                    master_map = {_row_key(r): i for i, r in master.iterrows()}
                    apply_cols = ["is_selected", "recommended_type", "recommended_scale", "table_role", "item_text"]

                    for _, r in edited.iterrows():
                        k = _row_key(r)
                        if k in master_map:
                            i = master_map[k]
                            for c in apply_cols:
                                master.at[i, c] = r.get(c, master.at[i, c])
                        else:
                            new_row = {c: "" for c in base_cols}
                            for c in apply_cols:
                                new_row[c] = r.get(c, "")
                            new_row["sq_id"] = r.get("sq_id", "")
                            new_row["priority"] = 3
                            new_row["is_selected"] = bool(r.get("is_selected", True))
                            master = pd.concat([master, pd.DataFrame([new_row])], ignore_index=True)

                    master["is_selected"] = master["is_selected"].apply(
                        lambda x: True if str(x).strip() in ["True", "true", "1", "yes", "Yes", "Y", "y"] else False
                    )
                    master["priority"] = pd.to_numeric(master["priority"], errors="coerce").fillna(3).astype(int).clip(1, 5)

                    st.session_state["survey_item_rows"] = master.to_dict(orient="records")

                    # if get_active_revision() is not None:
                    #     save_session_keys_to_active_revision()
                    maybe_autosave_active_revision()

            # =========================================================
            # PPT反映＆一時保存（中央ペイン）
            # =========================================================
            st.markdown("---")
            st.markdown("### 保存（PPT + 調査項目Excel + JSON をZIPで出力）")

            if st.button("保存（ZIP出力を作成）", use_container_width=True, key="btn_build_zip"):

                # ★ここを追加：テンプレパスが無ければ先に落とす or エラーにする
                if not st.session_state.get("pptx_path"):
                    st.error("pptx_path が未設定です。PPTテンプレートを先にセットしてください（BlobからDLする処理が必要です）。")
                else:
                    try:
                        zip_bytes, zip_name = _make_output_zip_bytes(st.session_state)
                        st.session_state["final_zip_bytes"] = zip_bytes
                        st.session_state["final_zip_name"] = zip_name
                        st.success("ZIPを作成しました。下のボタンからダウンロードできます。")
                    except Exception as e:
                        st.error(f"保存用ZIPの作成に失敗しました: {e}")

            if st.session_state.get("final_zip_bytes"):
                st.download_button(
                    "ZIPをダウンロード",
                    data=st.session_state["final_zip_bytes"],
                    file_name=st.session_state.get("final_zip_name", "output.zip"),
                    mime="application/zip",
                    use_container_width=True,
                    key="btn_download_zip",
                )


    elif mode == "case_review":
        render_case_review_screen() 
